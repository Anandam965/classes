import os
import uuid
import time
import json
import html
import requests
from datetime import date, timedelta

import streamlit as st
import streamlit.components.v1 as components
import pandas as pd
from supabase import create_client
import google.generativeai as genai

# =========================
# IMGBB IMAGE UPLOAD
# =========================
def upload_image_to_imgbb(image_file):
    import base64
    try:
        api_key = st.secrets.get("IMGBB_API_KEY", "")
        if not api_key:
            st.error("IMGBB_API_KEY secrets  !")
            return None
        img_data = base64.b64encode(image_file.read()).decode("utf-8")
        response = requests.post("https://api.imgbb.com/1/upload", data={"key": api_key, "image": img_data})
        result = response.json()
        if result.get("success"):
            return result["data"]["url"]
        else:
            st.error(f"Upload failed: {result}")
            return None
    except Exception as e:
        st.error(f"Image upload error: {e}")
        return None

# =========================
# PAGE CONFIG
# =========================
st.set_page_config(page_title="Advanced LMS Portal", layout="wide")

st.markdown("""
<style>
/* Streamlit/BaseWeb selectbox menus can extend below the viewport in the sidebar.
   Keep the popup visible by making the options panel scroll internally. */
div[data-baseweb="popover"] div[role="listbox"],
div[data-baseweb="popover"] ul[role="listbox"],
div[data-baseweb="popover"] [data-testid="stVirtualDropdown"] {
    max-height: min(420px, 55vh) !important;
    overflow-y: auto !important;
}
div[data-baseweb="popover"] {
    max-height: 70vh !important;
    z-index: 999999 !important;
}
[data-testid="stSidebar"] div[data-baseweb="select"] {
    z-index: 999999 !important;
}
</style>
""", unsafe_allow_html=True)

# =========================
# SUPABASE + GEMINI INIT
# =========================
try:
    SUPABASE_URL = st.secrets["SUPABASE_URL"]
    SUPABASE_KEY = st.secrets["SUPABASE_KEY"]
    supabase = create_client(SUPABASE_URL, SUPABASE_KEY)
except Exception as e:
    st.error("Secrets    . Settings -> Secrets   .")
    st.stop()

try:
    genai.configure(api_key=st.secrets["GEMINI_API_KEY"])
except Exception:
    pass

# =========================
# SESSION STATES
# =========================
defaults = {
    "logged_in": False,
    "role": "",
    "user_id": "",
    "question_index": 0,
    "answers": {},
    "start_exam": False,
    "exam_submitted": False,
    "exam_id": "",
    "exam_title": "",
    "exam_end_time": 0.0,
    "exam_section_end_time": 0.0,
    "exam_sections": [],
    "exam_section_index": 0,
    "exam_proctoring_enabled": False,
    "current_questions": [],
    "completed_ids": None,
    "admin_preview_mode": False,
    "pin_verified": False,
    "pin_setup_mode": False,
    "email_temp": "",
    "user_id_temp": "",
    "role_temp": "",
    "user_page": "My Classes",
    "card_user_page": "Monthly Bill",
    "focus_class_id": "",
    "focus_exam_id": "",
    "question_start_time": {},
    "question_time_log": {},
    "program_run_results": {},
    "program_custom_results": {},
    "program_submissions": {},
    "attendance_marked_date": "",
    "ai_generated_qs": None,
    "last_attempt_id": "",
    "selected_exam_detail_id": "",
    "selected_exam_series_id": "",
    "exam_lock_message": "",
    "view_solutions_attempt_id": "",
    "programming_session_loaded": False,
    "malpractice_reported_keys": set(),
    "suprabhatam_index": 0,
    "suprabhatam_language": "Telugu",
    "hidden_test_access": {},
    "hidden_test_edit_access": {},
    "hidden_test_viewed_cases": {},
}
for key, val in defaults.items():
    if key not in st.session_state:
        st.session_state[key] = val

# =========================
# PERSISTENT LOGIN
# =========================
def persist_browser_login():
    if not st.session_state.get("logged_in") or not st.session_state.get("user_id"):
        return
    uid = json.dumps(str(st.session_state.user_id))
    role = json.dumps(str(st.session_state.role))
    components.html(f"""
    <script>
      try {{
        window.parent.localStorage.setItem("lms_saved_uid", {uid});
        window.parent.localStorage.setItem("lms_saved_role", {role});
      }} catch (e) {{}}
    </script>
    """, height=0)


def restore_browser_login_bridge():
    components.html("""
    <script>
      try {
        const url = new URL(window.parent.location.href);
        const hasLogin = url.searchParams.get("uid") && url.searchParams.get("role");
        const uid = window.parent.localStorage.getItem("lms_saved_uid");
        const role = window.parent.localStorage.getItem("lms_saved_role");
        if (!hasLogin && uid && role) {
          url.searchParams.set("uid", uid);
          url.searchParams.set("role", role);
          window.parent.location.replace(url.toString());
        }
      } catch (e) {}
    </script>
    """, height=0)


def show_logout_redirect():
    st.session_state.logged_in = False
    st.session_state.role = ""
    st.session_state.user_id = ""
    try:
        st.query_params.clear()
    except Exception:
        pass
    components.html("""
    <script>
      try {
        window.parent.localStorage.removeItem("lms_saved_uid");
        window.parent.localStorage.removeItem("lms_saved_role");
        window.parent.location.replace(window.parent.location.pathname);
      } catch (e) {
        window.parent.location.reload();
      }
    </script>
    """, height=0)
    st.info("Logging out...")
    st.stop()


def user_account_status(user_row):
    status = str((user_row or {}).get("approval_status") or "approved").strip().lower()
    return status or "approved"


def is_user_account_approved(user_row):
    if (user_row or {}).get("role") == "admin":
        return True
    return user_account_status(user_row) not in {"pending", "rejected", "blocked"}


def apply_saved_login(saved_uid, saved_role):
    if not saved_uid or not saved_role:
        return False
    try:
        if saved_role == "card_user":
            card_rows = supabase.table("card_users").select("*").eq("id", saved_uid).execute().data
            if card_rows and card_rows[0].get("active", True):
                st.session_state.logged_in = True
                st.session_state.role = "card_user"
                st.session_state.user_id = saved_uid
                st.session_state.pin_verified = True
                return True
        else:
            urow_data = supabase.table("users").select("*").eq("id", saved_uid).execute().data
            if urow_data and urow_data[0]["role"] == saved_role and is_user_account_approved(urow_data[0]):
                st.session_state.logged_in = True
                st.session_state.role = saved_role
                st.session_state.user_id = saved_uid
                st.session_state.pin_verified = True
                return True
    except Exception:
        pass
    return False


if not st.session_state.logged_in:
    try:
        qp = st.query_params
        saved_uid = qp.get("uid", "")
        saved_role = qp.get("role", "")
        if not apply_saved_login(saved_uid, saved_role):
            restore_browser_login_bridge()
    except Exception:
        pass

# =========================
# PROGRAMMING CODE EVALUATOR
# =========================
def evaluate_java_code(user_code, input_data, expected_output):
    if not st.secrets.get("RAPIDAPI_KEY", ""):
        result = run_java_code(user_code, input_data)
        return result.get("stdout", "").strip() == expected_output.strip()
    url = "https://judge0-ce.p.rapidapi.com/submissions?base64_encoded=false&fields=*"
    payload = {"source_code": user_code, "language_id": 62, "stdin": input_data}
    headers = {"x-rapidapi-key": st.secrets.get("RAPIDAPI_KEY", ""), "Content-Type": "application/json"}
    try:
        response = requests.post(url, json=payload, headers=headers).json()
        token = response.get("token")
        if not token:
            return False
        time.sleep(2)
        result = requests.get(
            f"https://judge0-ce.p.rapidapi.com/submissions/{token}?base64_encoded=false&fields=*",
            headers=headers
        ).json()
        return result.get("stdout", "").strip() == expected_output.strip()
    except Exception:
        return False

PROGRAMMING_LANGUAGES = {
    "java": {
        "label": "Java",
        "wandbox_compiler": "openjdk-jdk-21+35",
        "wandbox_options": "--release=8",
        "file_name": "Main.java",
        "judge0_id": 62,
        "code_language": "java",
        "default_code": """import java.util.*;

public class Main {
    public static void main(String[] args) {
        Scanner sc = new Scanner(System.in);
        String name = sc.hasNextLine() ? sc.nextLine() : \"Student\";
        System.out.println(\"Hello, \" + name + \"!\");
    }
}""",
    },
    "c": {
        "label": "C",
        "wandbox_compiler": "gcc-13.2.0-c",
        "wandbox_options": "",
        "file_name": "main.c",
        "judge0_id": 50,
        "code_language": "c",
        "default_code": """#include <stdio.h>

int main() {
    char name[100] = \"Student\";
    if (fgets(name, sizeof(name), stdin)) {
        for (int i = 0; name[i] != '\\0'; i++) {
            if (name[i] == '\\n') {
                name[i] = '\\0';
                break;
            }
        }
    }
    printf(\"Hello, %s!\\n\", name);
    return 0;
}""",
    },
    "python": {
        "label": "Python",
        "wandbox_compiler": "cpython-3.10.15",
        "wandbox_options": "",
        "file_name": "main.py",
        "judge0_id": 71,
        "code_language": "python",
        "default_code": """try:
    name = input().strip()
except EOFError:
    name = \"Student\"
print(f\"Hello, {name or 'Student'}!\")""",
    },
}

PROGRAMMING_LANGUAGE_LABELS = {meta["label"]: key for key, meta in PROGRAMMING_LANGUAGES.items()}


def normalize_programming_language(language):
    lang = str(language or "java").strip().lower()
    if lang in ("py", "python3"):
        return "python"
    if lang in ("c", "gcc"):
        return "c"
    return lang if lang in PROGRAMMING_LANGUAGES else "java"


def get_programming_language_meta(language):
    return PROGRAMMING_LANGUAGES[normalize_programming_language(language)]


def compiler_service_unavailable(text):
    raw = str(text or "")
    busy_markers = [
        "Resource temporarily unavailable",
        "OCI runtime error",
        "Too Many Requests",
        "rate limit",
        "temporarily unavailable",
        "Service Unavailable",
        "Bad Gateway",
        "Gateway Timeout",
    ]
    return any(marker.lower() in raw.lower() for marker in busy_markers)


def normalize_compiler_stderr(stderr):
    warning_lines = [
        "warning: [options] source value 8 is obsolete",
        "warning: [options] target value 8 is obsolete",
        "warning: [options] To suppress warnings about obsolete options",
    ]
    return "\n".join(
        line for line in str(stderr or "").splitlines()
        if not any(line.startswith(w) for w in warning_lines) and line.strip() != "3 warnings"
    )


def run_code_with_wandbox(user_code, input_data="", language="java"):
    meta = get_programming_language_meta(language)
    code = user_code
    if normalize_programming_language(language) == "java":
        code = user_code.replace("public class Main", "class Main")
    wandbox_payload = {
        "compiler": meta["wandbox_compiler"],
        "code": code,
        "stdin": input_data or "",
    }
    if meta.get("wandbox_options"):
        wandbox_payload["compiler-option-raw"] = meta["wandbox_options"]
    result = requests.post("https://wandbox.org/api/compile.json", json=wandbox_payload, timeout=40).json()
    stdout = result.get("program_output") or ""
    status_code = str(result.get("status", ""))
    stderr = (
        result.get("compiler_error")
        or result.get("compiler_message")
        or result.get("program_error")
        or result.get("program_message")
        or result.get("message")
        or ""
    )
    stderr = normalize_compiler_stderr(stderr)
    label = meta["label"]
    return {
        "ok": status_code == "0",
        "status": f"Accepted ({label})" if status_code == "0" else f"Error ({label})",
        "stdout": stdout,
        "stderr": stderr,
        "time": None,
        "memory": None,
    }


def run_code_with_judge0_public(user_code, input_data="", language="java"):
    meta = get_programming_language_meta(language)
    url = "https://ce.judge0.com/submissions?base64_encoded=false&fields=*"
    payload = {"source_code": user_code, "language_id": meta["judge0_id"], "stdin": input_data or ""}
    headers = {"Content-Type": "application/json"}
    response = requests.post(url, json=payload, headers=headers, timeout=30).json()
    token = response.get("token")
    if not token:
        return {"ok": False, "status": "Submission failed", "stdout": "", "stderr": str(response), "time": None, "memory": None}
    time.sleep(2)
    result = requests.get(
        f"https://ce.judge0.com/submissions/{token}?base64_encoded=false&fields=*",
        headers=headers,
        timeout=30
    ).json()
    status = (result.get("status") or {}).get("description", "Unknown")
    stdout = result.get("stdout") or ""
    stderr = normalize_compiler_stderr(result.get("stderr") or result.get("compile_output") or result.get("message") or "")
    return {
        "ok": status == "Accepted",
        "status": f"{status} (Judge0 public)",
        "stdout": stdout,
        "stderr": stderr,
        "time": result.get("time"),
        "memory": result.get("memory"),
    }


def run_programming_code(user_code, input_data="", language="java"):
    lang = normalize_programming_language(language)
    if lang == "java":
        return run_java_code(user_code, input_data)

    meta = get_programming_language_meta(lang)
    rapidapi_key = st.secrets.get("RAPIDAPI_KEY", "")
    if not rapidapi_key:
        try:
            result = run_code_with_wandbox(user_code, input_data, lang)
            if result.get("ok") or not compiler_service_unavailable(result.get("stderr")):
                return result
        except Exception as e:
            result = {"stderr": str(e)}
        try:
            fallback = run_code_with_judge0_public(user_code, input_data, lang)
            if fallback.get("stderr"):
                fallback["stderr"] = ("Wandbox busy/error kabatti Judge0 public lo run chesanu.\n" + fallback["stderr"]).strip()
            return fallback
        except Exception as e:
            return {"ok": False, "status": f"{meta['label']} API Error", "stdout": "", "stderr": str(e), "time": None, "memory": None}

    url = "https://judge0-ce.p.rapidapi.com/submissions?base64_encoded=false&fields=*"
    payload = {"source_code": user_code, "language_id": meta["judge0_id"], "stdin": input_data or ""}
    headers = {"x-rapidapi-key": rapidapi_key, "Content-Type": "application/json"}
    try:
        response = requests.post(url, json=payload, headers=headers, timeout=30).json()
        token = response.get("token")
        if not token:
            return {"ok": False, "status": "Submission failed", "stdout": "", "stderr": str(response), "time": None, "memory": None}
        time.sleep(2)
        result = requests.get(
            f"https://judge0-ce.p.rapidapi.com/submissions/{token}?base64_encoded=false&fields=*",
            headers=headers,
            timeout=30
        ).json()
        status = (result.get("status") or {}).get("description", "Unknown")
        stdout = result.get("stdout") or ""
        stderr = normalize_compiler_stderr(result.get("stderr") or result.get("compile_output") or result.get("message") or "")
        return {
            "ok": status == "Accepted",
            "status": status,
            "stdout": stdout,
            "stderr": stderr,
            "time": result.get("time"),
            "memory": result.get("memory"),
        }
    except Exception as e:
        return {"ok": False, "status": "API Error", "stdout": "", "stderr": str(e), "time": None, "memory": None}

def run_java_code(user_code, input_data=""):
    rapidapi_key = st.secrets.get("RAPIDAPI_KEY", "")
    if not rapidapi_key:
        def run_with_judge0_public():
            return run_code_with_judge0_public(user_code, input_data, "java")

        java8_code = user_code.replace("public class Main", "class Main")
        wandbox_payload = {
            "compiler": "openjdk-jdk-21+35",
            "compiler-option-raw": "--release=8",
            "code": java8_code,
            "stdin": input_data or "",
        }
        try:
            result = {}
            for attempt in range(2):
                result = requests.post("https://wandbox.org/api/compile.json", json=wandbox_payload, timeout=40).json()
                raw_error = " ".join([
                    str(result.get("compiler_error") or ""),
                    str(result.get("program_error") or ""),
                    str(result.get("compiler_message") or ""),
                    str(result.get("program_message") or ""),
                ])
                if not compiler_service_unavailable(raw_error):
                    break
                if attempt == 0:
                    time.sleep(1)
            raw_error = " ".join([
                str(result.get("compiler_error") or ""),
                str(result.get("program_error") or ""),
                str(result.get("compiler_message") or ""),
                str(result.get("program_message") or ""),
            ])
            if compiler_service_unavailable(raw_error):
                fallback = run_with_judge0_public()
                fallback["status"] = f"{fallback.get('status', 'Unknown')} - Wandbox busy fallback"
                return fallback
            stdout = result.get("program_output") or ""
            status_code = str(result.get("status", ""))
            stderr = ""
            if status_code != "0":
                stderr = (
                    result.get("compiler_error")
                    or result.get("compiler_message")
                    or result.get("program_error")
                    or result.get("program_message")
                    or ""
                )
                stderr = normalize_compiler_stderr(stderr)
            status = "Accepted (Java 8 compatible)" if status_code == "0" else "Error"
            return {
                "ok": status_code == "0",
                "status": status,
                "stdout": stdout,
                "stderr": stderr,
                "time": None,
                "memory": None,
            }
        except Exception as e:
            try:
                fallback = run_with_judge0_public()
                fallback["status"] = f"{fallback.get('status', 'Unknown')} - Wandbox API fallback"
                return fallback
            except Exception:
                return {"ok": False, "status": "Java API Error", "stdout": "", "stderr": str(e)}

    url = "https://judge0-ce.p.rapidapi.com/submissions?base64_encoded=false&fields=*"
    payload = {"source_code": user_code, "language_id": 62, "stdin": input_data or ""}
    headers = {"x-rapidapi-key": rapidapi_key, "Content-Type": "application/json"}
    try:
        response = requests.post(url, json=payload, headers=headers, timeout=30).json()
        token = response.get("token")
        if not token:
            return {"ok": False, "status": "Submission failed", "stdout": "", "stderr": str(response)}
        time.sleep(2)
        result = requests.get(
            f"https://judge0-ce.p.rapidapi.com/submissions/{token}?base64_encoded=false&fields=*",
            headers=headers,
            timeout=30
        ).json()
        status = (result.get("status") or {}).get("description", "Unknown")
        stdout = result.get("stdout") or ""
        stderr = result.get("stderr") or result.get("compile_output") or result.get("message") or ""
        return {
            "ok": status == "Accepted",
            "status": status,
            "stdout": stdout,
            "stderr": stderr,
            "time": result.get("time"),
            "memory": result.get("memory"),
        }
    except Exception as e:
        return {"ok": False, "status": "API Error", "stdout": "", "stderr": str(e)}

PROGRAMMING_META_PREFIX = "__PROGRAMMING_META__"

def make_programming_meta(description, test_cases, language="java"):
    return PROGRAMMING_META_PREFIX + json.dumps({
        "description": description or "",
        "test_cases": test_cases or [],
        "language": normalize_programming_language(language),
    }, ensure_ascii=False)

def get_programming_meta(question):
    raw = question.get("explanation") or ""
    if isinstance(raw, str) and raw.startswith(PROGRAMMING_META_PREFIX):
        try:
            data = json.loads(raw[len(PROGRAMMING_META_PREFIX):])
            data["test_cases"] = data.get("test_cases") or []
            data["language"] = normalize_programming_language(data.get("language"))
            for tc in data["test_cases"]:
                tc["hidden"] = bool(tc.get("hidden", False))
            return data
        except Exception:
            return {"description": "", "test_cases": [], "language": "java"}
    return {"description": raw if question.get("type") == "programming" else "", "test_cases": [], "language": "java"}


def parse_csv_bool(value):
    return str(value or "").strip().lower() in {"1", "true", "yes", "y", "hidden"}


def parse_programming_csv_cases(row):
    raw_json = str(row.get("test_cases_json", "") or row.get("test_cases", "") or "").strip()
    if raw_json:
        try:
            parsed = json.loads(raw_json)
            if isinstance(parsed, list):
                return [
                    {
                        "input": str(tc.get("input", "")),
                        "expected_output": str(tc.get("expected_output", tc.get("output", ""))),
                        "marks": int(tc.get("marks", 1) or 1),
                        "hidden": bool(tc.get("hidden", False)),
                    }
                    for tc in parsed
                    if isinstance(tc, dict)
                ]
        except Exception:
            pass

    cases = []
    for idx in range(1, 11):
        inp = str(row.get(f"test_input_{idx}", "") or "")
        out = str(row.get(f"test_output_{idx}", "") or row.get(f"expected_output_{idx}", "") or "")
        if not inp and not out:
            continue
        try:
            marks = int(row.get(f"test_marks_{idx}", 1) or 1)
        except Exception:
            marks = 1
        cases.append({
            "input": inp,
            "expected_output": out,
            "marks": marks,
            "hidden": parse_csv_bool(row.get(f"test_hidden_{idx}", False)),
        })
    return cases


def build_question_payload_from_csv_row(row, exam_id):
    q_type = str(row.get("type", "mcq") or "mcq").strip().lower()
    if q_type not in {"mcq", "blank", "programming"}:
        q_type = "mcq"
    exp_val = str(row.get("explanation", "") or "").strip()
    correct_answer = str(row.get("correct_answer", "") or "").strip()
    if q_type == "programming":
        language = normalize_programming_language(row.get("language", row.get("programming_language", "java")))
        description = str(row.get("description", "") or row.get("programming_description", "") or exp_val).strip()
        test_cases = parse_programming_csv_cases(row)
        exp_val = make_programming_meta(description, test_cases, language)
        correct_answer = "AUTO"
    return {
        "exam_id": exam_id,
        "question": str(row.get("question", "") or "").strip(),
        "type": q_type,
        "option_a": str(row.get("option_a", "") or ""),
        "option_b": str(row.get("option_b", "") or ""),
        "option_c": str(row.get("option_c", "") or ""),
        "option_d": str(row.get("option_d", "") or ""),
        "correct_answer": correct_answer,
        "hint": str(row.get("hint", "") or ""),
        "image_url": str(row.get("image_url", "") or "").strip() or None,
        "explanation": exp_val if exp_val else None,
    }

def enable_textarea_tab_support():
    st.components.v1.html(
        """
        <script>
        const attachTabs = () => {
            const doc = window.parent.document;
            doc.querySelectorAll('textarea').forEach((ta) => {
                if (ta.dataset.tabInsertReady === '1') return;
                ta.dataset.tabInsertReady = '1';
                ta.addEventListener('keydown', (event) => {
                    if (event.key !== 'Tab') return;
                    event.preventDefault();
                    event.stopPropagation();
                    const start = ta.selectionStart;
                    const end = ta.selectionEnd;
                    ta.setRangeText('    ', start, end, 'end');
                    ta.dispatchEvent(new Event('input', { bubbles: true }));
                });
            });
        };
        attachTabs();
        setInterval(attachTabs, 1000);
        </script>
        """,
        height=0,
    )

def inject_programming_exam_shell(is_prog_exam=True):
    if not is_prog_exam:
        return
    st.markdown(
        """
        <style>
        header, footer, [data-testid="stSidebar"], [data-testid="stToolbar"], [data-testid="stDecoration"] {
            display: none !important;
        }
        .main .block-container {
            max-width: 100vw !important;
            padding: 0.55rem 1rem 0.8rem !important;
        }
        div[data-testid="column"]:nth-of-type(2) {
            position: sticky;
            top: 0.55rem;
            align-self: flex-start;
        }
        .exam-topbar {
            display: flex;
            align-items: center;
            justify-content: space-between;
            gap: 1rem;
            border-bottom: 1px solid #d7dde8;
            padding: 0.35rem 0 0.6rem;
            margin-bottom: 0.55rem;
            background: #ffffff;
        }
        .exam-title {
            color: #172033;
            font-size: 1.05rem;
            font-weight: 750;
        }
        .exam-status {
            display: flex;
            align-items: center;
            gap: 0.65rem;
            color: #526071;
            font-size: 0.88rem;
        }
        .lock-pill {
            border: 1px solid #b9c5d8;
            border-radius: 999px;
            padding: 0.25rem 0.65rem;
            color: #26384f;
            background: #f7f9fc;
            font-weight: 650;
        }
        .problem-pane {
            color: #182235;
            line-height: 1.55;
        }
        .problem-pane h3 {
            font-size: 1.08rem;
            margin: 0.2rem 0 0.7rem;
        }
        .sample-case {
            border: 1px solid #dbe2ef;
            border-radius: 8px;
            padding: 0.75rem;
            margin: 0.65rem 0;
            background: #fbfcff;
        }
        .vscode-shell {
            border: 1px solid #263243;
            border-radius: 8px;
            overflow: hidden;
            background: #0f1724;
            box-shadow: 0 10px 28px rgba(15, 23, 36, 0.18);
            margin-bottom: -0.45rem;
        }
        .vscode-titlebar {
            display: flex;
            align-items: center;
            justify-content: space-between;
            gap: 0.75rem;
            background: #1f2937;
            color: #cbd5e1;
            padding: 0.45rem 0.7rem;
            font-size: 0.78rem;
            border-bottom: 1px solid #334155;
        }
        .vscode-dots {
            display: flex;
            gap: 0.35rem;
        }
        .vscode-dots span {
            width: 10px;
            height: 10px;
            border-radius: 999px;
            display: inline-block;
        }
        .vscode-dots span:nth-child(1) { background: #ef4444; }
        .vscode-dots span:nth-child(2) { background: #f59e0b; }
        .vscode-dots span:nth-child(3) { background: #22c55e; }
        .vscode-file {
            color: #e2e8f0;
            font-weight: 650;
        }
        .vscode-hints {
            color: #94a3b8;
        }
        .ide-code-frame {
            position: relative !important;
            background: #10141d !important;
            border: 1px solid #2d3748 !important;
            border-radius: 0 0 8px 8px !important;
            overflow: hidden !important;
            min-height: 395px !important;
        }
        .ide-code-frame:focus-within {
            border-color: #4f8cff !important;
            box-shadow: 0 0 0 1px #4f8cff !important;
        }
        .ide-line-gutter {
            position: absolute;
            inset: 0 auto 0 0;
            width: 54px;
            padding: 12px 8px 12px 0;
            overflow: hidden;
            background: #0b1020;
            border-right: 1px solid #334155;
            color: #64748b;
            font-family: Consolas, "Cascadia Code", "Courier New", monospace;
            font-size: 14px;
            line-height: 1.55;
            text-align: right;
            user-select: none;
            white-space: pre;
            z-index: 1;
        }
        .ide-highlight-layer {
            position: absolute;
            inset: 0 0 0 54px;
            margin: 0;
            padding: 12px 14px;
            overflow: hidden;
            color: #edf2ff;
            background: transparent;
            font-family: Consolas, "Cascadia Code", "Courier New", monospace;
            font-size: 14px;
            line-height: 1.55;
            tab-size: 4;
            white-space: pre-wrap;
            word-break: normal;
            overflow-wrap: normal;
            pointer-events: none;
            z-index: 1;
        }
        .ide-highlight-layer code {
            font: inherit;
            color: inherit;
            white-space: inherit;
        }
        .ide-token-keyword { color: #c586c0; }
        .ide-token-type { color: #4ec9b0; }
        .ide-token-string { color: #ce9178; }
        .ide-token-number { color: #b5cea8; }
        .ide-token-comment { color: #6a9955; }
        .ide-token-function { color: #dcdcaa; }
        .ide-wrap textarea {
            position: relative !important;
            z-index: 2 !important;
            font-family: Consolas, "Cascadia Code", "Courier New", monospace !important;
            font-size: 14px !important;
            line-height: 1.55 !important;
            tab-size: 4;
            background: transparent !important;
            color: transparent !important;
            -webkit-text-fill-color: transparent !important;
            border: 0 !important;
            border-radius: 0 !important;
            padding: 12px 14px 12px 68px !important;
            caret-color: #70e1ff !important;
            resize: vertical !important;
            min-height: 395px !important;
            white-space: pre-wrap !important;
            overflow-wrap: normal !important;
        }
        .ide-wrap textarea::selection {
            background: rgba(79, 140, 255, 0.35) !important;
            -webkit-text-fill-color: transparent !important;
        }
        .ide-wrap textarea:focus {
            border-color: transparent !important;
            box-shadow: none !important;
            outline: none !important;
        }
        .console-card {
            border: 1px solid #d9e1ef;
            border-radius: 8px;
            padding: 0.75rem;
            background: #ffffff;
        }
        .console-log-pass {
            border-left: 4px solid #1f9d55;
            background: #f0fff4;
            padding: 0.55rem 0.7rem;
            border-radius: 6px;
            margin: 0.45rem 0;
        }
        .console-log-fail {
            border-left: 4px solid #d64545;
            background: #fff5f5;
            padding: 0.55rem 0.7rem;
            border-radius: 6px;
            margin: 0.45rem 0;
        }
        div.stButton > button[kind="primary"] {
            border-radius: 8px !important;
            font-weight: 750 !important;
        }
        </style>
        """,
        unsafe_allow_html=True,
    )



def inject_vscode_editor_styles():
    st.markdown(
        """
        <style>
        .vscode-shell {
            border: 1px solid #263243;
            border-radius: 8px;
            overflow: hidden;
            background: #0f1724;
            box-shadow: 0 10px 28px rgba(15, 23, 36, 0.18);
            margin-bottom: -0.45rem;
        }
        .vscode-titlebar {
            display: flex;
            align-items: center;
            justify-content: space-between;
            gap: 0.75rem;
            background: #1f2937;
            color: #cbd5e1;
            padding: 0.45rem 0.7rem;
            font-size: 0.78rem;
            border-bottom: 1px solid #334155;
        }
        .vscode-dots { display: flex; gap: 0.35rem; }
        .vscode-dots span {
            width: 10px;
            height: 10px;
            border-radius: 999px;
            display: inline-block;
        }
        .vscode-dots span:nth-child(1) { background: #ef4444; }
        .vscode-dots span:nth-child(2) { background: #f59e0b; }
        .vscode-dots span:nth-child(3) { background: #22c55e; }
        .vscode-file { color: #e2e8f0; font-weight: 650; }
        .vscode-hints { color: #94a3b8; }
        .ide-code-frame {
            position: relative !important;
            background: #10141d !important;
            border: 1px solid #2d3748 !important;
            border-radius: 0 0 8px 8px !important;
            overflow: hidden !important;
            min-height: 360px !important;
        }
        .ide-code-frame:focus-within {
            border-color: #4f8cff !important;
            box-shadow: 0 0 0 1px #4f8cff !important;
        }
        .ide-line-gutter {
            position: absolute;
            inset: 0 auto 0 0;
            width: 54px;
            padding: 12px 8px 12px 0;
            overflow: hidden;
            background: #0b1020;
            border-right: 1px solid #334155;
            color: #64748b;
            font-family: Consolas, "Cascadia Code", "Courier New", monospace;
            font-size: 14px;
            line-height: 1.55;
            text-align: right;
            user-select: none;
            white-space: pre;
            z-index: 1;
        }
        .ide-highlight-layer {
            position: absolute;
            inset: 0 0 0 54px;
            margin: 0;
            padding: 12px 14px;
            overflow: hidden;
            color: #edf2ff;
            background: transparent;
            font-family: Consolas, "Cascadia Code", "Courier New", monospace;
            font-size: 14px;
            line-height: 1.55;
            tab-size: 4;
            white-space: pre-wrap;
            word-break: normal;
            overflow-wrap: normal;
            pointer-events: none;
            z-index: 1;
        }
        .ide-highlight-layer code { font: inherit; color: inherit; white-space: inherit; }
        .ide-token-keyword { color: #c586c0; }
        .ide-token-type { color: #4ec9b0; }
        .ide-token-string { color: #ce9178; }
        .ide-token-number { color: #b5cea8; }
        .ide-token-comment { color: #6a9955; }
        .ide-token-function { color: #dcdcaa; }
        .ide-wrap textarea {
            position: relative !important;
            z-index: 2 !important;
            font-family: Consolas, "Cascadia Code", "Courier New", monospace !important;
            font-size: 14px !important;
            line-height: 1.55 !important;
            tab-size: 4;
            background: transparent !important;
            color: transparent !important;
            -webkit-text-fill-color: transparent !important;
            border: 0 !important;
            border-radius: 0 !important;
            padding: 12px 14px 12px 68px !important;
            caret-color: #70e1ff !important;
            resize: vertical !important;
            white-space: pre-wrap !important;
            overflow-wrap: normal !important;
        }
        .ide-wrap textarea::selection {
            background: rgba(79, 140, 255, 0.35) !important;
            -webkit-text-fill-color: transparent !important;
        }
        .ide-wrap textarea:focus {
            border-color: transparent !important;
            box-shadow: none !important;
            outline: none !important;
        }
        </style>
        """,
        unsafe_allow_html=True,
    )
def enable_fullscreen_exam_lock():
    components.html(
        """
        <script>
        (function () {
            const parentWindow = window.parent;
            const doc = parentWindow.document;
            if (parentWindow.__examLockReady) return;
            parentWindow.__examLockReady = true;
            let armedAt = Date.now();
            let reporting = false;

            function report(reason) {
                if (reporting || Date.now() - armedAt < 1200) return;
                reporting = true;
                try {
                    const url = new URL(parentWindow.location.href);
                    url.searchParams.set("malpractice", "1");
                    url.searchParams.set("mal_reason", reason);
                    parentWindow.location.href = url.toString();
                } catch (e) {
                    reporting = false;
                }
            }

            function requestFullScreen() {
                const root = doc.documentElement;
                if (doc.fullscreenElement) return;
                const fn = root.requestFullscreen || root.webkitRequestFullscreen || root.msRequestFullscreen;
                if (fn) {
                    try { fn.call(root).catch(function () {}); } catch (e) {}
                }
            }

            function lockHistory() {
                try {
                    parentWindow.history.pushState({ examLocked: true }, "", parentWindow.location.href);
                    parentWindow.addEventListener("popstate", function () {
                        parentWindow.history.pushState({ examLocked: true }, "", parentWindow.location.href);
                        report("Browser back/navigation attempted");
                    });
                } catch (e) {}
            }

            requestFullScreen();
            lockHistory();
            ["click", "keydown", "pointerdown"].forEach(function (evt) {
                doc.addEventListener(evt, requestFullScreen, true);
            });
            doc.addEventListener("fullscreenchange", function () {
                if (!doc.fullscreenElement) report("Fullscreen exited");
            });
            doc.addEventListener("visibilitychange", function () {
                if (doc.hidden) report("Tab switched or minimized");
            });
            parentWindow.addEventListener("blur", function () {
                report("Exam window lost focus");
            });
            parentWindow.addEventListener("beforeunload", function (event) {
                event.preventDefault();
                event.returnValue = "";
            });
        })();
        </script>
        """,
        height=0,
    )


def enable_ide_textarea_behaviour(question_id, language="java"):
    components.html(
        f"""
        <script>
        (function () {{
            const doc = window.parent.document;
            const key = "ideReady_{question_id}";
            const language = "{normalize_programming_language(language)}";
            function enhance() {{
                const areas = Array.from(doc.querySelectorAll('textarea'));
                const ta = areas.find(el => (el.getAttribute('aria-label') || '').includes('Program')) || areas[0];
                if (!ta || ta.dataset[key] === "1") return;
                ta.dataset[key] = "1";
                const wrapper = ta.closest('[data-testid="stTextArea"]');
                if (wrapper) wrapper.classList.add("ide-wrap");
                const frame = ta.parentElement;
                if (frame) frame.classList.add("ide-code-frame");
                let gutter = frame ? frame.querySelector(".ide-line-gutter") : null;
                if (!gutter && frame) {{
                    gutter = doc.createElement("div");
                    gutter.className = "ide-line-gutter";
                    frame.insertBefore(gutter, ta);
                }}
                let highlight = frame ? frame.querySelector(".ide-highlight-layer") : null;
                if (!highlight && frame) {{
                    highlight = doc.createElement("pre");
                    highlight.className = "ide-highlight-layer";
                    const codeEl = doc.createElement("code");
                    highlight.appendChild(codeEl);
                    frame.insertBefore(highlight, ta);
                }}
                ta.setAttribute("spellcheck", "false");
                ta.setAttribute("autocomplete", "off");
                ta.setAttribute("autocapitalize", "off");

                function escapeHtml(value) {{
                    return value
                        .replace(/&/g, "&amp;")
                        .replace(/</g, "&lt;")
                        .replace(/>/g, "&gt;");
                }}

                function highlightSyntax(source) {{
                    let safe = escapeHtml(source || " ");
                    safe = safe.replace(/(\\/\\/.*?$|#.*?$|\\/\\*[\\s\\S]*?\\*\\/)/gm, '<span class="ide-token-comment">$1</span>');
                    safe = safe.replace(/(&quot;.*?&quot;|'.*?')/g, '<span class="ide-token-string">$1</span>');
                    safe = safe.replace(/\\b(\\d+(?:\\.\\d+)?)\\b/g, '<span class="ide-token-number">$1</span>');
                    safe = safe.replace(/\\b(public|private|protected|class|static|void|int|long|double|float|boolean|char|String|if|else|for|while|do|switch|case|break|continue|return|new|try|catch|finally|import|package|include|using|namespace|def|elif|in|range|print|None|True|False)\\b/g, function(match) {{
                        return /^(int|long|double|float|boolean|char|String)$/.test(match)
                            ? '<span class="ide-token-type">' + match + '</span>'
                            : '<span class="ide-token-keyword">' + match + '</span>';
                    }});
                    safe = safe.replace(/\\b([A-Za-z_][A-Za-z0-9_]*)\\s*(?=\\()/g, '<span class="ide-token-function">$1</span>');
                    return safe;
                }}

                function syncEditorChrome() {{
                    const lines = Math.max(1, ta.value.split("\\n").length);
                    if (gutter) gutter.textContent = Array.from({{ length: lines }}, (_, idx) => idx + 1).join("\\n");
                    if (highlight) {{
                        highlight.firstChild.innerHTML = highlightSyntax(ta.value) + (ta.value.endsWith("\\n") ? "\\n" : "");
                        highlight.scrollTop = ta.scrollTop;
                        highlight.scrollLeft = ta.scrollLeft;
                    }}
                    if (gutter) gutter.scrollTop = ta.scrollTop;
                }}

                function autoPair(open, close) {{
                    const start = ta.selectionStart;
                    const end = ta.selectionEnd;
                    ta.setRangeText(open + ta.value.slice(start, end) + close, start, end, 'end');
                    ta.selectionStart = ta.selectionEnd = start + 1;
                    ta.dispatchEvent(new Event('input', {{ bubbles: true }}));
                }}

                function insertSnippet(text, backOffset) {{
                    const start = ta.selectionStart;
                    const end = ta.selectionEnd;
                    ta.setRangeText(text, start, end, "end");
                    const pos = start + text.length - (backOffset || 0);
                    ta.selectionStart = ta.selectionEnd = pos;
                    ta.dispatchEvent(new Event('input', {{ bubbles: true }}));
                }}

                const snippets = {{
                    java: {{
                        "main": "public static void main(String[] args) {{\\n    \\n}}",
                        "sout": "System.out.println();",
                        "fori": "for (int i = 0; i < n; i++) {{\\n    \\n}}"
                    }},
                    c: {{
                        "main": "int main() {{\\n    \\n    return 0;\\n}}",
                        "printf": "printf(\"%d\\\\n\", );",
                        "fori": "for (int i = 0; i < n; i++) {{\\n    \\n}}"
                    }},
                    python: {{
                        "main": "def main():\\n    \\n\\nif __name__ == \"__main__\":\\n    main()",
                        "fori": "for i in range(n):\\n    ",
                        "print": "print()"
                    }}
                }};

                function currentWord() {{
                    const start = ta.selectionStart;
                    const prefix = ta.value.slice(0, start);
                    const match = prefix.match(/[A-Za-z_][A-Za-z0-9_]*$/);
                    return match ? match[0] : "";
                }}

                ta.addEventListener("input", syncEditorChrome);
                ta.addEventListener("scroll", syncEditorChrome);
                syncEditorChrome();

                ta.addEventListener('keydown', function (event) {{
                    const pairs = {{ "(": ")", "[": "]", "{{": "}}", '"': '"', "'": "'" }};
                    if (event.key === "Tab") {{
                        event.preventDefault();
                        const word = currentWord();
                        const langSnippets = snippets[language] || {{}};
                        if (word && langSnippets[word]) {{
                            const start = ta.selectionStart - word.length;
                            const end = ta.selectionStart;
                            ta.setRangeText(langSnippets[word], start, end, "end");
                            ta.dispatchEvent(new Event('input', {{ bubbles: true }}));
                            return;
                        }}
                        const start = ta.selectionStart;
                        const end = ta.selectionEnd;
                        ta.setRangeText("    ", start, end, "end");
                        ta.dispatchEvent(new Event('input', {{ bubbles: true }}));
                    }} else if (pairs[event.key]) {{
                        event.preventDefault();
                        autoPair(event.key, pairs[event.key]);
                    }} else if (event.key === "Enter") {{
                        const before = ta.value.slice(0, ta.selectionStart).split("\\n").pop() || "";
                        const indent = (before.match(/^\\s+/) || [""])[0] + (/\\{{\\s*$/.test(before) ? "    " : "");
                        if (indent) {{
                            event.preventDefault();
                            const start = ta.selectionStart;
                            const end = ta.selectionEnd;
                            ta.setRangeText("\\n" + indent, start, end, "end");
                            ta.dispatchEvent(new Event('input', {{ bubbles: true }}));
                        }}
                    }}
                }});
            }}
            enhance();
            setTimeout(enhance, 500);
            setInterval(enhance, 1200);
        }})();
        </script>
        """,
        height=0,
    )

def inject_student_home_styles():
    st.markdown(
        """
        <style>
        .student-home-hero {
            border: 1px solid #dbe3ef;
            border-radius: 8px;
            padding: 1rem 1.1rem;
            background: linear-gradient(135deg, #ffffff 0%, #f6f9ff 58%, #f7fff9 100%);
            margin-bottom: 0.8rem;
        }
        .student-home-hero h2 {
            margin: 0 0 0.25rem;
            color: #142033;
            font-size: 1.35rem;
        }
        .student-home-hero p {
            margin: 0;
            color: #5b6778;
            font-size: 0.92rem;
        }
        .home-stat-row {
            display: grid;
            grid-template-columns: repeat(4, minmax(0, 1fr));
            gap: 0.75rem;
            margin: 0.75rem 0 1rem;
        }
        .home-stat {
            border: 1px solid #dde5f1;
            border-radius: 8px;
            padding: 0.85rem;
            background: #ffffff;
        }
        .home-stat small {
            color: #66758a;
            font-weight: 650;
        }
        .home-stat strong {
            display: block;
            color: #162238;
            font-size: 1.2rem;
            margin-top: 0.2rem;
        }
        .module-strip {
            display: flex;
            align-items: center;
            justify-content: space-between;
            gap: 1rem;
            padding: 0.2rem 0 0.55rem;
            border-bottom: 1px solid #edf1f7;
            margin-bottom: 0.65rem;
        }
        .module-strip h3 {
            margin: 0;
            color: #172033;
            font-size: 1.05rem;
        }
        .module-strip span {
            color: #536174;
            font-size: 0.86rem;
            font-weight: 650;
        }
        .class-card {
            border: 1px solid #e0e7f2;
            border-radius: 8px;
            padding: 0.85rem;
            background: #ffffff;
            margin: 0.55rem 0 0.75rem;
        }
        .class-title {
            color: #172033;
            font-weight: 750;
            font-size: 1rem;
            margin-bottom: 0.25rem;
        }
        .exam-chip {
            display: inline-flex;
            align-items: center;
            border: 1px solid #cbd7ea;
            border-radius: 999px;
            padding: 0.18rem 0.55rem;
            background: #f8fbff;
            color: #24364f;
            font-size: 0.82rem;
            font-weight: 650;
            margin: 0.2rem 0 0.45rem;
        }
        .student-exams-page {
            padding-top: 0.65rem;
        }
        .student-exams-page [data-testid="stHorizontalBlock"]:has(.folder-filter-anchor) {
            justify-content: center;
            gap: 0.75rem;
            margin: 0 auto 2.6rem;
            max-width: 520px;
        }
        .student-exams-page [data-testid="stHorizontalBlock"]:has(.folder-filter-anchor) [data-testid="column"] {
            width: auto !important;
            flex: 0 0 auto !important;
            min-width: 0 !important;
        }
        .student-exams-page [data-testid="stHorizontalBlock"]:has(.folder-filter-anchor) button {
            min-width: 70px;
            height: 46px;
            border-radius: 4px;
            border: 1px solid #d6dce6;
            background: #fff;
            color: #1b2433;
            font-weight: 800;
            letter-spacing: 0;
            box-shadow: none;
        }
        .student-exams-page [data-testid="stHorizontalBlock"]:has(.folder-filter-anchor) button[kind="primary"] {
            border-color: #315dff;
            color: #2458ff;
            background: #ffffff;
        }
        .student-exam-card-wrap div[data-testid="stVerticalBlockBorderWrapper"]:has(.student-exam-card-content) {
            min-height: 190px;
            border: 1px solid #d9dfe8;
            border-radius: 4px;
            background: #fbfcff;
            padding: 34px 30px 26px;
            box-shadow: 0 10px 24px rgba(18, 31, 53, 0.035);
            display: flex;
            flex-direction: column;
            justify-content: space-between;
        }
        .student-exam-tags {
            display: flex;
            align-items: center;
            gap: 6px;
            margin-bottom: 18px;
            flex-wrap: wrap;
        }
        .student-exam-tag {
            display: inline-flex;
            align-items: center;
            height: 27px;
            padding: 0 8px;
            border-radius: 999px;
            font-size: 0.91rem;
            line-height: 1;
            font-weight: 800;
            letter-spacing: 0;
        }
        .student-exam-tag.category {
            color: #526d95;
            background: #dce7f5;
            border: 1px solid #c6d3e7;
        }
        .student-exam-tag.free {
            color: #20c997;
            background: #cff8ed;
            border: 1px solid #a7eadb;
        }
        .student-exam-title {
            color: #20242c;
            font-size: 1.55rem;
            line-height: 1.18;
            font-weight: 800;
            letter-spacing: 0;
            min-height: 58px;
            margin: 0 0 22px;
        }
        .student-folder-title {
            color: #172033;
            font-size: 1.82rem;
            line-height: 1.15;
            font-weight: 850;
            letter-spacing: 0;
            min-height: 92px;
            display: flex;
            align-items: center;
            margin: 0 0 22px;
        }
        .student-exam-tag.blocked {
            color: #9f1239;
            background: #ffe4e6;
            border: 1px solid #fecdd3;
        }
        .student-exam-card-wrap button {
            height: 47px;
            border-radius: 4px;
            border: 1px solid #111827;
            background: #fbfcff;
            color: #172033;
            font-size: 1.18rem;
            font-weight: 500;
            letter-spacing: 0;
            box-shadow: none;
        }
        .student-exam-card-wrap button:hover {
            border-color: #111827;
            color: #111827;
            background: #ffffff;
        }
        .student-exam-card-wrap button:disabled,
        .student-exam-card-wrap button:disabled:hover {
            border-color: #d1d5db;
            color: #9f1239;
            background: #fff1f2;
            opacity: 1;
        }
        @media (max-width: 1100px) {
            .student-exam-card-wrap div[data-testid="stVerticalBlockBorderWrapper"]:has(.student-exam-card-content) { padding: 26px 22px 22px; }
            .student-exam-title { font-size: 1.28rem; }
            .student-folder-title { font-size: 1.55rem; min-height: 76px; }
        }
        @media (max-width: 800px) {
            .home-stat-row { grid-template-columns: repeat(2, minmax(0, 1fr)); }
            .student-exams-page [data-testid="stHorizontalBlock"]:has(.folder-filter-anchor) { margin-bottom: 1.4rem; }
            .student-exam-card-wrap div[data-testid="stVerticalBlockBorderWrapper"]:has(.student-exam-card-content) { min-height: 175px; }
        }
        </style>
        """,
        unsafe_allow_html=True,
    )

def get_question_max_marks(question):
    if question.get("type") == "programming":
        meta = get_programming_meta(question)
        total = 0
        for tc in meta.get("test_cases", []):
            try:
                total += int(tc.get("marks", 0) or 0)
            except Exception:
                pass
        return total if total > 0 else 1
    return 1

def get_exam_max_marks(questions):
    return sum(get_question_max_marks(q) for q in questions)

def run_programming_test_cases(question, code, language=None):
    meta = get_programming_meta(question)
    selected_language = normalize_programming_language(language or meta.get("language", "java"))
    results = []
    total_marks = get_question_max_marks(question)
    earned = 0
    for idx, tc in enumerate(meta.get("test_cases", []), start=1):
        inp = tc.get("input", "")
        expected = str(tc.get("expected_output", "")).strip()
        try:
            marks = int(tc.get("marks", 0) or 0)
        except Exception:
            marks = 0
        result = run_programming_code(code, inp, selected_language)
        actual = str(result.get("stdout", "")).strip()
        passed = result.get("ok") and actual == expected
        if passed:
            earned += marks
        results.append({
            "case": idx,
            "input": inp,
            "expected_output": expected,
            "actual_output": actual,
            "marks": marks,
            "hidden": bool(tc.get("hidden", False)),
            "passed": bool(passed),
            "status": result.get("status", "Unknown"),
            "error": result.get("stderr", ""),
        })
    pct = int((earned / total_marks) * 100) if total_marks else 0
    return {"earned": earned, "total": total_marks, "percentage": pct, "language": selected_language, "results": results}

def normalize_id_list(value):
    if value is None:
        return []
    if isinstance(value, list):
        return [str(v) for v in value if str(v).strip()]
    raw = str(value or "").strip()
    if not raw:
        return []
    try:
        parsed = json.loads(raw)
        if isinstance(parsed, list):
            return [str(v) for v in parsed if str(v).strip()]
    except Exception:
        pass
    return [item.strip() for item in raw.split(",") if item.strip()]

def user_account_status(user_row):
    status = str((user_row or {}).get("approval_status") or "approved").strip().lower()
    return status or "approved"

def is_user_account_approved(user_row):
    if (user_row or {}).get("role") == "admin":
        return True
    return user_account_status(user_row) not in {"pending", "rejected", "blocked"}

def generate_numeric_pin(length=6):
    return str(uuid.uuid4().int)[-length:].zfill(length)

def get_hidden_test_acl():
    try:
        admins = supabase.table("users").select("*").eq("role", "admin").execute().data or []
    except Exception:
        return {}
    for row in admins:
        if row.get("hidden_test_view_pin") or row.get("hidden_test_edit_pin"):
            return row
    return admins[0] if admins else {}

def current_user_hidden_access(question_id, mode="view"):
    if st.session_state.get("role") == "admin":
        return True
    qid = str(question_id)
    state_key = "hidden_test_edit_access" if mode == "edit" else "hidden_test_access"
    if st.session_state.get(state_key, {}).get(qid):
        return True
    acl = get_hidden_test_acl()
    user_id = str(st.session_state.get("user_id") or "")
    ids_key = "hidden_test_edit_user_ids" if mode == "edit" else "hidden_test_view_user_ids"
    allowed_ids = normalize_id_list(acl.get(ids_key))
    return bool(user_id and user_id in allowed_ids and st.session_state.get(state_key, {}).get("__global__"))

def verify_hidden_test_pin(pin, mode="view"):
    acl = get_hidden_test_acl()
    pin_key = "hidden_test_edit_pin" if mode == "edit" else "hidden_test_view_pin"
    allowed_ids_key = "hidden_test_edit_user_ids" if mode == "edit" else "hidden_test_view_user_ids"
    expected_pin = str(acl.get(pin_key) or "").strip()
    allowed_ids = normalize_id_list(acl.get(allowed_ids_key))
    user_id = str(st.session_state.get("user_id") or "")
    if not expected_pin:
        return False, "Admin inka PIN generate cheyyaledu."
    if user_id not in allowed_ids:
        return False, "Mee account ki hidden test access ivvaledu."
    if str(pin or "").strip() != expected_pin:
        return False, "Wrong PIN."
    return True, ""

def current_user_can_request_hidden_access(mode="view"):
    if st.session_state.get("role") == "admin":
        return True
    acl = get_hidden_test_acl()
    ids_key = "hidden_test_edit_user_ids" if mode == "edit" else "hidden_test_view_user_ids"
    return str(st.session_state.get("user_id") or "") in normalize_id_list(acl.get(ids_key))

def get_hidden_case_view_limit(user_id=None):
    user_id = str(user_id or st.session_state.get("user_id") or "")
    if not user_id:
        return 0
    try:
        rows = supabase.table("users").select("hidden_test_view_limit").eq("id", user_id).limit(1).execute().data or []
        return int(rows[0].get("hidden_test_view_limit") or 0) if rows else 0
    except Exception:
        return 0

def get_hidden_case_views_used(user_id=None):
    user_id = str(user_id or st.session_state.get("user_id") or "")
    if not user_id:
        return 0
    try:
        rows = supabase.table("hidden_test_case_views").select("id").eq("user_id", user_id).execute().data or []
        return len(rows)
    except Exception:
        return 0

def get_hidden_case_views_remaining(user_id=None):
    return max(0, get_hidden_case_view_limit(user_id) - get_hidden_case_views_used(user_id))

def hidden_case_view_key(question_id, case_no):
    return f"{question_id}:{case_no}"

def has_hidden_case_viewed(question_id, case_no, user_id=None):
    if st.session_state.get("role") == "admin":
        return True
    user_id = str(user_id or st.session_state.get("user_id") or "")
    key = hidden_case_view_key(question_id, case_no)
    if st.session_state.hidden_test_viewed_cases.get(key):
        return True
    try:
        rows = supabase.table("hidden_test_case_views").select("id").eq("user_id", user_id).eq("question_id", str(question_id)).eq("case_no", int(case_no)).limit(1).execute().data or []
        if rows:
            st.session_state.hidden_test_viewed_cases[key] = True
            return True
    except Exception:
        pass
    return False

def consume_hidden_case_view(question_id, case_no):
    if st.session_state.get("role") == "admin":
        return True, ""
    if has_hidden_case_viewed(question_id, case_no):
        return True, ""
    remaining = get_hidden_case_views_remaining()
    if remaining <= 0:
        return False, "Hidden test case view limit ayipoyindi. Admin limit increase cheyyali."
    try:
        supabase.table("hidden_test_case_views").insert({
            "user_id": str(st.session_state.user_id),
            "question_id": str(question_id),
            "case_no": int(case_no),
        }).execute()
        st.session_state.hidden_test_viewed_cases[hidden_case_view_key(question_id, case_no)] = True
        return True, ""
    except Exception as e:
        return False, f"View record save avvaledu: {e}"

def get_user_question_hint(question_id):
    try:
        rows = supabase.table("hint_requests").select("*").eq("user_id", str(st.session_state.user_id)).eq("question_id", str(question_id)).order("created_at", desc=True).limit(1).execute().data or []
    except Exception:
        return None
    return rows[0] if rows else None

def submit_hint_request(exam_id, question_id):
    existing = get_user_question_hint(question_id)
    if existing and existing.get("status") == "pending":
        return False, "Hint request already pending lo undi."
    try:
        supabase.table("hint_requests").insert({
            "user_id": str(st.session_state.user_id),
            "exam_id": str(exam_id),
            "question_id": str(question_id),
            "status": "pending",
        }).execute()
        return True, "Hint request admin ki sent."
    except Exception as e:
        return False, f"Hint request failed: {e}"

def submit_test_case_report(exam_id, question_id, case_no, report_text, result_payload=None):
    try:
        supabase.table("test_case_reports").insert({
            "user_id": str(st.session_state.user_id),
            "exam_id": str(exam_id),
            "question_id": str(question_id),
            "case_no": int(case_no),
            "report_text": str(report_text or "").strip(),
            "result_payload": result_payload or {},
            "status": "pending",
        }).execute()
        return True, "Report admin ki sent."
    except Exception as e:
        return False, f"Report failed: {e}"

def generate_programming_questions_ppt(questions, exam_title):
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    import io

    def rgb(h):
        h = h.lstrip("#")
        return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

    def txt(sl, text, x, y, w, h, sz=13, bold=False, color="1A1A2E", align=PP_ALIGN.LEFT):
        tb = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        r = p.add_run()
        r.text = str(text or "")
        r.font.size = Pt(sz)
        r.font.bold = bold
        r.font.color.rgb = rgb(color)
        return tb

    def box(sl, x, y, w, h, fill, border="D0D8E8"):
        s = sl.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
        s.fill.solid()
        s.fill.fore_color.rgb = rgb(fill)
        s.line.color.rgb = rgb(border)
        return s

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    blank = prs.slide_layouts[6]

    title_slide = prs.slides.add_slide(blank)
    box(title_slide, 0, 0, 10, 5.625, "1E2761", "1E2761")
    txt(title_slide, exam_title or "Programming Questions", 0.6, 1.4, 8.8, 1.0, sz=32, bold=True, color="FFFFFF", align=PP_ALIGN.CENTER)
    txt(title_slide, f"{len(questions)} programming questions with all test cases", 0.8, 2.65, 8.4, 0.5, sz=18, color="CADCFC", align=PP_ALIGN.CENTER)

    for q_index, q in enumerate(questions, start=1):
        meta = get_programming_meta(q)
        test_cases = meta.get("test_cases", []) or []
        overview = prs.slides.add_slide(blank)
        box(overview, 0, 0, 10, 5.625, "F4F6FB", "F4F6FB")
        txt(overview, f"Q{q_index}. {q.get('question', '')}", 0.35, 0.2, 9.3, 0.55, sz=18, bold=True)
        txt(overview, f"Language: {get_programming_language_meta(meta.get('language', 'java'))['label']}  |  Marks: {get_question_max_marks(q)}", 0.35, 0.82, 9.2, 0.3, sz=11, color="4B587C")
        box(overview, 0.35, 1.22, 9.3, 3.85, "FFFFFF")
        txt(overview, meta.get("description", ""), 0.55, 1.38, 8.9, 3.45, sz=12)
        txt(overview, f"{q_index}/{len(questions)}", 8.7, 5.25, 0.9, 0.2, sz=9, color="7F8C8D", align=PP_ALIGN.RIGHT)

        for tc_idx, tc in enumerate(test_cases, start=1):
            slide = prs.slides.add_slide(blank)
            box(slide, 0, 0, 10, 5.625, "F4F6FB", "F4F6FB")
            badge = "Hidden" if tc.get("hidden") else "Sample"
            txt(slide, f"Q{q_index} - Test Case {tc_idx} ({badge})", 0.35, 0.2, 9.3, 0.45, sz=18, bold=True)
            txt(slide, f"Marks: {tc.get('marks', 1)}", 0.35, 0.7, 9, 0.25, sz=11, color="4B587C")
            box(slide, 0.35, 1.08, 4.5, 3.95, "FFFFFF")
            box(slide, 5.15, 1.08, 4.5, 3.95, "FFFFFF")
            txt(slide, "Input", 0.55, 1.23, 4.1, 0.25, sz=12, bold=True, color="1E2761")
            txt(slide, tc.get("input", ""), 0.55, 1.58, 4.1, 3.1, sz=11)
            txt(slide, "Expected Output", 5.35, 1.23, 4.1, 0.25, sz=12, bold=True, color="1E2761")
            txt(slide, tc.get("expected_output", ""), 5.35, 1.58, 4.1, 3.1, sz=11)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()

def get_account_access_sql():
    return """
alter table users add column if not exists approval_status text not null default 'approved';
alter table users add column if not exists username text;
alter table users add column if not exists requested_pin text;
alter table users add column if not exists hidden_test_view_pin text;
alter table users add column if not exists hidden_test_edit_pin text;
alter table users add column if not exists hidden_test_view_user_ids jsonb not null default '[]'::jsonb;
alter table users add column if not exists hidden_test_edit_user_ids jsonb not null default '[]'::jsonb;
alter table users add column if not exists hidden_test_view_limit integer not null default 0;

create table if not exists hidden_test_case_views (
  id uuid primary key default gen_random_uuid(),
  user_id uuid not null references users(id) on delete cascade,
  question_id uuid not null references questions(id) on delete cascade,
  case_no integer not null,
  viewed_at timestamptz default now(),
  unique(user_id, question_id, case_no)
);

create table if not exists hint_requests (
  id uuid primary key default gen_random_uuid(),
  user_id uuid not null references users(id) on delete cascade,
  exam_id uuid references exams(id) on delete cascade,
  question_id uuid not null references questions(id) on delete cascade,
  status text not null default 'pending',
  admin_hint text,
  created_at timestamptz default now(),
  reviewed_at timestamptz
);

create table if not exists test_case_reports (
  id uuid primary key default gen_random_uuid(),
  user_id uuid not null references users(id) on delete cascade,
  exam_id uuid references exams(id) on delete cascade,
  question_id uuid not null references questions(id) on delete cascade,
  case_no integer not null,
  report_text text,
  result_payload jsonb not null default '{}'::jsonb,
  status text not null default 'pending',
  admin_note text,
  created_at timestamptz default now(),
  reviewed_at timestamptz
);
"""

def is_programming_exam(exam_id):
    try:
        q_rows = supabase.table("questions").select("type").eq("exam_id", exam_id).execute().data or []
        return bool(q_rows) and all(q.get("type") == "programming" for q in q_rows)
    except Exception:
        return False

def get_exam_folders_sql():
    return """
create table if not exists exam_folders (
  id uuid primary key default gen_random_uuid(),
  parent_id uuid references exam_folders(id) on delete cascade,
  title text not null,
  display_order integer not null default 0,
  created_at timestamptz default now()
);

alter table exams
add column if not exists folder_id uuid references exam_folders(id) on delete set null;

insert into exam_folders (title, display_order)
select seed.title, seed.display_order
from (values ('Aptitude', 1), ('Reasoning', 2), ('Programming', 3)) as seed(title, display_order)
where not exists (
  select 1 from exam_folders f
  where f.parent_id is null and lower(f.title) = lower(seed.title)
);
"""

def get_exam_enhancements_sql():
    return """
alter table exams
add column if not exists proctoring_enabled boolean not null default false;

alter table exams
add column if not exists section_config jsonb not null default '[]'::jsonb;
"""

def get_exam_series_sql():
    return """
create table if not exists exam_series (
  id uuid primary key default gen_random_uuid(),
  title text not null,
  description text,
  enabled boolean not null default true,
  rounds jsonb not null default '[]'::jsonb,
  created_at timestamptz default now(),
  updated_at timestamptz default now()
);
"""

def fetch_exam_folders(show_warning=False):
    try:
        return supabase.table("exam_folders").select("*").order("display_order").order("title").execute().data or []
    except Exception as e:
        if show_warning:
            st.warning(f"Exam folders table ready ledu. SQL setup run cheyyandi: {e}")
        return []

def build_exam_folder_options(folders, include_uncategorized=True):
    folder_by_id = {str(folder.get("id")): folder for folder in folders}
    def path_for(folder):
        names = []
        current = folder
        seen = set()
        while current and str(current.get("id")) not in seen:
            seen.add(str(current.get("id")))
            names.append(str(current.get("title") or "Folder"))
            current = folder_by_id.get(str(current.get("parent_id") or ""))
        return " / ".join(reversed(names))
    options = {}
    if include_uncategorized:
        options["No Folder / Uncategorized"] = None
    for folder in sorted(folders, key=path_for):
        options[path_for(folder)] = folder.get("id")
    return options

def create_exam_record(payload):
    try:
        return supabase.table("exams").insert(payload).execute().data
    except Exception as e:
        optional_keys = ["folder_id", "proctoring_enabled", "section_config"]
        if any(key in payload for key in optional_keys):
            fallback = dict(payload)
            for key in optional_keys:
                fallback.pop(key, None)
            st.warning("New exam columns database lo inka levu. Exam create ayyindi, kani folders/proctoring/sections save avvakapovachu. SQL setup run cheyyandi.")
            return supabase.table("exams").insert(fallback).execute().data
        raise e

def update_exam_record(exam_id, payload):
    try:
        return supabase.table("exams").update(payload).eq("id", exam_id).execute()
    except Exception as e:
        optional_keys = ["proctoring_enabled", "section_config"]
        if any(key in payload for key in optional_keys):
            fallback = dict(payload)
            for key in optional_keys:
                fallback.pop(key, None)
            st.warning("Proctoring/section columns database lo inka levu. SQL setup run chesi malli save cheyyandi.")
            if fallback:
                return supabase.table("exams").update(fallback).eq("id", exam_id).execute()
            return None
        raise e

def set_exam_folder(exam_id, folder_id):
    try:
        supabase.table("exams").update({"folder_id": folder_id}).eq("id", exam_id).execute()
        return True
    except Exception as e:
        st.error(f"Exam folder update avvaledu. SQL setup run chesara? {e}")
        return False

def format_duration(seconds):
    seconds = int(seconds or 0)
    hours, rem = divmod(seconds, 3600)
    mins, secs = divmod(rem, 60)
    if hours:
        return f"{hours} hour, {mins} Mins, {secs} Sec"
    if mins:
        return f"{mins} Mins, {secs} Sec"
    return f"{secs} Sec"

def parse_exam_section_config(raw):
    if not raw:
        return []
    if isinstance(raw, str):
        try:
            raw = json.loads(raw)
        except Exception:
            return []
    if not isinstance(raw, list):
        return []
    sections = []
    for idx, item in enumerate(raw, start=1):
        if not isinstance(item, dict):
            continue
        title = str(item.get("title") or f"Section {idx}").strip() or f"Section {idx}"
        try:
            duration = int(item.get("duration_mins") or 0)
        except Exception:
            duration = 0
        try:
            question_count = int(item.get("question_count") or 0)
        except Exception:
            question_count = 0
        if duration > 0 and question_count > 0:
            sections.append({"title": title, "duration_mins": duration, "question_count": question_count})
    return sections

def normalize_exam_sections(exam, total_questions):
    total_questions = int(total_questions or 0)
    sections = parse_exam_section_config(exam.get("section_config"))
    if not sections:
        return [{
            "title": str(exam.get("title") or "Exam"),
            "duration_mins": int(exam.get("duration_mins", 30) or 30),
            "question_count": total_questions,
            "start": 0,
            "end": total_questions,
        }]
    normalized = []
    cursor = 0
    for section in sections:
        if cursor >= total_questions:
            break
        end = min(cursor + int(section["question_count"]), total_questions)
        normalized.append({**section, "start": cursor, "end": end})
        cursor = end
    if cursor < total_questions and normalized:
        normalized[-1]["question_count"] += total_questions - cursor
        normalized[-1]["end"] = total_questions
    elif cursor < total_questions:
        normalized.append({
            "title": "Remaining",
            "duration_mins": int(exam.get("duration_mins", 30) or 30),
            "question_count": total_questions,
            "start": 0,
            "end": total_questions,
        })
    return normalized

def get_current_section():
    sections = st.session_state.get("exam_sections") or []
    if not sections:
        return {"title": st.session_state.get("exam_title", "Exam"), "start": 0, "end": len(st.session_state.current_questions)}
    idx = min(int(st.session_state.get("exam_section_index", 0) or 0), len(sections) - 1)
    return sections[idx]

def move_to_exam_section(index):
    sections = st.session_state.get("exam_sections") or []
    if not sections or index >= len(sections):
        return False
    section = sections[index]
    st.session_state.exam_section_index = index
    st.session_state.question_index = int(section.get("start", 0) or 0)
    st.session_state.exam_section_end_time = time.time() + (int(section.get("duration_mins", 1) or 1) * 60)
    return True

def is_exam_proctoring_enabled(exam):
    return bool(exam.get("proctoring_enabled"))

def parse_exam_series_rounds(raw):
    if not raw:
        return []
    if isinstance(raw, str):
        try:
            raw = json.loads(raw)
        except Exception:
            return []
    if not isinstance(raw, list):
        return []
    rounds = []
    for idx, item in enumerate(raw, start=1):
        if not isinstance(item, dict):
            continue
        exam_id = str(item.get("exam_id") or "").strip()
        if not exam_id:
            continue
        try:
            qualify_pct = int(item.get("qualify_pct", 60) or 60)
        except Exception:
            qualify_pct = 60
        rounds.append({
            "exam_id": exam_id,
            "title": str(item.get("title") or f"Round {idx}").strip() or f"Round {idx}",
            "qualify_pct": max(0, min(100, qualify_pct)),
        })
    return rounds

def fetch_exam_series(show_warning=False):
    try:
        return supabase.table("exam_series").select("*").order("created_at", desc=True).execute().data or []
    except Exception as e:
        if show_warning:
            st.warning(f"Multi round exams table ready ledu. SQL setup run cheyyandi: {e}")
        return []

def get_exam_best_attempt(user_id, exam_id):
    attempts = get_exam_attempt_rows(user_id, exam_id)
    if not attempts:
        return None
    return max(attempts, key=lambda attempt: int(attempt.get("score") or 0))

def get_exam_score_pct(exam, attempt):
    if not attempt:
        return 0
    try:
        questions = supabase.table("questions").select("*").eq("exam_id", exam["id"]).execute().data or []
        total = get_exam_max_marks(questions)
    except Exception:
        total = 0
    score = int(attempt.get("score") or 0)
    return int((score / total) * 100) if total else 0

def get_series_round_states(series, exams, user_id):
    exam_by_id = {str(exam.get("id")): exam for exam in exams}
    rounds = parse_exam_series_rounds(series.get("rounds"))
    states = []
    previous_qualified = True
    for idx, round_cfg in enumerate(rounds, start=1):
        exam = exam_by_id.get(str(round_cfg.get("exam_id")))
        attempt = get_exam_best_attempt(user_id, round_cfg.get("exam_id")) if exam else None
        pct = get_exam_score_pct(exam, attempt) if exam and attempt else 0
        qualified = bool(attempt) and pct >= int(round_cfg.get("qualify_pct", 60) or 60)
        unlocked = previous_qualified and bool(exam) and bool(exam.get("enabled"))
        states.append({
            "index": idx,
            "round": round_cfg,
            "exam": exam,
            "attempt": attempt,
            "pct": pct,
            "qualified": qualified,
            "unlocked": unlocked,
        })
        previous_qualified = qualified
    return states

def get_exam_series_lock_message(exam, user_id):
    series_list = fetch_exam_series(show_warning=False)
    if not series_list:
        return ""
    try:
        exams = supabase.table("exams").select("*").execute().data or []
    except Exception:
        exams = [exam]
    exam_id = str(exam.get("id"))
    for series in series_list:
        if not bool(series.get("enabled")):
            continue
        states = get_series_round_states(series, exams, user_id)
        for state in states:
            if str(state["round"].get("exam_id")) == exam_id and not state["unlocked"]:
                return f"Locked: '{series.get('title', 'Multi Round Exam')}' lo previous round qualify ayyaka ee round open avuthundi."
    return ""

def format_attempt_date(value):
    if not value:
        return "Attempt time not saved"
    text = str(value).replace("T", " ").replace("Z", "")
    if "." in text:
        text = text.split(".", 1)[0]
    return text[:19]

def get_exam_attempt_rows(user_id, exam_id):
    try:
        rows = supabase.table("exam_attempts").select("*").eq("user_id", user_id).eq("exam_id", exam_id).execute().data or []
        return sorted(rows, key=lambda r: str(r.get("created_at") or r.get("submitted_at") or r.get("id") or ""), reverse=True)
    except Exception:
        return []

def get_attempt_time_spent_seconds(attempt_id):
    try:
        rows = supabase.table("user_answers").select("time_spent_seconds").eq("attempt_id", attempt_id).execute().data or []
        return sum(int(row.get("time_spent_seconds") or 0) for row in rows)
    except Exception:
        return 0

def get_attempt_answer_map(attempt_id):
    try:
        rows = supabase.table("user_answers").select("*").eq("attempt_id", attempt_id).execute().data or []
        return {a["question_id"]: a.get("answer", "") for a in rows}
    except Exception:
        return {}

def render_exam_result_summary(exam, questions, attempt, show_return=True):
    max_marks = get_exam_max_marks(questions)
    score = int(attempt.get("score") or 0)
    pct = (score / max_marks * 100) if max_marks else 0
    answer_map = get_attempt_answer_map(attempt["id"])
    attempted = len([q for q in questions if str(answer_map.get(q["id"], "")).strip()])
    wrong = max(max_marks - score, 0)
    sections = normalize_exam_sections(exam, len(questions))
    section_rows = []
    for section in sections:
        section_questions = questions[int(section.get("start", 0) or 0):int(section.get("end", len(questions)) or len(questions))]
        section_attempted = len([q for q in section_questions if str(answer_map.get(q["id"], "")).strip()])
        section_marks = get_exam_max_marks(section_questions)
        section_score = 0
        for q in section_questions:
            saved_answer = answer_map.get(q["id"], "")
            if q.get("type") == "programming":
                try:
                    section_score += int((json.loads(saved_answer) if saved_answer else {}).get("earned", 0) or 0)
                except Exception:
                    pass
            elif q.get("type") == "mcq":
                if check_mcq_correct(saved_answer, q):
                    section_score += 1
            elif str(saved_answer).strip().lower() == str(q.get("correct_answer", "")).strip().lower():
                section_score += 1
        section_pct = (section_score / section_marks * 100) if section_marks else 0
        section_rows.append(
            f"<tr><td><strong>{html.escape(str(section.get('title') or 'Section'))}</strong></td>"
            f"<td>{section_attempted} / {len(section_questions)}</td>"
            f"<td>{section_score} / {max(section_marks - section_score, 0)}</td>"
            f"<td>{section_score:.2f} / {section_marks}</td>"
            f"<td>{section_pct:.2f}</td><td>{section_pct:.2f}</td></tr>"
        )
    section_rows_html = "".join(section_rows)
    time_spent = get_attempt_time_spent_seconds(attempt.get("id"))
    qualifying_pct = 80
    status = "Qualified" if pct >= qualifying_pct else "Not Qualified"
    st.markdown(
        f"""
        <style>
        .result-hero {{display:grid;grid-template-columns:1.05fr .95fr;gap:36px;align-items:center;padding:34px 28px 26px;border:1px solid #e5eaf2;border-radius:8px;background:#fff;}}
        .result-title {{font-size:2rem;font-weight:800;color:#10243d;margin-bottom:22px;}}
        .result-metrics {{display:grid;grid-template-columns:repeat(3,minmax(130px,1fr));gap:28px 42px;}}
        .result-metric small {{display:block;color:#0f172a;font-size:1rem;margin-bottom:8px;}}
        .result-metric strong {{font-size:1.45rem;color:#07111f;line-height:1.25;}}
        .result-art {{min-height:245px;border-radius:8px;background:linear-gradient(135deg,#f7fafc,#eef6ff);display:flex;align-items:center;justify-content:center;color:#1f7a4d;font-size:7rem;font-weight:900;}}
        .result-table table {{width:100%;border-collapse:collapse;margin-top:10px;}}
        .result-table th,.result-table td {{border:1px solid #d8dee8;padding:12px 16px;text-align:left;}}
        .result-table th {{font-size:.85rem;color:#526070;background:#fbfcfe;}}
        @media(max-width:900px){{.result-hero{{grid-template-columns:1fr;}}.result-metrics{{grid-template-columns:1fr 1fr;}}}}
        </style>
        <div class="result-hero">
            <div>
                <div class="result-title">{html.escape(str(exam.get('title') or st.session_state.exam_title))}</div>
                <div class="result-metrics">
                    <div class="result-metric"><small>Time Spent</small><strong>{format_duration(time_spent)}</strong></div>
                    <div class="result-metric"><small>Marks</small><strong>{score} / {max_marks}</strong></div>
                    <div class="result-metric"><small>Percentage</small><strong>{pct:.2f}</strong></div>
                    <div class="result-metric"><small>Accuracy</small><strong>{pct:.2f}</strong></div>
                    <div class="result-metric"><small>Status</small><strong>{status}</strong></div>
                    <div class="result-metric"><small>Qualifying Percentage</small><strong>{qualifying_pct:.2f}%</strong></div>
                </div>
            </div>
            <div class="result-art">&#10003;</div>
        </div>
        <h3 style="margin-top:34px;">Sectional Summary</h3>
        <div class="result-table">
        <table>
            <thead><tr><th>SECTION</th><th>ATTEMPTED</th><th>CORRECT / WRONG</th><th>MARKS</th><th>PERCENTAGE</th><th>ACCURACY</th></tr></thead>
            <tbody>{section_rows_html}</tbody>
        </table>
        </div>
        """,
        unsafe_allow_html=True,
    )
    c1, c2 = st.columns([1, 4])
    with c1:
        if st.button("View Solutions", type="primary", use_container_width=True, key=f"view_solutions_{attempt['id']}"):
            st.session_state.view_solutions_attempt_id = attempt["id"]
            st.rerun()
    if show_return:
        with c2:
            if st.button("Return to Exams", use_container_width=False):
                save_programming_exam_session(status="submitted", force_submit=False)
                st.session_state.start_exam = False
                st.session_state.exam_submitted = False
                st.session_state.answers = {}
                st.session_state.question_index = 0
                st.session_state.current_questions = []
                st.session_state.selected_exam_detail_id = str(exam.get("id") or "")
                st.rerun()
    if st.session_state.get("view_solutions_attempt_id") == attempt.get("id"):
        exam_data = supabase.table("exams").select("*").eq("id", exam.get("id")).execute().data
        if exam_data and exam_data[0].get("show_answers"):
            st.subheader("Solutions")
            render_review_sheet(questions, answer_map, [attempt])
        else:
            st.info("Solutions visibility admin disable chesaru.")

def render_exam_detail_view(exam):
    st.markdown(f"### {exam.get('title', 'Exam')}")
    questions = supabase.table("questions").select("*").eq("exam_id", exam["id"]).execute().data or []
    attempts = get_exam_attempt_rows(st.session_state.user_id, exam["id"])
    active_session = get_active_programming_session(st.session_state.user_id, exam["id"]) if is_programming_exam(exam["id"]) else None
    tab_details, tab_attempts = st.tabs(["Assessment Details", "Previous Attempts"])
    with tab_details:
        sections = normalize_exam_sections(exam, len(questions))
        c1, c2, c3, c4 = st.columns(4)
        c1.metric("Duration", f"{int(exam.get('duration_mins') or 30)} mins")
        c2.metric("Questions", len(questions))
        c3.metric("Marks", get_exam_max_marks(questions))
        c4.metric("Proctoring", "Camera On" if is_exam_proctoring_enabled(exam) else "Off")
        if len(sections) > 1:
            st.markdown("#### Sections")
            for section in sections:
                st.caption(f"{section['title']} - {section['duration_mins']} mins - Q{section['start'] + 1} to Q{section['end']}")
        has_pwd = exam.get("password") and str(exam["password"]).strip()
        entered_pwd = st.text_input(f"Access Code for {exam['title']}", type="password", key=f"detail_pwd_{exam['id']}") if has_pwd else ""
        if st.button("Start Exam", key=f"detail_start_{exam['id']}", type="primary", use_container_width=True):
            if has_pwd and entered_pwd.strip() != str(exam["password"]).strip():
                st.error("Wrong Password!")
            else:
                start_exam_with_questions(exam, questions)
                st.rerun()
    with tab_attempts:
        cards = []
        if active_session:
            cards.append(("In Progress", active_session, None))
        cards.extend(("Completed", attempt, attempt) for attempt in attempts)
        if not cards:
            st.info("Previous attempts levu.")
        cols = st.columns(3)
        for idx, (status, row, attempt) in enumerate(cards):
            with cols[idx % 3]:
                with st.container(border=True):
                    st.markdown("**Attempted on:**")
                    st.caption(format_attempt_date(row.get("updated_at") or row.get("created_at") or row.get("submitted_at")))
                    st.markdown(f"**Status:** {status}")
                    if status == "In Progress":
                        if st.button("Resume Test", key=f"resume_session_{exam['id']}", use_container_width=True):
                            start_exam_with_questions(exam, questions)
                            st.rerun()
                    else:
                        score = int(attempt.get("score") or 0)
                        total = get_exam_max_marks(questions)
                        st.caption(f"Score: {score}/{total}")
                        if st.button("View Results", key=f"attempt_result_{attempt['id']}", use_container_width=True):
                            st.session_state.exam_id = exam["id"]
                            st.session_state.exam_title = exam["title"]
                            st.session_state.current_questions = questions
                            st.session_state.last_attempt_id = attempt["id"]
                            st.session_state.start_exam = True
                            st.session_state.exam_submitted = True
                            st.rerun()

def render_student_exam_card(exam, folder_label="Programming", key_prefix="exam"):
    title = html.escape(str(exam.get("title") or "Exam"))
    is_enabled = bool(exam.get("enabled"))
    status_tag = '<span class="student-exam-tag free">Free</span>' if is_enabled else '<span class="student-exam-tag blocked">Blocked</span>'
    with st.container(border=True):
        st.markdown(
            f"""
            <div class="student-exam-card-content">
                <div class="student-exam-tags">
                    {status_tag}
                </div>
                <div class="student-exam-title">{title}</div>
            </div>
            """,
            unsafe_allow_html=True,
        )
        if is_enabled:
            if st.button("Open", key=f"{key_prefix}_open_{exam['id']}", use_container_width=True):
                st.session_state.selected_exam_detail_id = str(exam["id"])
                st.rerun()
        else:
            st.button("Blocked", key=f"{key_prefix}_blocked_{exam['id']}", use_container_width=True, disabled=True)

def render_student_folder_card(folder, subfolder_count=0, exam_count=0, key_prefix="folder"):
    title = html.escape(str(folder.get("title") or "Folder"))
    with st.container(border=True):
        st.markdown(
            f"""
            <div class="student-exam-card-content">
                <div class="student-folder-title">{title}</div>
            </div>
            """,
            unsafe_allow_html=True,
        )
        if st.button("Open", key=f"{key_prefix}_open_{folder['id']}", use_container_width=True):
            st.session_state.student_exam_current_folder_id = str(folder["id"])
            st.rerun()

def render_student_exam_series_card(series, exams, user_id, key_prefix="series"):
    title = html.escape(str(series.get("title") or "Multi Round Exam"))
    rounds = parse_exam_series_rounds(series.get("rounds"))
    states = get_series_round_states(series, exams, user_id)
    completed = len([state for state in states if state["qualified"]])
    with st.container(border=True):
        st.markdown(
            f"""
            <div class="student-exam-card-content">
                <div class="student-exam-tags">
                    <span class="student-exam-tag category">Multi Round</span>
                    <span class="student-exam-tag free">{completed}/{len(rounds)} Qualified</span>
                </div>
                <div class="student-exam-title">{title}</div>
            </div>
            """,
            unsafe_allow_html=True,
        )
        desc = str(series.get("description") or "").strip()
        if desc:
            st.caption(desc)
        if st.button("Open Rounds", key=f"{key_prefix}_open_{series['id']}", use_container_width=True):
            st.session_state.selected_exam_series_id = str(series["id"])
            st.rerun()

def render_student_exam_series_detail(series, exams, user_id):
    st.markdown(f"### {series.get('title', 'Multi Round Exam')}")
    desc = str(series.get("description") or "").strip()
    if desc:
        st.caption(desc)
    states = get_series_round_states(series, exams, user_id)
    if not states:
        st.info("Rounds setup avvaledu.")
        return
    for state in states:
        round_cfg = state["round"]
        exam = state["exam"]
        qualify_pct = int(round_cfg.get("qualify_pct", 60) or 60)
        with st.container(border=True):
            c1, c2, c3 = st.columns([2.2, 1, 1])
            with c1:
                st.markdown(f"**Round {state['index']}: {round_cfg.get('title', 'Round')}**")
                st.caption(exam.get("title", "Exam missing") if exam else "Linked exam missing")
            with c2:
                status = "Qualified" if state["qualified"] else ("Open" if state["unlocked"] else "Locked")
                st.metric("Status", status)
            with c3:
                score_text = f"{state['pct']}%" if state["attempt"] else "-"
                st.metric(f"Need {qualify_pct}%", score_text)
            if not exam:
                st.error("Ee round linked exam dorakaledu.")
            elif state["unlocked"]:
                questions = supabase.table("questions").select("*").eq("exam_id", exam["id"]).execute().data or []
                if state["qualified"]:
                    st.success("Ee round qualify ayyaru. Next round unlock ayyindi.")
                has_pwd = exam.get("password") and str(exam["password"]).strip()
                entered_pwd = st.text_input(f"Access Code for {exam['title']}", type="password", key=f"series_pwd_{series['id']}_{exam['id']}") if has_pwd else ""
                btn_label = "Retake Round" if state["attempt"] else "Start Round"
                if st.button(btn_label, key=f"series_start_{series['id']}_{exam['id']}", type="primary", use_container_width=True):
                    if has_pwd and entered_pwd.strip() != str(exam["password"]).strip():
                        st.error("Wrong Password!")
                    else:
                        start_exam_with_questions(exam, questions)
                        st.rerun()
            else:
                st.info("Previous round qualify ayyaka ee round open avuthundi.")

def show_student_exams_tab(user_id):
    folders = fetch_exam_folders(show_warning=True)
    series_list = fetch_exam_series(show_warning=False)
    try:
        exams = supabase.table("exams").select("*").execute().data or []
    except Exception as e:
        st.error(f"Exams load avvaledu: {e}")
        return
    if not exams and not folders and not series_list:
        st.info("Active exams levu.")
        return

    folder_by_id = {str(folder.get("id")): folder for folder in folders}
    series_exam_ids = {
        str(round_cfg.get("exam_id"))
        for series in series_list if bool(series.get("enabled"))
        for round_cfg in parse_exam_series_rounds(series.get("rounds"))
    }
    exams_by_folder = {}
    for exam in exams:
        if str(exam.get("id")) in series_exam_ids:
            continue
        exams_by_folder.setdefault(str(exam.get("folder_id") or ""), []).append(exam)

    def folder_path(folder):
        names = []
        current = folder
        seen = set()
        while current and str(current.get("id")) not in seen:
            seen.add(str(current.get("id")))
            names.append(str(current.get("title") or "Folder"))
            current = folder_by_id.get(str(current.get("parent_id") or ""))
        return " / ".join(reversed(names))

    def sort_folder(folder):
        title = str(folder.get("title") or "")
        programming_first = 0 if title.strip().lower() == "programming" else 1
        return (programming_first, int(folder.get("display_order") or 0), title.lower())

    selected_exam_id = str(st.session_state.get("selected_exam_detail_id") or "")
    if selected_exam_id:
        selected_exam = next((exam for exam in exams if str(exam.get("id")) == selected_exam_id), None)
        if selected_exam:
            if st.button("Back to Exam Folders", key="back_to_exam_folders"):
                st.session_state.selected_exam_detail_id = ""
                st.rerun()
            render_exam_detail_view(selected_exam)
            return
        st.session_state.selected_exam_detail_id = ""

    selected_series_id = str(st.session_state.get("selected_exam_series_id") or "")
    if selected_series_id:
        selected_series = next((series for series in series_list if str(series.get("id")) == selected_series_id), None)
        if selected_series:
            if st.button("Back to Exam Folders", key="back_to_exam_series"):
                st.session_state.selected_exam_series_id = ""
                st.rerun()
            render_student_exam_series_detail(selected_series, exams, user_id)
            return
        st.session_state.selected_exam_series_id = ""

    current_folder_id = str(st.session_state.get("student_exam_current_folder_id") or "")
    if current_folder_id and current_folder_id not in folder_by_id:
        current_folder_id = ""
        st.session_state.student_exam_current_folder_id = ""

    st.markdown('<div class="student-exams-page">', unsafe_allow_html=True)

    enabled_series = [series for series in series_list if bool(series.get("enabled"))]
    if enabled_series and not current_folder_id:
        st.markdown("### Multi Round Exams")
        series_cols = st.columns(3)
        for idx, series in enumerate(enabled_series):
            with series_cols[idx % 3]:
                st.markdown('<div class="student-exam-card-wrap">', unsafe_allow_html=True)
                render_student_exam_series_card(series, exams, user_id, key_prefix=f"student_series_{idx}")
                st.markdown("</div>", unsafe_allow_html=True)
        st.divider()

    if current_folder_id:
        current_folder = folder_by_id[current_folder_id]
        back_label = "Back to Folders"
        if st.button(back_label, key=f"student_folder_back_{current_folder_id}"):
            parent_id = str(current_folder.get("parent_id") or "")
            st.session_state.student_exam_current_folder_id = parent_id
            st.rerun()
        st.markdown(f"### {html.escape(folder_path(current_folder))}")
        child_folders = [folder for folder in folders if str(folder.get("parent_id") or "") == current_folder_id]
        visible_exams = exams_by_folder.get(current_folder_id, []) if not child_folders else []
    else:
        child_folders = [folder for folder in folders if not folder.get("parent_id")]
        visible_exams = exams_by_folder.get("", []) if not child_folders else []

    child_folders = sorted(child_folders, key=sort_folder)
    folder_label_by_id = {str(folder.get("id")): folder_path(folder) for folder in folders}

    if child_folders:
        cols = st.columns(4)
        for idx, folder in enumerate(child_folders):
            folder_id = str(folder.get("id"))
            nested_count = len([child for child in folders if str(child.get("parent_id") or "") == folder_id])
            exam_count = len(exams_by_folder.get(folder_id, []))
            with cols[idx % 4]:
                st.markdown('<div class="student-exam-card-wrap">', unsafe_allow_html=True)
                render_student_folder_card(folder, subfolder_count=nested_count, exam_count=exam_count, key_prefix=f"student_folder_{idx}")
                st.markdown('</div>', unsafe_allow_html=True)
    elif visible_exams:
        cols = st.columns(4)
        for idx, exam in enumerate(visible_exams):
            folder_id = str(exam.get("folder_id") or "")
            folder_label = folder_label_by_id.get(folder_id, "Programming")
            with cols[idx % 4]:
                st.markdown('<div class="student-exam-card-wrap">', unsafe_allow_html=True)
                render_student_exam_card(exam, folder_label=folder_label, key_prefix=f"student_exam_{current_folder_id or 'root'}_{idx}")
                st.markdown('</div>', unsafe_allow_html=True)
    else:
        st.caption("I folder lo inka sub folders/exams levu.")
    st.markdown('</div>', unsafe_allow_html=True)

def show_admin_exam_folders_tab():
    st.subheader("Exam Folders")
    with st.expander("Database SQL setup", expanded=False):
        st.code(get_exam_folders_sql(), language="sql")
        st.code(get_exam_enhancements_sql(), language="sql")
    folders = fetch_exam_folders(show_warning=True)
    folder_options = build_exam_folder_options(folders)
    folder_parent_map = {str(folder.get("id")): str(folder.get("parent_id") or "") for folder in folders}
    def is_descendant_folder(candidate_id, folder_id):
        current = str(candidate_id or "")
        seen = set()
        while current and current not in seen:
            if current == str(folder_id):
                return True
            seen.add(current)
            current = folder_parent_map.get(current, "")
        return False
    col_add, col_move = st.columns([1, 2])
    with col_add:
        st.markdown("#### Create Folder")
        with st.form("create_exam_folder_form", clear_on_submit=True):
            parent_label = st.selectbox("Parent Folder", list(folder_options.keys()), key="new_exam_folder_parent")
            folder_title = st.text_input("Folder Name", placeholder="Aptitude / Reasoning / Programming")
            display_order = st.number_input("Display Order", min_value=0, value=0, step=1)
            if st.form_submit_button("Create Folder", type="primary"):
                if not folder_title.strip():
                    st.error("Folder name enter cheyyandi.")
                else:
                    try:
                        supabase.table("exam_folders").insert({
                            "title": folder_title.strip(),
                            "parent_id": folder_options.get(parent_label),
                            "display_order": int(display_order),
                        }).execute()
                        st.success("Folder create ayyindi.")
                        st.rerun()
                    except Exception as e:
                        st.error(f"Folder create avvaledu. SQL setup run cheyyandi: {e}")
        st.markdown("#### Manage Folders")
        for folder in folders:
            with st.expander(folder.get("title", "Folder")):
                new_title = st.text_input("Title", value=folder.get("title") or "", key=f"exam_folder_title_{folder['id']}")
                new_order = st.number_input("Order", min_value=0, value=int(folder.get("display_order") or 0), step=1, key=f"exam_folder_order_{folder['id']}")
                parent_choices = {
                    label: fid for label, fid in folder_options.items()
                    if str(fid) != str(folder["id"]) and not is_descendant_folder(fid, folder["id"])
                }
                current_parent_label = next((label for label, fid in parent_choices.items() if str(fid) == str(folder.get("parent_id"))), "No Folder / Uncategorized")
                new_parent_label = st.selectbox(
                    "Parent Folder",
                    list(parent_choices.keys()),
                    index=list(parent_choices.keys()).index(current_parent_label) if current_parent_label in parent_choices else 0,
                    key=f"exam_folder_parent_{folder['id']}",
                )
                save_col, del_col = st.columns(2)
                with save_col:
                    if st.button("Save", key=f"exam_folder_save_{folder['id']}", use_container_width=True):
                        supabase.table("exam_folders").update({
                            "title": new_title.strip(),
                            "display_order": int(new_order),
                            "parent_id": parent_choices.get(new_parent_label),
                        }).eq("id", folder["id"]).execute()
                        st.success("Folder update ayyindi."); st.rerun()
                with del_col:
                    if st.button("Delete", key=f"exam_folder_delete_{folder['id']}", use_container_width=True):
                        supabase.table("exam_folders").delete().eq("id", folder["id"]).execute()
                        st.warning("Folder delete ayyindi."); st.rerun()
    with col_move:
        st.markdown("#### Move Existing Exams")
        exams = supabase.table("exams").select("*").execute().data or []
        if not exams:
            st.info("Exams levu.")
            return
        for exam in exams:
            current_folder = exam.get("folder_id")
            labels = list(folder_options.keys())
            current_label = next((label for label, fid in folder_options.items() if str(fid) == str(current_folder)), labels[0])
            with st.container(border=True):
                st.markdown(f"**{exam.get('title', 'Untitled Exam')}**")
                selected = st.selectbox("Move to folder", labels, index=labels.index(current_label), key=f"move_exam_folder_{exam['id']}")
                if st.button("Move Exam", key=f"move_exam_btn_{exam['id']}", use_container_width=True):
                    if set_exam_folder(exam["id"], folder_options.get(selected)):
                        st.success("Exam folder update ayyindi.")
                        st.rerun()

def show_admin_exam_series_tab():
    st.subheader("Multi Round Exam Builder")
    with st.expander("Database SQL setup", expanded=False):
        st.code(get_exam_series_sql(), language="sql")

    try:
        exams = supabase.table("exams").select("*").order("title").execute().data or []
    except Exception as e:
        st.error(f"Exams load avvaledu: {e}")
        return
    series_list = fetch_exam_series(show_warning=True)
    if not exams:
        st.info("First individual exams create cheyyandi. Tarvatha ikkada rounds ga join cheyyachu.")
        return

    exam_options = {f"{exam.get('title', 'Untitled')}  |  {exam.get('id')}": exam for exam in exams}
    st.markdown("#### Create Multi Round Exam")
    title = st.text_input("Series Name", placeholder="Campus Hiring - 3 Rounds", key="series_new_title")
    description = st.text_area("Description", placeholder="Round 1 qualify ayithe Round 2 open avuthundi.", key="series_new_desc", height=80)
    selected_labels = st.multiselect(
        "Rounds order lo individual exams select cheyyandi",
        list(exam_options.keys()),
        key="series_new_rounds",
    )
    rounds = []
    if selected_labels:
        st.caption("Qualifying percentage per round")
        for idx, label in enumerate(selected_labels, start=1):
            exam = exam_options[label]
            qualify_pct = st.number_input(
                f"Round {idx}: {exam.get('title')} qualify %",
                min_value=0,
                max_value=100,
                value=60,
                step=1,
                key=f"series_new_qpct_{idx}_{exam['id']}",
            )
            rounds.append({"exam_id": str(exam["id"]), "title": str(exam.get("title") or f"Round {idx}"), "qualify_pct": int(qualify_pct)})
    enabled = st.checkbox("Enable Multi Round Exam", value=True, key="series_new_enabled")
    if st.button("Create Multi Round Exam", type="primary", use_container_width=True, key="series_create_btn"):
        if not title.strip():
            st.error("Series name enter cheyyandi.")
        elif len(rounds) < 2:
            st.error("At least 2 individual exams select cheyyandi.")
        else:
            try:
                supabase.table("exam_series").insert({
                    "title": title.strip(),
                    "description": description.strip() if description.strip() else None,
                    "enabled": enabled,
                    "rounds": rounds,
                    "updated_at": "now()",
                }).execute()
                st.success("Multi round exam create ayyindi.")
                st.rerun()
            except Exception as e:
                st.error(f"Create avvaledu. SQL setup run chesara? {e}")

    st.divider()
    st.markdown("#### Manage Multi Round Exams")
    if not series_list:
        st.info("Multi round exams levu.")
        return
    for series in series_list:
        with st.expander(series.get("title", "Multi Round Exam")):
            current_rounds = parse_exam_series_rounds(series.get("rounds"))
            new_title = st.text_input("Series Name", value=series.get("title") or "", key=f"series_title_{series['id']}")
            new_desc = st.text_area("Description", value=series.get("description") or "", key=f"series_desc_{series['id']}", height=70)
            new_enabled = st.toggle("Enabled", value=bool(series.get("enabled")), key=f"series_enabled_{series['id']}")
            st.caption("Rounds JSON advanced edit")
            rounds_text = json.dumps(current_rounds, ensure_ascii=False, indent=2)
            new_rounds_text = st.text_area(
                "Rounds",
                value=rounds_text,
                key=f"series_rounds_{series['id']}",
                height=160,
                help='Format: [{"exam_id":"...","title":"Aptitude","qualify_pct":60}]',
            )
            for idx, round_cfg in enumerate(current_rounds, start=1):
                exam = next((ex for ex in exams if str(ex.get("id")) == str(round_cfg.get("exam_id"))), None)
                label = exam.get("title") if exam else "Exam missing"
                st.caption(f"Round {idx}: {label} | Need {round_cfg.get('qualify_pct', 60)}%")
            c_save, c_delete = st.columns(2)
            with c_save:
                if st.button("Save Series", key=f"series_save_{series['id']}", type="primary", use_container_width=True):
                    parsed_rounds = parse_exam_series_rounds(new_rounds_text)
                    if len(parsed_rounds) < 2:
                        st.error("At least 2 valid rounds kavali.")
                    else:
                        supabase.table("exam_series").update({
                            "title": new_title.strip() or series.get("title"),
                            "description": new_desc.strip() if new_desc.strip() else None,
                            "enabled": new_enabled,
                            "rounds": parsed_rounds,
                            "updated_at": "now()",
                        }).eq("id", series["id"]).execute()
                        st.success("Series update ayyindi.")
                        st.rerun()
            with c_delete:
                if st.button("Delete Series", key=f"series_delete_{series['id']}", use_container_width=True):
                    supabase.table("exam_series").delete().eq("id", series["id"]).execute()
                    st.warning("Series delete ayyindi.")
                    st.rerun()

def get_programming_session_sql():
    return """
create table if not exists programming_exam_sessions (
  user_id uuid not null references users(id) on delete cascade,
  exam_id uuid not null references exams(id) on delete cascade,
  answers jsonb not null default '{}'::jsonb,
  question_index integer not null default 0,
  exam_end_time double precision,
  program_submissions jsonb not null default '{}'::jsonb,
  question_time_log jsonb not null default '{}'::jsonb,
  status text not null default 'active',
  force_submit boolean not null default false,
  malpractice_count integer not null default 0,
  last_malpractice_reason text,
  updated_at timestamptz default now(),
  primary key (user_id, exam_id)
);
"""


def get_active_programming_session(user_id, exam_id):
    try:
        rows = supabase.table("programming_exam_sessions").select("*").eq("user_id", str(user_id)).eq("exam_id", str(exam_id)).eq("status", "active").limit(1).execute().data
        return rows[0] if rows else None
    except Exception:
        return None


def save_programming_exam_session(status="active", force_submit=None, malpractice_reason=None):
    if not st.session_state.get("exam_id") or not is_programming_exam(st.session_state.exam_id):
        return
    payload = {
        "user_id": str(st.session_state.user_id),
        "exam_id": str(st.session_state.exam_id),
        "answers": st.session_state.get("answers", {}),
        "question_index": int(st.session_state.get("question_index", 0) or 0),
        "exam_end_time": float(st.session_state.get("exam_end_time", 0) or 0),
        "program_submissions": st.session_state.get("program_submissions", {}),
        "question_time_log": st.session_state.get("question_time_log", {}),
        "status": status,
        "updated_at": "now()",
    }
    if force_submit is not None:
        payload["force_submit"] = bool(force_submit)
    if malpractice_reason:
        payload["last_malpractice_reason"] = malpractice_reason
    try:
        supabase.table("programming_exam_sessions").upsert(payload, on_conflict="user_id,exam_id").execute()
    except Exception:
        pass


def load_programming_exam_session(user_id, exam, questions):
    session = get_active_programming_session(user_id, exam["id"])
    if not session:
        return None
    sections = normalize_exam_sections(exam, len(questions))
    duration = sum(int(section.get("duration_mins", 0) or 0) for section in sections) or int(exam.get("duration_mins", 30) or 30)
    end_time = float(session.get("exam_end_time") or 0)
    if end_time < 0:
        end_time = time.time() + abs(end_time)
    elif end_time <= time.time():
        end_time = time.time() + (duration * 60)
    return {
        "exam_id": exam["id"],
        "exam_title": exam["title"],
        "start_exam": True,
        "exam_submitted": False,
        "answers": session.get("answers") or {},
        "question_index": min(int(session.get("question_index") or 0), max(len(questions) - 1, 0)),
        "current_questions": questions,
        "exam_end_time": end_time,
        "exam_section_end_time": time.time() + (int(sections[0].get("duration_mins", duration) or duration) * 60),
        "exam_sections": sections,
        "exam_section_index": 0,
        "exam_proctoring_enabled": is_exam_proctoring_enabled(exam),
        "program_run_results": {},
        "program_custom_results": {},
        "program_submissions": session.get("program_submissions") or {},
        "question_time_log": session.get("question_time_log") or {},
        "question_start_time": {},
        "programming_session_loaded": True,
    }


def start_exam_with_questions(exam, questions):
    lock_message = get_exam_series_lock_message(exam, st.session_state.user_id)
    if lock_message:
        st.session_state.exam_lock_message = lock_message
        st.error(lock_message)
        return
    resumed = load_programming_exam_session(st.session_state.user_id, exam, questions) if is_programming_exam(exam["id"]) else None
    if resumed:
        st.session_state.update(resumed)
        return
    sections = normalize_exam_sections(exam, len(questions))
    duration = sum(int(section.get("duration_mins", 0) or 0) for section in sections) or int(exam.get("duration_mins", 30) or 30)
    first_section_duration = int(sections[0].get("duration_mins", duration) or duration) if sections else duration
    st.session_state.update({
        "exam_id": exam["id"],
        "exam_title": exam["title"],
        "start_exam": True,
        "exam_submitted": False,
        "answers": {},
        "question_index": 0,
        "current_questions": questions,
        "exam_end_time": time.time() + (duration * 60),
        "exam_section_end_time": time.time() + (first_section_duration * 60),
        "exam_sections": sections,
        "exam_section_index": 0,
        "exam_proctoring_enabled": is_exam_proctoring_enabled(exam),
        "program_run_results": {},
        "program_custom_results": {},
        "program_submissions": {},
        "question_time_log": {},
        "question_start_time": {},
        "programming_session_loaded": True,
    })
    save_programming_exam_session(status="active", force_submit=False)


def report_programming_malpractice(reason):
    if not st.session_state.get("exam_id"):
        return
    is_prog = is_programming_exam(st.session_state.exam_id)
    key = f"{st.session_state.user_id}:{st.session_state.exam_id}:{reason}"
    if key in st.session_state.malpractice_reported_keys:
        return
    st.session_state.malpractice_reported_keys.add(key)
    try:
        uinfo = supabase.table("users").select("name, email").eq("id", st.session_state.user_id).execute().data or [{}]
        uname = uinfo[0].get("name") or uinfo[0].get("email") or "Student"
        send_admin_notification(f"Malpractice alert: {uname} - {st.session_state.exam_title} - {reason}")
        if not is_prog:
            return
        session = get_active_programming_session(st.session_state.user_id, st.session_state.exam_id) or {}
        count = int(session.get("malpractice_count") or 0) + 1
        save_programming_exam_session(status="active", malpractice_reason=reason)
        supabase.table("programming_exam_sessions").update({
            "malpractice_count": count,
            "last_malpractice_reason": reason,
            "updated_at": "now()",
        }).eq("user_id", str(st.session_state.user_id)).eq("exam_id", str(st.session_state.exam_id)).execute()
    except Exception:
        pass

def render_live_proctoring_panel(enabled=True):
    if not enabled:
        return
    st.markdown("### Live Proctoring")
    components.html(
        """
        <div style="border:1px solid #d8dee8;border-radius:8px;padding:10px;background:#ffffff;font-family:system-ui,Segoe UI,Arial,sans-serif;">
          <video id="proctorVideo" autoplay playsinline muted
            style="width:100%;aspect-ratio:4/3;background:#111827;border-radius:6px;object-fit:cover;"></video>
          <div id="proctorStatus" style="margin-top:8px;font-size:13px;color:#475569;font-weight:600;">
            Camera permission waiting...
          </div>
        </div>
        <script>
        const statusEl = document.getElementById("proctorStatus");
        const videoEl = document.getElementById("proctorVideo");
        async function startProctorCamera() {
          try {
            const stream = await navigator.mediaDevices.getUserMedia({video: true, audio: false});
            videoEl.srcObject = stream;
            statusEl.textContent = "Camera live";
            statusEl.style.color = "#047857";
          } catch (err) {
            statusEl.textContent = "Camera access blocked. Please allow camera permission and refresh.";
            statusEl.style.color = "#b91c1c";
            try {
              const url = new URL(window.parent.location.href);
              url.searchParams.set("malpractice", "1");
              url.searchParams.set("mal_reason", "Camera permission blocked");
              window.parent.location.href = url.toString();
            } catch(e) {}
          }
        }
        if (navigator.mediaDevices && navigator.mediaDevices.getUserMedia) {
          startProctorCamera();
        } else {
          statusEl.textContent = "Camera not supported in this browser.";
          statusEl.style.color = "#b91c1c";
        }
        </script>
        """,
        height=270,
    )

def user_attempted_question(user_id, question_id):
    try:
        attempts = supabase.table("exam_attempts").select("id").eq("user_id", user_id).execute().data or []
        attempt_ids = [a["id"] for a in attempts]
        if not attempt_ids:
            return False
        answers = supabase.table("user_answers").select("id").eq("question_id", question_id).in_("attempt_id", attempt_ids).limit(1).execute().data
        return bool(answers)
    except Exception:
        return False

def parse_program_answer(value, default_language="java"):
    if isinstance(value, dict):
        return str(value.get("code", "")), normalize_programming_language(value.get("language", default_language))
    return str(value or ""), normalize_programming_language(default_language)


def get_cached_program_submission(question_id, code, language="java"):
    saved = st.session_state.program_submissions.get(str(question_id), {})
    if (
        saved.get("code") == code
        and normalize_programming_language(saved.get("language", language)) == normalize_programming_language(language)
        and saved.get("score_data")
    ):
        return saved["score_data"]
    return None

def get_unsubmitted_program_questions(questions, answers):
    pending = []
    for idx, q in enumerate(questions, start=1):
        if q.get("type") != "programming":
            continue
        qid = q["id"]
        meta = get_programming_meta(q)
        code, language = parse_program_answer(answers.get(qid, ""), meta.get("language", "java"))
        if not get_cached_program_submission(qid, code, language):
            pending.append(idx)
    return pending

def score_exam_answers(questions, answers, use_cached_programming=True):
    total_marks = get_exam_max_marks(questions)
    earned_marks = 0
    answer_payloads = {}
    for q in questions:
        qid = q["id"]
        user_val = answers.get(qid, "")
        if q.get("type") == "programming":
            meta = get_programming_meta(q)
            code, language = parse_program_answer(user_val, meta.get("language", "java"))
            score_data = get_cached_program_submission(qid, code, language) if use_cached_programming else None
            if score_data is None and not use_cached_programming:
                score_data = run_programming_test_cases(q, code, language)
            if score_data is None:
                score_data = {
                    "earned": 0,
                    "total": get_question_max_marks(q),
                    "percentage": 0,
                    "results": [],
                    "note": "Program was not individually submitted before final exam submit.",
                }
            earned_marks += int(score_data.get("earned", 0) or 0)
            answer_payloads[qid] = json.dumps({"code": code, "language": language, **score_data}, ensure_ascii=False)
        elif q.get("type") == "mcq":
            if check_mcq_correct(user_val, q):
                earned_marks += 1
            answer_payloads[qid] = user_val
        else:
            if str(user_val).strip().lower() == str(q.get("correct_answer", "")).strip().lower():
                earned_marks += 1
            answer_payloads[qid] = user_val
    percentage = int((earned_marks / total_marks) * 100) if total_marks else 0
    return earned_marks, total_marks, percentage, answer_payloads

def compact_answer_for_save(answer, max_chars=3500):
    text = answer if isinstance(answer, str) else str(answer)
    if len(text) <= max_chars:
        return text
    try:
        data = json.loads(text)
        if isinstance(data, dict) and "code" in data:
            code = str(data.get("code", ""))
            earned = data.get("earned", 0)
            total = data.get("total", 0)
            return json.dumps({
                "code": code[:max_chars],
                "earned": earned,
                "total": total,
                "language": data.get("language", "java"),
                "note": "Detailed test-case output was trimmed while saving.",
            }, ensure_ascii=False)
    except Exception:
        pass
    return text[:max_chars] + "\n...trimmed while saving..."

def insert_user_answer_row(attempt_id, question_id, answer, time_spent_seconds=None):
    base_payload = {"attempt_id": attempt_id, "question_id": question_id, "answer": answer}
    payloads = []
    if time_spent_seconds is not None:
        payloads.append({**base_payload, "time_spent_seconds": int(time_spent_seconds or 0)})
    payloads.append(base_payload)
    payloads.append({**base_payload, "answer": compact_answer_for_save(answer)})

    last_error = None
    for payload in payloads:
        try:
            supabase.table("user_answers").insert(payload).execute()
            return
        except Exception as e:
            last_error = e
    raise last_error

def insert_exam_attempt_row(payload):
    payloads = [{**payload, "created_at": "now()"}, payload]
    last_error = None
    for item in payloads:
        try:
            supabase.table("exam_attempts").insert(item).execute()
            return
        except Exception as e:
            last_error = e
    raise last_error

def submit_exam_attempt(questions, include_time=False, require_programming_submitted=False):
    if require_programming_submitted:
        pending = get_unsubmitted_program_questions(questions, st.session_state.answers)
        if pending:
            nums = ", ".join(str(n) for n in pending)
            raise ValueError(f"Please click Submit Program for question(s): {nums}")
    attempt_uuid = str(uuid.uuid4())
    final_score, total_marks, final_percentage, answer_payloads = score_exam_answers(questions, st.session_state.answers, use_cached_programming=True)
    insert_exam_attempt_row({
        "id": attempt_uuid,
        "user_id": st.session_state.user_id,
        "exam_id": st.session_state.exam_id,
        "score": final_score,
    })
    for q in questions:
        t_spent = st.session_state.question_time_log.get(q["id"], 0) if include_time else None
        insert_user_answer_row(attempt_uuid, q["id"], answer_payloads.get(q["id"], ""), t_spent)
    st.session_state.last_attempt_id = attempt_uuid
    if is_programming_exam(st.session_state.exam_id):
        save_programming_exam_session(status="submitted", force_submit=False)
    return attempt_uuid

def get_answer_time_spent(attempt_id, question_id):
    try:
        ua_data = supabase.table("user_answers").select("time_spent_seconds") \
            .eq("attempt_id", attempt_id).eq("question_id", question_id).execute().data
        return ua_data[0]["time_spent_seconds"] if ua_data and ua_data[0].get("time_spent_seconds") else 0
    except Exception:
        return 0

# =========================
# OCR: IMAGE  QUESTION EXTRACT
# =========================
def extract_question_from_image(image_source):
    """Image nundi OCR.space API use chesi question + options extract cheyyali"""
    import re

    def parse_ocr_text(raw_text):
        # Normalize line endings and clean up
        lines = [l.strip() for l in raw_text.replace("\r\n", "\n").replace("\r", "\n").splitlines()]
        lines = [l for l in lines if l]

        options = {"A": "", "B": "", "C": "", "D": ""}
        question_lines = []
        option_found_at = None

        # Pattern: line starts with A) / A. / A: / (A) / 1) / 1. etc.
        opt_line_pat = re.compile(
            r"^(?:option\s*)?([A-Da-d]|[1-4])\s*[\)\]\.:\-]\s*(.+)", re.IGNORECASE
        )

        for i, line in enumerate(lines):
            m = opt_line_pat.match(line)
            if m:
                lbl = m.group(1).upper()
                if lbl in ["1","2","3","4"]:
                    lbl = chr(ord("A") + int(lbl) - 1)
                val = m.group(2).strip()
                if lbl in options and not options[lbl]:
                    options[lbl] = val
                    if option_found_at is None:
                        option_found_at = i
            else:
                # Only add to question if we haven't hit options yet
                if option_found_at is None:
                    # Skip "Answer:" lines
                    if not re.match(r"(?i)^(answer|correct\s*answer|ans)\s*[:\-]", line):
                        question_lines.append(line)

        question = " ".join(question_lines).strip()

        # Try to detect correct answer from text
        answer = ""
        full_text = "\n".join(lines)
        ans_m = re.search(
            r"(?im)^(?:answer|correct\s*answer|ans)\s*[:\-]\s*([A-D]|[1-4]|.+)$",
            full_text
        )
        if ans_m:
            raw_ans = ans_m.group(1).strip()
            lbl = raw_ans.upper()
            if lbl in ["1","2","3","4"]:
                lbl = chr(ord("A") + int(lbl) - 1)
            if lbl in options:
                answer = lbl  # store as label A/B/C/D
            else:
                answer = raw_ans

        has_options = any(v for v in options.values())
        return {
            "question": question if question else full_text[:300],
            "type": "mcq" if has_options else "blank",
            "option_a": options["A"],
            "option_b": options["B"],
            "option_c": options["C"],
            "option_d": options["D"],
            "correct_answer": answer,
            "hint": "",
        }

    try:
        api_key = st.secrets.get("OCR_SPACE_API_KEY", "")
        if not api_key:
            st.error("OCR_SPACE_API_KEY Streamlit secrets  add .")
            return None
        payload = {"apikey": api_key, "language": "eng", "OCREngine": "2",
                   "isOverlayRequired": "false", "scale": "true"}
        if isinstance(image_source, str):
            resp = requests.post("https://api.ocr.space/parse/image",
                                 data={**payload, "url": image_source}, timeout=40)
        else:
            resp = requests.post("https://api.ocr.space/parse/image", data=payload,
                                 files={"file": (image_source.name, image_source.getvalue(), image_source.type)},
                                 timeout=40)
        result = resp.json()
        if result.get("IsErroredOnProcessing"):
            errs = result.get("ErrorMessage") or result.get("ErrorDetails") or "OCR failed"
            if isinstance(errs, list): errs = " ".join(str(e) for e in errs)
            st.error(f"OCR error: {errs}")
            return None
        parsed = result.get("ParsedResults") or []
        raw_text = "\n".join(p.get("ParsedText","") for p in parsed if p.get("ParsedText")).strip()
        if not raw_text:
            st.error("Image  text clear  .")
            return None
        with st.expander("OCR raw text"):
            st.code(raw_text)
        return parse_ocr_text(raw_text)
    except Exception as e:
        st.error(f"Image question extract failed: {e}")
        return None


def login():
    st.title("LMS Login")
    login_tab, signup_tab = st.tabs(["Login", "Signup"])

    with login_tab:
        login_method = st.radio("Login method:", ["Email & Password", "PIN Login"], horizontal=True, key="login_method_radio")

        if login_method == "Email & Password":
            email = st.text_input("Email", key="login_email")
            password = st.text_input("Password", type="password", key="login_password")
            if st.button("Login", type="primary", use_container_width=True, key="email_login_btn"):
                try:
                    response = supabase.auth.sign_in_with_password({"email": email, "password": password})
                    user = response.user
                    user_data = supabase.table("users").select("*").eq("email", email).execute()
                    if user_data.data:
                        urow = user_data.data[0]
                        if not is_user_account_approved(urow):
                            st.error("Mee account admin approval kosam pending lo undi.")
                            try:
                                supabase.auth.sign_out()
                            except Exception:
                                pass
                            return
                        st.session_state.email_temp = email
                        st.session_state.user_id_temp = user.id
                        st.session_state.role_temp = urow["role"]
                        st.session_state.pin_setup_mode = not bool(urow.get("app_pin"))
                        st.session_state.pin_verified = False
                        st.rerun()
                    else:
                        st.error("User record not found.")
                except Exception as e:
                    st.error(str(e))
        else:
            st.caption("Enter your unique PIN.")
            pin = st.text_input("PIN", type="password", max_chars=6, key="pin_only_input")
            if st.button("Enter App", type="primary", use_container_width=True, key="pin_only_btn"):
                if not pin:
                    st.error("Enter PIN.")
                else:
                    try:
                        user_data = supabase.table("users").select("*").eq("app_pin", pin).execute()
                        if user_data.data:
                            urow = user_data.data[0]
                            if not is_user_account_approved(urow):
                                st.error("Mee account admin approval kosam pending lo undi.")
                                return
                            st.session_state.logged_in = True
                            st.session_state.role = urow["role"]
                            st.session_state.user_id = urow["id"]
                            st.session_state.pin_verified = True
                            st.query_params["uid"] = str(urow["id"])
                            st.query_params["role"] = urow["role"]
                            st.rerun()
                        card_user_data = supabase.table("card_users").select("*").eq("app_pin", pin).execute()
                        if card_user_data.data:
                            urow = card_user_data.data[0]
                            st.session_state.logged_in = True
                            st.session_state.role = "card_user"
                            st.session_state.user_id = urow["id"]
                            st.session_state.pin_verified = True
                            st.query_params["uid"] = str(urow["id"])
                            st.query_params["role"] = "card_user"
                            st.rerun()
                        st.error("Wrong PIN. Try again.")
                    except Exception as e:
                        st.error(str(e))

    with signup_tab:
        st.subheader("Create Account")
        st.caption("Details submit chesaka admin approve chestadu. Approval taruvatha login avvachu.")
        with st.form("signup_request_form", clear_on_submit=True):
            full_name = st.text_input("Name")
            signup_email = st.text_input("Mail")
            signup_username = st.text_input("Username")
            signup_password = st.text_input("Password", type="password")
            requested_pin = st.text_input("PIN for account creation", type="password", max_chars=6)
            submitted = st.form_submit_button("Request Account", type="primary", use_container_width=True)
        if submitted:
            if not all([full_name.strip(), signup_email.strip(), signup_username.strip(), signup_password, requested_pin.strip()]):
                st.error("Name, mail, username, password, PIN anni enter cheyyandi.")
            elif not requested_pin.isdigit() or len(requested_pin) < 4:
                st.error("PIN 4-6 digits undali.")
            else:
                try:
                    existing_pin = supabase.table("users").select("id").eq("app_pin", requested_pin).execute().data
                    existing_req_pin = supabase.table("users").select("id").eq("requested_pin", requested_pin).execute().data
                    if existing_pin or existing_req_pin:
                        st.error("Ee PIN already use lo undi. Vere PIN try cheyyandi.")
                    else:
                        auth_resp = supabase.auth.sign_up({"email": signup_email.strip(), "password": signup_password})
                        auth_user = auth_resp.user
                        if not auth_user:
                            st.error("Signup create avvaledu. Email/password check cheyyandi.")
                            return
                        supabase.table("users").insert({
                            "id": auth_user.id,
                            "name": full_name.strip(),
                            "email": signup_email.strip(),
                            "username": signup_username.strip(),
                            "role": "user",
                            "approval_status": "pending",
                            "requested_pin": requested_pin.strip(),
                            "app_pin": None,
                        }).execute()
                        st.success("Account request sent. Admin approve chesaka login cheyyandi.")
                except Exception as e:
                    st.error(f"Signup failed: {e}")


def pin_screen():
    st.title("App Lock")
    if st.session_state.pin_setup_mode:
        st.subheader("Set PIN (4-6 digits)")
        st.caption("PIN must be globally unique.")
        pin1 = st.text_input("New PIN (4-6 digits)", type="password", max_chars=6, key="pin_new")
        pin2 = st.text_input("Confirm PIN", type="password", max_chars=6, key="pin_confirm")
        if st.button("Set PIN", type="primary", use_container_width=True):
            if not pin1.isdigit() or len(pin1) < 4:
                st.error("Enter 4-6 digits only.")
            elif pin1 != pin2:
                st.error("PINs do not match.")
            else:
                existing = supabase.table("users").select("id").eq("app_pin", pin1).execute().data
                if existing and existing[0]["id"] != st.session_state.user_id_temp:
                    st.error("This PIN is already used by another user.")
                else:
                    supabase.table("users").update({"app_pin": pin1}).eq("id", st.session_state.user_id_temp).execute()
                    st.session_state.logged_in = True
                    st.session_state.role = st.session_state.role_temp
                    st.session_state.user_id = st.session_state.user_id_temp
                    st.session_state.pin_verified = True
                    st.query_params["uid"] = str(st.session_state.user_id_temp)
                    st.query_params["role"] = st.session_state.role_temp
                    st.success("PIN set. Welcome!")
                    st.rerun()
    else:
        st.subheader("Enter PIN")
        pin_input = st.text_input("4-digit PIN", type="password", max_chars=4, key="pin_entry")
        col1, col2 = st.columns(2)
        with col1:
            if st.button("Enter App", type="primary", use_container_width=True):
                urow = supabase.table("users").select("app_pin").eq("id", st.session_state.user_id_temp).execute().data
                if urow and urow[0]["app_pin"] == pin_input:
                    st.session_state.logged_in = True
                    st.session_state.role = st.session_state.role_temp
                    st.session_state.user_id = st.session_state.user_id_temp
                    st.session_state.pin_verified = True
                    st.query_params["uid"] = str(st.session_state.user_id_temp)
                    st.query_params["role"] = st.session_state.role_temp
                    st.rerun()
                else:
                    st.error("Wrong PIN.")
        with col2:
            if st.button("Login with another account", use_container_width=True):
                st.session_state.email_temp = ""
                st.session_state.user_id_temp = ""
                st.session_state.role_temp = ""
                st.session_state.pin_verified = False
                st.session_state.pin_setup_mode = False
                st.rerun()

# =========================
# LEADERBOARD
# =========================
def get_exam_leaderboard(exam_id):
    try:
        attempts = supabase.table("exam_attempts").select("*").eq("exam_id", exam_id).execute().data
        user_best_scores = {}
        for att in attempts:
            uid = att["user_id"]
            score = att["score"]
            if uid not in user_best_scores or score > user_best_scores[uid]:
                user_best_scores[uid] = score
        leaderboard = []
        for uid, max_score in user_best_scores.items():
            user_info = supabase.table("users").select("name, email").eq("id", uid).execute().data
            if user_info:
                leaderboard.append({"Name": user_info[0]["name"], "Email": user_info[0]["email"], "Score": max_score})
        return sorted(leaderboard, key=lambda x: x["Score"], reverse=True)
    except Exception:
        return []

# =========================
# NOTIFICATIONS
# =========================
def send_notification(message, user_id=None):
    supabase.table("notifications").insert({
        "user_id": str(user_id) if user_id else None,
        "message": message,
    }).execute()

def send_admin_notification(message):
    try:
        admins = supabase.table("users").select("id").eq("role", "admin").execute().data or []
        if admins:
            for admin in admins:
                send_notification(message, admin["id"])
        else:
            send_notification(message)
    except Exception:
        send_notification(message)

def get_unread_notifications(user_id):
    try:
        uid = str(user_id)
        all_notifs = supabase.table("notifications").select("*").order("created_at", desc=True).execute().data
        read_data = supabase.table("notification_reads").select("notification_id").eq("user_id", uid).execute().data
        read_ids = {r["notification_id"] for r in read_data}
        return [n for n in all_notifs if (n["user_id"] is None or n["user_id"] == uid) and n["id"] not in read_ids]
    except Exception:
        return []

def mark_notifications_read(user_id):
    try:
        uid = str(user_id)
        unread = get_unread_notifications(uid)
        for n in unread:
            supabase.table("notification_reads").upsert(
                {"user_id": uid, "notification_id": n["id"]}, on_conflict="user_id,notification_id"
            ).execute()
    except Exception:
        pass

def clean_ui_text(value):
    text = str(value or "")
    return "".join(ch for ch in text if ch == "\n" or ch == "\t" or 32 <= ord(ch) <= 126).strip()

def show_notification_banner(user_id):
    notifs = get_unread_notifications(user_id)
    if not notifs:
        return
    st.markdown("""
        <style>
        .notif-wrap { border-left:5px solid #ff6b6b;background:#fff5f5;padding:10px 16px;
          border-radius:0 8px 8px 0;margin-bottom:5px;font-size:0.92rem;color:#c0392b;font-weight:500; }
        .notif-wrap:first-child { border-left-color:#e74c3c;background:#ffeaea;font-weight:700;font-size:0.97rem; }
        </style>""", unsafe_allow_html=True)
    notif_html = "".join(f"<div class='notif-wrap'>Alert: {clean_ui_text(n.get('message'))}</div>" for n in notifs)
    st.markdown(notif_html, unsafe_allow_html=True)
    if st.button(f"Mark all as Read ({len(notifs)})", key="mark_notif_read", type="secondary"):
        mark_notifications_read(user_id)
        st.rerun()

# =========================
# GROUP CHAT
# =========================
def get_unread_count(user_id):
    try:
        read_data = supabase.table("message_reads").select("last_read_at").eq("user_id", str(user_id)).execute().data
        if not read_data:
            total = supabase.table("messages").select("id").execute().data
            return len(total)
        last_read = read_data[0]["last_read_at"]
        unread = supabase.table("messages").select("id").gt("created_at", last_read).neq("user_id", str(user_id)).execute().data
        return len(unread)
    except Exception:
        return 0

def group_chat():
    st.title("Group Chat")
    user_id = str(st.session_state.user_id)
    messages = supabase.table("messages").select("*").order("created_at", desc=False).limit(50).execute().data
    chat_container = st.container(height=450)
    with chat_container:
        if not messages:
            st.info("No messages yet. Send the first message.")
        for msg in messages:
            is_me = msg["user_id"] == user_id
            time_str = str(msg.get("created_at", ""))[:16]
            if is_me:
                col1, col2 = st.columns([2, 5])
                with col2:
                    st.markdown(
                        f"<div style='background:#dcf8c6;padding:10px 14px;border-radius:12px 12px 0px 12px;margin:4px 0;'>"
                        f"<small style='color:#555;font-weight:600'>You</small><br>{msg['message']}"
                        f"<br><small style='color:#999;font-size:10px'>{time_str}</small></div>",
                        unsafe_allow_html=True)
            else:
                col1, col2 = st.columns([5, 2])
                with col1:
                    st.markdown(
                        f"<div style='background:#f1f0f0;padding:10px 14px;border-radius:12px 12px 12px 0px;margin:4px 0;'>"
                        f"<small style='color:#0084ff;font-weight:600'>{msg.get('user_name','Unknown')}</small><br>{msg['message']}"
                        f"<br><small style='color:#999;font-size:10px'>{time_str}</small></div>",
                        unsafe_allow_html=True)
    st.divider()
    col_input, col_send, col_read = st.columns([5, 1, 1])
    with col_input:
        new_msg = st.text_input("Message ...", key="chat_input", label_visibility="collapsed")
    with col_send:
        send = st.button("Send", use_container_width=True, type="primary")
    with col_read:
        mark_read = st.button("Read", use_container_width=True)
    if send and new_msg.strip():
        uinfo = supabase.table("users").select("name").eq("id", user_id).execute().data
        uname = uinfo[0]["name"] if uinfo else "Unknown"
        supabase.table("messages").insert({"user_id": user_id, "user_name": uname, "message": new_msg.strip()}).execute()
        supabase.table("message_reads").upsert({"user_id": user_id, "last_read_at": "now()"}, on_conflict="user_id").execute()
        st.rerun()
    if mark_read:
        supabase.table("message_reads").upsert({"user_id": user_id, "last_read_at": "now()"}, on_conflict="user_id").execute()
        st.success("All messages marked as read.")
        st.rerun()

# =========================
# ATTENDANCE (LeetCode/GitHub style)
# =========================
def mark_today_attendance(user_id):
    today = date.today().isoformat()
    if st.session_state.get("attendance_marked_date") == today:
        return
    try:
        supabase.table("attendance").upsert(
            {"user_id": str(user_id), "attendance_date": today},
            on_conflict="user_id,attendance_date"
        ).execute()
        st.session_state.attendance_marked_date = today
    except Exception:
        pass

def show_attendance_tab(user_id):
    st.title("My Attendance")
    today = date.today()
    start_day = today - timedelta(days=364)

    try:
        rows = supabase.table("attendance").select("attendance_date") \
            .eq("user_id", str(user_id)) \
            .gte("attendance_date", start_day.isoformat()) \
            .lte("attendance_date", today.isoformat()) \
            .execute().data
        attended = {str(r["attendance_date"]) for r in rows if r.get("attendance_date")}
    except Exception as e:
        st.warning("Attendance table database  . Admin SQL run .")
        with st.expander(" SQL to create attendance table"):
            st.code("""
CREATE TABLE IF NOT EXISTS attendance (
  id uuid PRIMARY KEY DEFAULT gen_random_uuid(),
  user_id uuid NOT NULL REFERENCES users(id) ON DELETE CASCADE,
  attendance_date date NOT NULL,
  created_at timestamptz DEFAULT now(),
  UNIQUE(user_id, attendance_date)
);""", language="sql")
        return

    total = len(attended)
    streak = 0
    cursor = today
    while cursor.isoformat() in attended:
        streak += 1
        cursor -= timedelta(days=1)

    # Longest streak
    longest = 0
    run = 0
    prev_d = None
    for ds in sorted(attended):
        d2 = date.fromisoformat(ds)
        if prev_d and (d2 - prev_d).days == 1:
            run += 1
        else:
            run = 1
        longest = max(longest, run)
        prev_d = d2

    c1, c2, c3, c4 = st.columns(4)
    c1.metric(" Total Days", total)
    c2.metric(" Current Streak", f"{streak} days")
    c3.metric(" Longest Streak", f"{longest} days")
    c4.metric(" Today", " Present" if today.isoformat() in attended else " Pending")

    st.divider()

    # Build heatmap  weeks as columns, days as rows (Sun=0 .. Sat=6)
    # Align start to Sunday
    offset = start_day.isoweekday() % 7  # Sun=0
    grid_start = start_day - timedelta(days=offset)
    all_days = [grid_start + timedelta(days=i) for i in range((today - grid_start).days + 1)]
    weeks = [all_days[i:i+7] for i in range(0, len(all_days), 7)]

    # Month labels row
    month_labels = []
    last_m = ""
    for week in weeks:
        lbl = ""
        for d in week:
            if d.day <= 7 and d.strftime("%b") != last_m:
                lbl = d.strftime("%b")
                last_m = lbl
                break
        month_labels.append(lbl)

    label_html = "".join(
        f"<div style='width:14px;font-size:10px;color:#57606a;text-align:left;'>{lbl}</div>"
        for lbl in month_labels
    )

    week_html = ""
    for week in weeks:
        cells = ""
        for d in week:
            if d > today:
                color = "transparent"
                border = "none"
            elif d.isoformat() in attended:
                color = "#2ea043"
                border = "1px solid rgba(27,31,36,0.1)"
            else:
                color = "#ebedf0"
                border = "1px solid rgba(27,31,36,0.06)"
            title = f"{d.strftime('%d %b %Y')}  {'Present' if d.isoformat() in attended else 'Absent'}"
            cells += (
                f"<div title='{title}' style='width:14px;height:14px;border-radius:2px;"
                f"background:{color};border:{border};'></div>"
            )
        week_html += f"<div style='display:flex;flex-direction:column;gap:3px;'>{cells}</div>"

    st.markdown(
        f"""
        <style>
        .att-card{{background:#fff;border:1px solid #d0d7de;border-radius:10px;padding:18px 20px;margin-bottom:16px;}}
        .att-title{{font-size:1.1rem;font-weight:700;margin-bottom:4px;}}
        .att-sub{{color:#57606a;font-size:0.88rem;margin-bottom:14px;}}
        </style>
        <div class="att-card">
          <div class="att-title">Daily Login Activity</div>
          <div class="att-sub">Login    green   (last 1 year)</div>
          <div style="display:flex;gap:3px;margin-bottom:4px;">{label_html}</div>
          <div style="display:flex;gap:3px;overflow-x:auto;">{week_html}</div>
          <div style="display:flex;align-items:center;gap:5px;justify-content:flex-end;color:#57606a;font-size:11px;margin-top:8px;">
            <span>Less</span>
            <div style="width:12px;height:12px;border-radius:2px;background:#ebedf0;border:1px solid rgba(27,31,36,0.06);"></div>
            <div style="width:12px;height:12px;border-radius:2px;background:#2ea043;border:1px solid rgba(27,31,36,0.1);"></div>
            <span>More</span>
          </div>
        </div>
        """,
        unsafe_allow_html=True
    )

# =========================
# STUDENT PROGRESS
# =========================
def get_student_completed_ids(user_id):
    if st.session_state.completed_ids is None:
        try:
            rows = supabase.table("class_completions").select("class_id").eq("user_id", user_id).execute().data
            st.session_state.completed_ids = {str(r["class_id"]) for r in rows}
        except Exception:
            st.session_state.completed_ids = set()
    return st.session_state.completed_ids

def focus_student_class(class_id, exam_id=""):
    st.session_state.focus_class_id = str(class_id) if class_id else ""
    st.session_state.focus_exam_id = str(exam_id) if exam_id else ""
    st.session_state.user_page = "Exams" if exam_id else "My Classes"
    st.rerun()

def start_student_exam(exam):
    has_pwd = exam.get("password") and str(exam["password"]).strip()
    if has_pwd:
        focus_student_class(exam.get("class_id"), exam.get("id"))
        return
    q_data = supabase.table("questions").select("*").eq("exam_id", exam["id"]).execute().data
    start_exam_with_questions(exam, q_data)
    st.rerun()

def collect_student_progress(user_id):
    completed_ids = get_student_completed_ids(user_id)
    modules = supabase.table("modules").select("*").execute().data or []
    submodules = supabase.table("submodules").select("*").execute().data or []
    classes = supabase.table("classes").select("*").execute().data or []
    exams = supabase.table("exams").select("*").execute().data or []
    attempts = supabase.table("exam_attempts").select("*").eq("user_id", user_id).execute().data or []

    sub_by_module = {}
    for sub in submodules:
        sub_by_module.setdefault(str(sub.get("module_id")), []).append(sub)

    classes_by_sub = {}
    class_module = {}
    for cls in classes:
        sid = str(cls.get("submodule_id"))
        classes_by_sub.setdefault(sid, []).append(cls)

    sub_module = {str(s.get("id")): str(s.get("module_id")) for s in submodules}
    for cls in classes:
        class_module[str(cls.get("id"))] = sub_module.get(str(cls.get("submodule_id")), "")

    best_scores = {}
    for att in attempts:
        eid = str(att.get("exam_id"))
        score = int(att.get("score") or 0)
        if eid not in best_scores or score > best_scores[eid]:
            best_scores[eid] = score

    active_exams = [e for e in exams if e.get("enabled")]
    question_counts = {}
    for exam in active_exams:
        try:
            q_rows = supabase.table("questions").select("*").eq("exam_id", exam["id"]).execute().data
            question_counts[str(exam["id"])] = get_exam_max_marks(q_rows)
        except Exception:
            question_counts[str(exam["id"])] = 0

    modules_data = []
    overall_total_classes = len(classes)
    overall_done_classes = sum(1 for cls in classes if str(cls.get("id")) in completed_ids)
    overall_marks_scored = 0
    overall_marks_total = 0
    pending_classes = []
    pending_exams = []

    for module in modules:
        mid = str(module.get("id"))
        module_classes = []
        for sub in sub_by_module.get(mid, []):
            for cls in classes_by_sub.get(str(sub.get("id")), []):
                cls["_submodule_title"] = sub.get("title", "")
                module_classes.append(cls)

        module_exam_rows = [e for e in active_exams if class_module.get(str(e.get("class_id"))) == mid]
        module_done = sum(1 for cls in module_classes if str(cls.get("id")) in completed_ids)
        module_total = len(module_classes)
        module_scored = 0
        module_total_marks = 0

        for exam in module_exam_rows:
            eid = str(exam.get("id"))
            total_q = question_counts.get(eid, 0)
            if total_q > 0:
                module_total_marks += total_q
                overall_marks_total += total_q
                best = best_scores.get(eid, 0)
                module_scored += best
                overall_marks_scored += best
            if eid not in best_scores:
                pending_exams.append({"exam": exam, "total_q": total_q})

        for cls in module_classes:
            if str(cls.get("id")) not in completed_ids:
                pending_classes.append(cls)

        modules_data.append({
            "module": module,
            "class_done": module_done,
            "class_total": module_total,
            "class_pct": int(module_done / module_total * 100) if module_total else 0,
            "marks_scored": module_scored,
            "marks_total": module_total_marks,
            "marks_pct": int(module_scored / module_total_marks * 100) if module_total_marks else 0,
        })

    return {
        "modules": modules_data,
        "overall_class_pct": int(overall_done_classes / overall_total_classes * 100) if overall_total_classes else 0,
        "overall_done_classes": overall_done_classes,
        "overall_total_classes": overall_total_classes,
        "overall_marks_pct": int(overall_marks_scored / overall_marks_total * 100) if overall_marks_total else 0,
        "overall_marks_scored": overall_marks_scored,
        "overall_marks_total": overall_marks_total,
        "pending_classes": pending_classes,
        "pending_exams": pending_exams,
        "best_scores": best_scores,
        "question_counts": question_counts,
        "active_exams": active_exams,
    }

def show_student_progress_tab(user_id):
    st.title("My Progress")
    try:
        progress = collect_student_progress(user_id)
    except Exception as e:
        st.error(f"Progress load avvaledu: {e}")
        return

    c1, c2, c3, c4 = st.columns(4)
    c1.metric("Overall Marks", f"{progress['overall_marks_pct']}%")
    c2.metric("Marks", f"{progress['overall_marks_scored']}/{progress['overall_marks_total']}")
    c3.metric("Class Progress", f"{progress['overall_class_pct']}%")
    c4.metric("Pending", f"{len(progress['pending_classes'])} classes / {len(progress['pending_exams'])} exams")

    st.progress(progress["overall_marks_pct"] / 100, text=f"Overall marks: {progress['overall_marks_pct']}%")
    st.progress(progress["overall_class_pct"] / 100, text=f"Classes completed: {progress['overall_done_classes']}/{progress['overall_total_classes']}")

    st.subheader("Module-wise Progress")
    for item in progress["modules"]:
        module = item["module"]
        with st.container(border=True):
            st.markdown(f"#### {module.get('title', 'Module')}")
            m1, m2 = st.columns(2)
            with m1:
                st.caption(f"Marks: {item['marks_scored']}/{item['marks_total']} ({item['marks_pct']}%)")
                st.progress(item["marks_pct"] / 100)
            with m2:
                st.caption(f"Classes: {item['class_done']}/{item['class_total']} ({item['class_pct']}%)")
                st.progress(item["class_pct"] / 100)

    st.subheader("Pending Classes")
    if not progress["pending_classes"]:
        st.success("All classes complete ayyayi!")
    else:
        for cls in progress["pending_classes"]:
            with st.container(border=True):
                col_info, col_btn = st.columns([4, 1])
                with col_info:
                    st.markdown(f"**{cls.get('title', 'Class')}**")
                    if cls.get("_submodule_title"):
                        st.caption(cls["_submodule_title"])
                with col_btn:
                    if st.button("Open", key=f"prog_open_cls_{cls['id']}", use_container_width=True):
                        focus_student_class(cls.get("id"))

    st.subheader("Pending Exams")
    if not progress["pending_exams"]:
        st.success("Pending exams levu.")
    else:
        for row in progress["pending_exams"]:
            exam = row["exam"]
            with st.container(border=True):
                col_info, col_btn = st.columns([4, 1])
                with col_info:
                    st.markdown(f"**{exam.get('title', 'Exam')}**")
                    st.caption(f"{row['total_q']} questions  {exam.get('duration_mins', 30)} mins")
                with col_btn:
                    if st.button("Open", key=f"prog_open_exam_{exam['id']}", use_container_width=True, type="primary"):
                        start_student_exam(exam)

def show_code_practice_tab():
    st.title("Code Practice")
    if "practice_language" not in st.session_state:
        st.session_state.practice_language = "java"
    if "practice_code_by_language" not in st.session_state:
        st.session_state.practice_code_by_language = {}
    if "practice_input" not in st.session_state:
        st.session_state.practice_input = ""

    selected_label = st.selectbox(
        "Language",
        list(PROGRAMMING_LANGUAGE_LABELS.keys()),
        index=list(PROGRAMMING_LANGUAGE_LABELS.values()).index(st.session_state.practice_language),
        key="practice_language_selector",
    )
    selected_language = PROGRAMMING_LANGUAGE_LABELS[selected_label]
    st.session_state.practice_language = selected_language
    lang_meta = get_programming_language_meta(selected_language)
    inject_vscode_editor_styles()
    enable_textarea_tab_support()
    enable_ide_textarea_behaviour(f"practice_{selected_language}", selected_language)

    if selected_language not in st.session_state.practice_code_by_language:
        st.session_state.practice_code_by_language[selected_language] = lang_meta["default_code"]

    st.markdown(
        f"""
        <div class="vscode-shell">
            <div class="vscode-titlebar">
                <div class="vscode-dots"><span></span><span></span><span></span></div>
                <div class="vscode-file">{lang_meta.get('file_name', 'solution.txt')}</div>
                <div class="vscode-hints">Tab: indent/snippet | brackets auto-close</div>
            </div>
        </div>
        """,
        unsafe_allow_html=True,
    )
    st.session_state.practice_code_by_language[selected_language] = st.text_area(
        f"{lang_meta['label']} Code",
        value=st.session_state.practice_code_by_language[selected_language],
        height=360,
        key=f"practice_code_editor_{selected_language}"
    )
    st.session_state.practice_input = st.text_area(
        "Input",
        value=st.session_state.practice_input,
        height=120,
        key="practice_input_editor",
        placeholder="Program ki kavalsina input ikkada type cheyyandi..."
    )

    col_run, col_reset = st.columns([1, 1])
    with col_run:
        run_clicked = st.button(f"Run {lang_meta['label']}", type="primary", use_container_width=True)
    with col_reset:
        if st.button("Reset Sample", use_container_width=True):
            st.session_state.practice_code_by_language[selected_language] = lang_meta["default_code"]
            st.session_state.practice_input = ""
            st.rerun()

    if run_clicked:
        code = st.session_state.practice_code_by_language[selected_language]
        if not code.strip():
            st.warning(f"{lang_meta['label']} code enter cheyyandi.")
            return
        with st.spinner(f"{lang_meta['label']} program run avuthundi..."):
            result = run_programming_code(code, st.session_state.practice_input, selected_language)
        st.subheader("Output")
        st.code(result.get("stdout", ""), language="text")
        if result.get("stderr"):
            st.subheader("Errors / Compiler Messages")
            st.code(result["stderr"], language="text")
        meta = []
        if result.get("time") is not None:
            meta.append(f"Time: {result['time']}s")
        if result.get("memory") is not None:
            meta.append(f"Memory: {result['memory']} KB")
        status_text = result.get("status", "Unknown")
        if result.get("ok"):
            st.success(f"Status: {status_text}" + (f" | {' | '.join(meta)}" if meta else ""))
        else:
            st.error(f"Status: {status_text}" + (f" | {' | '.join(meta)}" if meta else ""))


def show_java_practice_tab():
    show_code_practice_tab()


def show_programming_questions_tab(user_id):
    st.title("Programming Questions")
    exams = supabase.table("exams").select("*").eq("enabled", True).execute().data or []
    exam_rows = []
    for exam in exams:
        q_rows = supabase.table("questions").select("*").eq("exam_id", exam["id"]).eq("type", "programming").execute().data or []
        if q_rows:
            exam_rows.append({"exam": exam, "questions": q_rows})

    if not exam_rows:
        st.info("Programming questions levu.")
        return

    st.caption("Admin oka exam lo 5 programs add cheste, ikkada single exam card laga kanipistundi. Start chesthe aa 5 questions same exam workspace lo vastayi.")
    for idx, row in enumerate(exam_rows, start=1):
        exam = row["exam"]
        q_rows = row["questions"]
        total_marks = get_exam_max_marks(q_rows)
        attempted_count = sum(1 for q in q_rows if user_attempted_question(user_id, q["id"]))
        attempted = attempted_count == len(q_rows)
        with st.container(border=True):
            c_name, c_meta, c_action = st.columns([5, 2, 2])
            with c_name:
                st.markdown(f"### {idx}. {exam.get('title', 'Programming Exam')}")
                preview_names = ", ".join(str(q.get("question", "Untitled")) for q in q_rows[:3])
                if len(q_rows) > 3:
                    preview_names += f" + {len(q_rows) - 3} more"
                st.caption(preview_names)
            with c_meta:
                st.metric("Programs", len(q_rows))
                st.caption(f"Marks: {total_marks}")
                st.caption(f"Duration: {exam.get('duration_mins', 30)} mins")
                st.caption(f"Done: {attempted_count}/{len(q_rows)}")
            with c_action:
                label = "Attempted / Solve Again" if attempted else "Start Exam"
                has_pwd = exam.get("password") and str(exam.get("password")).strip()
                entered_pwd = ""
                if has_pwd:
                    entered_pwd = st.text_input("Access Code", type="password", key=f"prog_pwd_exam_{exam['id']}", label_visibility="collapsed")
                try:
                    ppt_bytes = generate_programming_questions_ppt(q_rows, exam.get("title", "Programming Questions"))
                    st.download_button(
                        "Download Questions PPT",
                        data=ppt_bytes,
                        file_name=f"{str(exam.get('title','programming'))[:25]}_programming_questions.pptx",
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                        key=f"download_prog_ppt_{exam['id']}",
                        use_container_width=True,
                    )
                except Exception as e:
                    st.caption(f"PPT download ready avvaledu: {e}")
                if st.button(label, key=f"solve_prog_exam_{exam['id']}", use_container_width=True, type="primary" if not attempted else "secondary"):
                    if has_pwd and entered_pwd.strip() != str(exam.get("password")).strip():
                        st.error("Wrong Password!")
                    else:
                        start_exam_with_questions(exam, q_rows)
                        st.rerun()

# =========================
# REVIEW SHEET (styled like screenshot)
# =========================
def render_review_sheet(questions, ans_map, db_attempt):
    if "explain_selected" not in st.session_state:
        st.session_state.explain_selected = set()

    for i, q in enumerate(questions):
        u_ans = ans_map.get(q["id"], "Not Answered")
        c_ans = str(q.get("correct_answer", "")).strip()
        is_correct = str(u_ans).strip().lower() == c_ans.lower()

        with st.container(border=True):
            # Header row: Q number + badge +  button
            hcol, bcol = st.columns([7, 1])
            with hcol:
                badge_style = (
                    "background:#e8f8ef;color:#1e7e45;font-weight:700;font-size:0.75rem;"
                    "padding:2px 10px;border-radius:20px;margin-left:8px;vertical-align:middle;"
                ) if is_correct else (
                    "background:#fdecea;color:#c0392b;font-weight:700;font-size:0.75rem;"
                    "padding:2px 10px;border-radius:20px;margin-left:8px;vertical-align:middle;"
                )
                badge_txt = "Correct" if is_correct else "Incorrect"
                st.markdown(
                    f"<div style='font-weight:700;font-size:1rem;margin-bottom:6px;'>"
                    f"Question {i+1} &nbsp;<span style='{badge_style}'>{badge_txt}</span></div>",
                    unsafe_allow_html=True
                )
                st.markdown(q["question"])
                if q.get("image_url"):
                    st.image(q["image_url"], width=320)
            with bcol:
                qid = q["id"]
                if qid in st.session_state.explain_selected:
                    if st.button("Marked", key=f"exp_{qid}", use_container_width=True, type="primary"):
                        st.session_state.explain_selected.discard(qid)
                        st.rerun()
                else:
                    if st.button("Explain", key=f"exp_{qid}", use_container_width=True):
                        st.session_state.explain_selected.add(qid)
                        st.rerun()

            # Time spent
            t_spent = get_answer_time_spent(db_attempt[0]["id"], q["id"])
            if t_spent and t_spent > 0:
                mins_s, secs_s = divmod(t_spent, 60)
                tstr = f"{mins_s}m {secs_s}s" if mins_s > 0 else f"{secs_s}s"
                st.caption(f" Time spent: **{tstr}**")

            # MCQ styled options
            if q["type"] == "mcq":
                opts = [("A", q.get("option_a","")), ("B", q.get("option_b","")),
                        ("C", q.get("option_c","")), ("D", q.get("option_d",""))]
                correct_display = c_ans
                opts_html = ""
                for lbl, otxt in opts:
                    # Is this the correct option?
                    is_opt_correct = (
                        c_ans.upper() == lbl
                        or c_ans.lower() == str(otxt).strip().lower()
                    )
                    # Is this what the user picked?
                    is_user_pick = (
                        str(u_ans).strip().upper() == lbl
                        or str(u_ans).strip().lower() == str(otxt).strip().lower()
                    )
                    if is_opt_correct:
                        correct_display = f"{lbl}. {otxt}"
                        bg, br, col, suffix, fw = "#eafaf0","#27ae60","#1b5e34","  ","700"
                    elif is_user_pick and not is_correct:
                        bg, br, col, suffix, fw = "#fdecea","#e74c3c","#c0392b","  (Your Answer) ","700"
                    else:
                        bg, br, col, suffix, fw = "#ffffff","#dfe6ee","#2c3e50","","400"
                    opts_html += (
                        f"<div style='background:{bg};border:1.5px solid {br};border-radius:10px;"
                        f"padding:12px 16px;margin-bottom:8px;color:{col};font-weight:{fw};'>"
                        f"{lbl}. {otxt}{suffix}</div>"
                    )
                st.markdown(opts_html, unsafe_allow_html=True)
                st.markdown(
                    f"<div style='background:#eaf2fb;border:1.5px solid #4a90d9;border-radius:10px;"
                    f"padding:12px 16px;margin:8px 0;color:#1a5276;font-weight:700;'>"
                    f"Correct Answer: {correct_display}</div>",
                    unsafe_allow_html=True
                )

            elif q["type"] == "programming":
                try:
                    prog_ans = json.loads(u_ans) if isinstance(u_ans, str) and u_ans.strip().startswith("{") else {"code": u_ans}
                except Exception:
                    prog_ans = {"code": u_ans}
                st.caption(f"Score: {prog_ans.get('earned', 0)}/{prog_ans.get('total', get_question_max_marks(q))} ({prog_ans.get('percentage', 0)}%)")
                answer_language = prog_ans.get("language", get_programming_meta(q).get("language", "java"))
                st.code(prog_ans.get("code", ""), language=get_programming_language_meta(answer_language)["code_language"])
                for res in prog_ans.get("results", []):
                    badge = " Passed" if res.get("passed") else " Failed"
                    st.caption(f"Test Case {res.get('case')}: {badge}  {res.get('marks', 0)} marks")

            else:
                ans_bg = "#eafaf0" if is_correct else "#fdecea"
                ans_br = "#27ae60" if is_correct else "#e74c3c"
                ans_col = "#1b5e34" if is_correct else "#c0392b"
                ans_sfx = " " if is_correct else " (Your Answer) "
                st.markdown(
                    f"<div style='background:{ans_bg};border:1.5px solid {ans_br};border-radius:10px;"
                    f"padding:12px 16px;margin-bottom:8px;color:{ans_col};font-weight:700;'>"
                    f"Your Answer: {u_ans}{ans_sfx}</div>",
                    unsafe_allow_html=True
                )
                if not is_correct:
                    st.markdown(
                        f"<div style='background:#eaf2fb;border:1.5px solid #4a90d9;border-radius:10px;"
                        f"padding:12px 16px;margin:8px 0;color:#1a5276;font-weight:700;'>"
                        f"Correct Answer: {c_ans}</div>",
                        unsafe_allow_html=True
                    )

            # Answer Explanation box
            explanation = str(q.get("explanation", "") or "").strip()
            if explanation:
                st.markdown(
                    f"<div style='background:#f5f6fa;border:1px solid #d0d8e8;border-radius:10px;"
                    f"padding:14px 16px;margin-top:6px;'>"
                    f"<div style='font-weight:700;margin-bottom:5px;'>Answer Explanation</div>"
                    f"<div style='color:#444;line-height:1.6;'>{explanation}</div></div>",
                    unsafe_allow_html=True
                )

            # Report Question link style
            st.markdown(
                "<div style='margin-top:8px;'>"
                "<span style='color:#e74c3c;font-size:0.82rem;cursor:pointer;'> Report Question</span>"
                "</div>",
                unsafe_allow_html=True
            )

    # Explain request send
    st.divider()
    selected_count = len(st.session_state.explain_selected)
    if selected_count > 0:
        st.info(f" {selected_count} questions marked for explanation")
        if st.button(f" Admin  Explain Request  ({selected_count} questions)", type="primary", use_container_width=True):
            try:
                supabase.table("explain_requests").insert({
                    "user_id": str(st.session_state.user_id),
                    "exam_id": str(st.session_state.exam_id),
                    "question_ids": json.dumps(list(st.session_state.explain_selected)),
                    "status": "pending"
                }).execute()
                uinfo = supabase.table("users").select("name").eq("id", st.session_state.user_id).execute().data
                uname = uinfo[0]["name"] if uinfo else "Student"
                send_notification(f" {uname} {selected_count} questions  explanation request !")
                st.session_state.explain_selected = set()
                st.success("Request sent.")
                st.rerun()
            except Exception as e:
                st.error(f"Request error: {e}")

# =========================
# PPT GENERATOR
# =========================
def generate_exam_ppt(questions, exam_title, q_requesters=None):
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    import io, requests as _req

    def rgb(h):
        h = h.lstrip("#")
        return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

    def box(sl, x, y, w, h, fill, border=None, bw=Pt(1)):
        s = sl.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
        s.fill.solid(); s.fill.fore_color.rgb = rgb(fill)
        if border: s.line.color.rgb = rgb(border); s.line.width = bw
        else: s.line.fill.background()
        return s

    def txt(sl, text, x, y, w, h, sz=13, bold=False, color="1A1A2E", align=PP_ALIGN.LEFT, italic=False):
        tb = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = tb.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.alignment = align
        r = p.add_run(); r.text = str(text)
        r.font.size = Pt(sz); r.font.bold = bold
        r.font.italic = italic; r.font.color.rgb = rgb(color)
        return tb

    NAV="1E2761"; LGT="F4F6FB"; ACC="4A90D9"; GRN="27AE60"; GRD="1E8449"; WHT="FFFFFF"; DRK="1A1A2E"; MUT="7F8C8D"
    prs = Presentation(); prs.slide_width = Inches(10); prs.slide_height = Inches(5.625)
    BL = prs.slide_layouts[6]

    ts = prs.slides.add_slide(BL)
    box(ts, 0, 0, 10, 5.625, NAV); box(ts, 0, 2.5, 10, 0.06, ACC)
    txt(ts, exam_title or "Exam Review", 0.5, 1.0, 9, 1.2, sz=36, bold=True, color=WHT, align=PP_ALIGN.CENTER)
    txt(ts, f"{len(questions)} Questions", 0.5, 2.65, 9, 0.6, sz=20, color="CADCFC", align=PP_ALIGN.CENTER)
    txt(ts, "Correct  Green  |  Wrong  White", 0.5, 4.6, 9, 0.4, sz=12, italic=True, color="8899CC", align=PP_ALIGN.CENTER)

    for idx, q in enumerate(questions):
        sl = prs.slides.add_slide(BL); box(sl, 0, 0, 10, 5.625, LGT)
        cur_y = 0.18
        box(sl, 0.3, cur_y, 0.7, 0.42, ACC)
        txt(sl, f"Q{idx+1}", 0.3, cur_y, 0.7, 0.42, sz=14, bold=True, color=WHT, align=PP_ALIGN.CENTER)
        txt(sl, q.get("question",""), 1.12, cur_y, 8.55, 0.72, sz=15, bold=True, color=DRK)
        cur_y += 0.78
        if q_requesters and q.get("id") in q_requesters:
            names = q_requesters[q["id"]][:4]
            ns = " " + ",  ".join(names)
            if len(q_requesters[q["id"]]) > 4:
                ns += f"  +{len(q_requesters[q['id']])-4} more"
            box(sl, 0.3, cur_y, 9.4, 0.3, "FFF9C4", "F9A825", Pt(1))
            txt(sl, ns, 0.45, cur_y+0.02, 9.1, 0.28, sz=10, italic=True, color="7B5800")
            cur_y += 0.35
        box(sl, 0.3, cur_y, 9.4, 0.03, "D0D8E8"); cur_y += 0.1
        img_available = False
        if q.get("image_url"):
            try:
                import io as _io
                r2 = _req.get(q["image_url"], timeout=10, headers={"User-Agent":"Mozilla/5.0"})
                if r2.status_code == 200 and len(r2.content) > 500:
                    sl.shapes.add_picture(_io.BytesIO(r2.content), Inches(6.0), Inches(cur_y), Inches(3.7), Inches(2.5))
                    img_available = True
            except Exception:
                pass
        correct_ans = str(q.get("correct_answer","")).strip()
        if q.get("type","mcq") == "mcq":
            opts = [("A", q.get("option_a","")), ("B", q.get("option_b","")), ("C", q.get("option_c","")), ("D", q.get("option_d",""))]
            for i, (lbl, otxt) in enumerate(opts):
                ox = 0.3 if img_available else (0.3 if i % 2 == 0 else 5.2)
                oy = cur_y + i * 0.82 if img_available else cur_y + (i // 2) * 0.95
                ow, oh = (5.5, 0.72) if img_available else (4.5, 0.82)
                is_cor = (correct_ans.upper() == lbl or correct_ans.strip().lower() == str(otxt).strip().lower())
                bg = GRN if is_cor else WHT; tc = WHT if is_cor else DRK; br = GRN if is_cor else "C8D6E5"
                box(sl, ox, oy, ow, oh, bg, br, Pt(1.5))
                box(sl, ox+0.1, oy+0.16, 0.46, 0.46, GRD if is_cor else ACC)
                txt(sl, lbl, ox+0.1, oy+0.16, 0.46, 0.46, sz=12, bold=True, color=WHT, align=PP_ALIGN.CENTER)
                txt(sl, str(otxt), ox+0.68, oy+0.08, ow-0.8, oh-0.16, sz=13, color=tc)
            txt(sl, f"  Correct: {correct_ans}", 0.3, 5.15, 9, 0.35, sz=11, bold=True, color=GRN)
        else:
            box(sl, 0.3, cur_y, 9.4, 1.1, "EAF7EE", GRN, Pt(2))
            txt(sl, correct_ans, 0.5, cur_y+0.1, 9.0, 0.9, sz=14, bold=True, color=GRN)
        hint = str(q.get("hint","") or "").strip()
        if hint:
            txt(sl, f" {hint}", 0.3, 5.38, 9, 0.25, sz=10, italic=True, color=MUT)
        txt(sl, f"{idx+1}/{len(questions)}", 8.6, 5.38, 1.1, 0.25, sz=9, color=MUT, align=PP_ALIGN.RIGHT)

    es = prs.slides.add_slide(BL); box(es, 0, 0, 10, 5.625, NAV)
    txt(es, "End of Review", 0.5, 1.8, 9, 1.5, sz=38, bold=True, color=WHT, align=PP_ALIGN.CENTER)
    txt(es, "Keep improving!", 0.5, 3.4, 9, 0.6, sz=18, italic=True, color="CADCFC", align=PP_ALIGN.CENTER)

    buf = io.BytesIO(); prs.save(buf); buf.seek(0)
    return buf.read()

def check_mcq_correct(user_val, q):
    """MCQ answer check  user_val can be label (A/B/C/D) or full text"""
    correct = str(q.get("correct_answer","")).strip()
    user = str(user_val).strip()
    if not user or not correct:
        return False
    # Direct match
    if user.lower() == correct.lower():
        return True
    # User answered as label, correct stored as text
    label_map = {
        "A": str(q.get("option_a","")), "B": str(q.get("option_b","")),
        "C": str(q.get("option_c","")), "D": str(q.get("option_d","")),
    }
    if user.upper() in label_map:
        return label_map[user.upper()].strip().lower() == correct.lower()
    # User answered as text, correct stored as label
    if correct.upper() in label_map:
        return label_map[correct.upper()].strip().lower() == user.lower()
    return False


SUPRABHATAM_LANGUAGES = {
    "Telugu": "telugu_text",
    "English": "english_text",
    "Sanskrit": "sanskrit_text",
    "Hindi": "hindi_text",
    "Tamil": "tamil_text",
    "Kannada": "kannada_text",
}


def get_suprabhatam_sql():
    return """
create table if not exists suprabhatam_slokas (
  id uuid primary key default gen_random_uuid(),
  display_order integer not null,
  title text,
  image_url text,
  telugu_text text,
  english_text text,
  sanskrit_text text,
  hindi_text text,
  tamil_text text,
  kannada_text text,
  meaning text,
  created_at timestamptz default now()
);

create index if not exists suprabhatam_slokas_display_order_idx
on suprabhatam_slokas(display_order);

create table if not exists suprabhatam_access (
  user_id uuid primary key references users(id) on delete cascade,
  enabled boolean not null default true,
  updated_at timestamptz default now()
);
"""


def show_suprabhatam_styles():
    st.markdown("""
    <style>
    .supra-reader {
        border: 1px solid #d9dde7;
        background: #ffffff;
        border-radius: 8px;
        padding: 18px;
        margin: 12px 0 18px 0;
    }
    .supra-frame {
        display: grid;
        grid-template-columns: minmax(140px, 260px) 1fr;
        gap: 18px;
        align-items: stretch;
    }
    .supra-image {
        width: 100%;
        min-height: 180px;
        border: 5px solid #f49a73;
        border-radius: 8px;
        overflow: hidden;
        background: #f7f0e8;
    }
    .supra-image img {
        width: 100%;
        height: 100%;
        min-height: 180px;
        object-fit: cover;
        display: block;
    }
    .supra-quote {
        min-height: 180px;
        border: 1px solid #cfd3dc;
        border-right: 8px solid #f05a28;
        border-bottom: 8px solid #f05a28;
        padding: 26px 28px;
        display: flex;
        align-items: center;
        justify-content: center;
        text-align: center;
        color: #1168c4;
        font-size: 1.55rem;
        line-height: 1.85;
        white-space: pre-wrap;
    }
    .supra-title {
        margin: 0 0 12px 0;
        color: #263047;
        font-size: 1.15rem;
        font-weight: 700;
    }
    .supra-meaning {
        margin-top: 14px;
        color: #3b4254;
        line-height: 1.7;
        white-space: pre-wrap;
    }
    @media (max-width: 760px) {
        .supra-frame { grid-template-columns: 1fr; }
        .supra-quote { font-size: 1.15rem; padding: 20px; }
    }
    </style>
    """, unsafe_allow_html=True)


def fetch_suprabhatam_slokas():
    try:
        return supabase.table("suprabhatam_slokas").select("*").order("display_order").execute().data or []
    except Exception as e:
        st.warning("Suprabhatam tables database lo create cheyyali. Admin SQL run cheyyandi.")
        with st.expander("Suprabhatam database SQL"):
            st.code(get_suprabhatam_sql(), language="sql")
        return []


def user_has_suprabhatam_access(user_id):
    if st.session_state.get("role") == "admin":
        return True
    try:
        rows = supabase.table("suprabhatam_access").select("enabled").eq("user_id", user_id).eq("enabled", True).limit(1).execute().data
        return bool(rows)
    except Exception:
        return False


def render_suprabhatam_reader(slokas=None):
    slokas = slokas if slokas is not None else fetch_suprabhatam_slokas()
    if not slokas:
        st.info("Inka Suprabhatam slokas add cheyyaledu.")
        return

    st.subheader("Suprabhatam")
    sloka_payload = []
    for i, sloka in enumerate(slokas):
        sloka_payload.append({
            "title": sloka.get("title") or f"Slokam {i + 1}",
            "image_url": sloka.get("image_url") or "",
            "meaning": sloka.get("meaning") or "",
            "languages": {
                lang: sloka.get(col) or ""
                for lang, col in SUPRABHATAM_LANGUAGES.items()
            },
        })

    component_html = """
    <div class="book-wrap">
      <div class="toolbar" id="langButtons"></div>
      <div class="counter" id="counter"></div>
      <div class="book" id="book">
        <div class="page" id="page">
          <div class="imageBox" id="imageBox"></div>
          <div class="quoteBox">
            <div class="slokaTitle" id="slokaTitle"></div>
            <div class="slokaText" id="slokaText"></div>
            <div class="meaning" id="meaning"></div>
          </div>
        </div>
      </div>
      <div class="nav">
        <button id="prevBtn" type="button">Previous Slokam</button>
        <button id="nextBtn" type="button">Next Slokam</button>
      </div>
    </div>

    <style>
      * { box-sizing: border-box; }
      body {
        margin: 0;
        font-family: "Segoe UI", Arial, sans-serif;
        color: #263047;
        background: transparent;
      }
      .book-wrap {
        width: 100%;
        padding: 8px 4px 18px;
      }
      .toolbar {
        display: flex;
        gap: 8px;
        flex-wrap: wrap;
        margin-bottom: 10px;
      }
      .toolbar button, .nav button {
        border: 1px solid #cfd6e4;
        background: #ffffff;
        color: #263047;
        border-radius: 8px;
        padding: 9px 13px;
        cursor: pointer;
        font-weight: 650;
      }
      .toolbar button.active {
        background: #1168c4;
        border-color: #1168c4;
        color: #ffffff;
      }
      .counter {
        font-size: 0.95rem;
        color: #5f6678;
        margin-bottom: 10px;
      }
      .book {
        perspective: 1800px;
      }
      .page {
        min-height: 315px;
        display: grid;
        grid-template-columns: minmax(155px, 280px) 1fr;
        gap: 18px;
        align-items: stretch;
        padding: 16px;
        border: 1px solid #d9dde7;
        border-radius: 8px;
        background:
          linear-gradient(90deg, rgba(0,0,0,0.08), rgba(255,255,255,0) 32px),
          #fffdf9;
        box-shadow: 0 10px 28px rgba(38, 48, 71, 0.12);
        transform-origin: left center;
      }
      .page.flip-next {
        animation: pageNext 520ms ease both;
      }
      .page.flip-prev {
        animation: pagePrev 520ms ease both;
      }
      @keyframes pageNext {
        0% { transform: rotateY(0deg); opacity: 1; }
        48% { transform: rotateY(-78deg); opacity: 0.45; }
        100% { transform: rotateY(0deg); opacity: 1; }
      }
      @keyframes pagePrev {
        0% { transform: rotateY(0deg); opacity: 1; }
        48% { transform: rotateY(72deg); opacity: 0.45; }
        100% { transform: rotateY(0deg); opacity: 1; }
      }
      .imageBox {
        min-height: 230px;
        border: 5px solid #f49a73;
        border-radius: 8px;
        overflow: hidden;
        background: #f7f0e8;
      }
      .imageBox img {
        width: 100%;
        height: 100%;
        min-height: 230px;
        object-fit: cover;
        display: block;
      }
      .quoteBox {
        min-height: 230px;
        border: 1px solid #cfd3dc;
        border-right: 8px solid #f05a28;
        border-bottom: 8px solid #f05a28;
        padding: 24px 28px;
        display: flex;
        flex-direction: column;
        justify-content: center;
        text-align: center;
        background: #ffffff;
      }
      .slokaTitle {
        margin-bottom: 12px;
        color: #263047;
        font-weight: 750;
        font-size: 1.05rem;
      }
      .slokaText {
        color: #1168c4;
        font-size: 1.55rem;
        line-height: 1.8;
        white-space: pre-wrap;
      }
      .meaning {
        margin-top: 14px;
        color: #3b4254;
        line-height: 1.6;
        white-space: pre-wrap;
        text-align: left;
      }
      .nav {
        display: grid;
        grid-template-columns: 1fr 1fr;
        gap: 12px;
        margin-top: 14px;
      }
      .nav button:last-child {
        background: #1168c4;
        border-color: #1168c4;
        color: #ffffff;
      }
      button:disabled {
        opacity: 0.45;
        cursor: not-allowed;
      }
      @media (max-width: 760px) {
        .page { grid-template-columns: 1fr; }
        .slokaText { font-size: 1.15rem; }
        .quoteBox { padding: 20px; }
      }
    </style>

    <script>
      const slokas = __SLOKAS__;
      const languageOrder = __LANGUAGES__;
      let index = 0;
      let language = languageOrder.includes("Telugu") ? "Telugu" : languageOrder[0];

      const page = document.getElementById("page");
      const titleEl = document.getElementById("slokaTitle");
      const textEl = document.getElementById("slokaText");
      const meaningEl = document.getElementById("meaning");
      const imageBox = document.getElementById("imageBox");
      const counter = document.getElementById("counter");
      const prevBtn = document.getElementById("prevBtn");
      const nextBtn = document.getElementById("nextBtn");
      const langButtons = document.getElementById("langButtons");

      slokas.forEach((sloka) => {
        if (sloka.image_url) {
          const img = new Image();
          img.src = sloka.image_url;
        }
      });

      function availableLanguages() {
        const current = slokas[index];
        return languageOrder.filter((lang) => (current.languages[lang] || "").trim().length > 0);
      }

      function drawLanguageButtons() {
        const available = availableLanguages();
        if (!available.includes(language)) {
          language = available.includes("Telugu") ? "Telugu" : (available[0] || languageOrder[0]);
        }
        langButtons.innerHTML = "";
        available.forEach((lang) => {
          const btn = document.createElement("button");
          btn.type = "button";
          btn.textContent = lang;
          btn.className = lang === language ? "active" : "";
          btn.addEventListener("click", () => {
            language = lang;
            render();
          });
          langButtons.appendChild(btn);
        });
      }

      function render() {
        const sloka = slokas[index];
        drawLanguageButtons();
        titleEl.textContent = sloka.title || `Slokam ${index + 1}`;
        const text = (sloka.languages[language] || "").trim() || "Selected language text inka add cheyyaledu.";
        textEl.textContent = `â€œ${text}â€`;
        meaningEl.textContent = sloka.meaning ? `Meaning: ${sloka.meaning}` : "";
        imageBox.innerHTML = sloka.image_url ? `<img src="${sloka.image_url}" alt="">` : "";
        counter.textContent = `Slokam ${index + 1} / ${slokas.length}`;
        prevBtn.disabled = index === 0;
        nextBtn.disabled = index === slokas.length - 1;
      }

      function turnPage(direction) {
        const nextIndex = index + direction;
        if (nextIndex < 0 || nextIndex >= slokas.length) return;
        page.classList.remove("flip-next", "flip-prev");
        void page.offsetWidth;
        page.classList.add(direction > 0 ? "flip-next" : "flip-prev");
        setTimeout(() => {
          index = nextIndex;
          render();
        }, 230);
      }

      prevBtn.addEventListener("click", () => turnPage(-1));
      nextBtn.addEventListener("click", () => turnPage(1));
      render();
    </script>
    """
    component_html = component_html.replace("__SLOKAS__", json.dumps(sloka_payload, ensure_ascii=False))
    component_html = component_html.replace("__LANGUAGES__", json.dumps(list(SUPRABHATAM_LANGUAGES.keys()), ensure_ascii=False))
    components.html(component_html, height=690, scrolling=True)


def show_suprabhatam_admin():
    st.title("Suprabhatam")
    st.caption("Manual ga slokas add cheyyandi. Display order prakaram book laga users ki kanipistundi.")

    with st.expander("Database SQL setup"):
        st.code(get_suprabhatam_sql(), language="sql")

    slokas = fetch_suprabhatam_slokas()
    tab_add, tab_manage, tab_access, tab_preview = st.tabs(["Add Slokam", "Manage Slokas", "User Access", "Preview"])

    with tab_add:
        next_order = (max([int(s.get("display_order") or 0) for s in slokas]) + 1) if slokas else 1
        with st.form("add_suprabhatam_slokam"):
            order = st.number_input("Display Order", min_value=1, value=next_order, step=1)
            title = st.text_input("Slokam Title", value=f"Slokam {next_order}")
            image_file = st.file_uploader("Slokam mundu image upload", type=["png", "jpg", "jpeg", "webp"])
            image_url = st.text_input("Leda image URL paste cheyyandi")
            telugu_text = st.text_area("Telugu", height=120)
            english_text = st.text_area("English", height=120)
            sanskrit_text = st.text_area("Sanskrit", height=120)
            hindi_text = st.text_area("Hindi", height=100)
            tamil_text = st.text_area("Tamil", height=100)
            kannada_text = st.text_area("Kannada", height=100)
            meaning = st.text_area("Meaning / Notes", height=90)
            submitted = st.form_submit_button("Add Slokam", type="primary")
        if submitted:
            final_image_url = image_url.strip()
            if image_file:
                uploaded_url = upload_image_to_imgbb(image_file)
                final_image_url = uploaded_url or final_image_url
            payload = {
                "display_order": int(order),
                "title": title.strip() or f"Slokam {int(order)}",
                "image_url": final_image_url,
                "telugu_text": telugu_text.strip(),
                "english_text": english_text.strip(),
                "sanskrit_text": sanskrit_text.strip(),
                "hindi_text": hindi_text.strip(),
                "tamil_text": tamil_text.strip(),
                "kannada_text": kannada_text.strip(),
                "meaning": meaning.strip(),
            }
            try:
                supabase.table("suprabhatam_slokas").insert(payload).execute()
                st.success("Slokam add ayyindi.")
                st.rerun()
            except Exception as e:
                st.error(f"Slokam add avvaledu: {e}")

    with tab_manage:
        if not slokas:
            st.info("Inka slokas levu.")
        for sloka in slokas:
            with st.expander(f"{sloka.get('display_order')}. {sloka.get('title') or 'Slokam'}"):
                with st.form(f"edit_suprabhatam_{sloka['id']}"):
                    order = st.number_input("Display Order", min_value=1, value=int(sloka.get("display_order") or 1), step=1, key=f"order_{sloka['id']}")
                    title = st.text_input("Slokam Title", value=sloka.get("title") or "", key=f"title_{sloka['id']}")
                    image_file = st.file_uploader("Replace image", type=["png", "jpg", "jpeg", "webp"], key=f"img_{sloka['id']}")
                    image_url = st.text_input("Image URL", value=sloka.get("image_url") or "", key=f"url_{sloka['id']}")
                    telugu_text = st.text_area("Telugu", value=sloka.get("telugu_text") or "", height=100, key=f"te_{sloka['id']}")
                    english_text = st.text_area("English", value=sloka.get("english_text") or "", height=100, key=f"en_{sloka['id']}")
                    sanskrit_text = st.text_area("Sanskrit", value=sloka.get("sanskrit_text") or "", height=100, key=f"sa_{sloka['id']}")
                    hindi_text = st.text_area("Hindi", value=sloka.get("hindi_text") or "", height=80, key=f"hi_{sloka['id']}")
                    tamil_text = st.text_area("Tamil", value=sloka.get("tamil_text") or "", height=80, key=f"ta_{sloka['id']}")
                    kannada_text = st.text_area("Kannada", value=sloka.get("kannada_text") or "", height=80, key=f"ka_{sloka['id']}")
                    meaning = st.text_area("Meaning / Notes", value=sloka.get("meaning") or "", height=80, key=f"meaning_{sloka['id']}")
                    save_col, delete_col = st.columns(2)
                    save_clicked = save_col.form_submit_button("Save Changes", type="primary")
                    delete_clicked = delete_col.form_submit_button("Delete")
                if save_clicked:
                    final_image_url = image_url.strip()
                    if image_file:
                        uploaded_url = upload_image_to_imgbb(image_file)
                        final_image_url = uploaded_url or final_image_url
                    try:
                        supabase.table("suprabhatam_slokas").update({
                            "display_order": int(order),
                            "title": title.strip(),
                            "image_url": final_image_url,
                            "telugu_text": telugu_text.strip(),
                            "english_text": english_text.strip(),
                            "sanskrit_text": sanskrit_text.strip(),
                            "hindi_text": hindi_text.strip(),
                            "tamil_text": tamil_text.strip(),
                            "kannada_text": kannada_text.strip(),
                            "meaning": meaning.strip(),
                        }).eq("id", sloka["id"]).execute()
                        st.success("Slokam update ayyindi.")
                        st.rerun()
                    except Exception as e:
                        st.error(f"Update avvaledu: {e}")
                if delete_clicked:
                    try:
                        supabase.table("suprabhatam_slokas").delete().eq("id", sloka["id"]).execute()
                        st.success("Slokam delete ayyindi.")
                        st.rerun()
                    except Exception as e:
                        st.error(f"Delete avvaledu: {e}")

    with tab_access:
        try:
            users = supabase.table("users").select("id, name, email, role").eq("role", "user").order("name").execute().data or []
            access_rows = supabase.table("suprabhatam_access").select("*").execute().data or []
            access_map = {str(row["user_id"]): bool(row.get("enabled")) for row in access_rows}
        except Exception as e:
            users = []
            access_map = {}
            st.warning(f"Access table ready ledu: {e}")
        if not users:
            st.info("Users dorakaledu.")
        for user in users:
            uid = str(user["id"])
            label = f"{user.get('name') or 'No name'} - {user.get('email') or ''}"
            enabled = st.checkbox(label, value=access_map.get(uid, False), key=f"supra_access_{uid}")
            current = access_map.get(uid, False)
            if enabled != current:
                try:
                    if enabled:
                        supabase.table("suprabhatam_access").upsert({"user_id": uid, "enabled": True}, on_conflict="user_id").execute()
                    else:
                        supabase.table("suprabhatam_access").delete().eq("user_id", uid).execute()
                    st.rerun()
                except Exception as e:
                    st.error(f"Access update avvaledu: {e}")

    with tab_preview:
        render_suprabhatam_reader(slokas)


def admin_dashboard():
    st.sidebar.title("Admin Workspace")
    persist_browser_login()
    if st.sidebar.button("Logout", use_container_width=True):
        for key in defaults:
            st.session_state[key] = defaults[key]
        show_logout_redirect()

    st.sidebar.divider()
    if st.session_state.admin_preview_mode:
        if st.sidebar.button("Back to Admin View", use_container_width=True, type="primary"):
            st.session_state.admin_preview_mode = False
            st.rerun()
        user_dashboard(preview_mode=True)
        return
    else:
        show_notification_banner(st.session_state.user_id)
        if st.sidebar.button("Student View Preview", use_container_width=True):
            st.session_state.admin_preview_mode = True
            st.rerun()
    st.sidebar.divider()

    
    unread_admin = get_unread_count(st.session_state.user_id)
    label = f"Group Chat ({unread_admin})" if unread_admin > 0 else "Group Chat"
    
    menu = st.sidebar.selectbox("Navigation Control",
        ["Manage Course Content", "Manage Exams & Questions", "Student Results & Ranks", "Users & Access", "Credit Cards", "Suprabhatam", label],
        key="admin_navigation")
    if "Group Chat" in menu:
        menu = "Group Chat"
    if menu == "Credit Cards":
        admin_credit_cards_dashboard()
        return
    if menu == "Suprabhatam":
        show_suprabhatam_admin()
        return

    if menu == "Users & Access":
        st.subheader("Users & Access")
        with st.expander("Required Supabase columns - run once if signup/access errors vasthe"):
            st.code(get_account_access_sql(), language="sql")

        try:
            users = supabase.table("users").select("*").order("name").execute().data or []
        except Exception as e:
            st.error(f"Users load avvaledu: {e}")
            users = []

        pending_users = [u for u in users if user_account_status(u) == "pending"]
        st.markdown("#### Pending Account Requests")
        if not pending_users:
            st.info("Pending signup requests levu.")
        for u in pending_users:
            with st.container(border=True):
                st.markdown(f"**{u.get('name','')}**  |  {u.get('email','')}  |  @{u.get('username','') or '-'}")
                st.caption(f"Requested account PIN: {u.get('requested_pin') or '-'}")
                approve_col, reject_col = st.columns(2)
                with approve_col:
                    if st.button("Approve Login", key=f"approve_user_{u['id']}", type="primary", use_container_width=True):
                        pin = str(u.get("requested_pin") or "").strip()
                        if not pin:
                            pin = generate_numeric_pin()
                        supabase.table("users").update({
                            "approval_status": "approved",
                            "app_pin": pin,
                            "requested_pin": None,
                        }).eq("id", u["id"]).execute()
                        st.success("User approved.")
                        st.rerun()
                with reject_col:
                    if st.button("Reject", key=f"reject_user_{u['id']}", use_container_width=True):
                        supabase.table("users").update({"approval_status": "rejected"}).eq("id", u["id"]).execute()
                        st.warning("User rejected.")
                        st.rerun()

        st.divider()
        st.markdown("#### Hidden Test Case PIN Access")
        admin_row = get_hidden_test_acl()
        view_pin = str(admin_row.get("hidden_test_view_pin") or "")
        edit_pin = str(admin_row.get("hidden_test_edit_pin") or "")
        st.caption(f"See PIN: {view_pin or 'Not generated'}")
        st.caption(f"Edit PIN: {edit_pin or 'Not generated'}")
        pin_col1, pin_col2 = st.columns(2)
        with pin_col1:
            if st.button("Generate See PIN", use_container_width=True):
                supabase.table("users").update({"hidden_test_view_pin": generate_numeric_pin()}).eq("id", st.session_state.user_id).execute()
                st.success("See PIN generated.")
                st.rerun()
        with pin_col2:
            if st.button("Generate Edit PIN", use_container_width=True):
                supabase.table("users").update({"hidden_test_edit_pin": generate_numeric_pin()}).eq("id", st.session_state.user_id).execute()
                st.success("Edit PIN generated.")
                st.rerun()

        student_users = [u for u in users if u.get("role") == "user" and user_account_status(u) == "approved"]
        user_labels = {f"{u.get('name','')} ({u.get('email','')})": str(u.get("id")) for u in student_users}
        current_view_ids = normalize_id_list(admin_row.get("hidden_test_view_user_ids"))
        current_edit_ids = normalize_id_list(admin_row.get("hidden_test_edit_user_ids"))
        current_view_labels = [label for label, uid in user_labels.items() if uid in current_view_ids]
        current_edit_labels = [label for label, uid in user_labels.items() if uid in current_edit_ids]
        view_labels = st.multiselect("Hidden test cases chudagalige accounts", list(user_labels.keys()), default=current_view_labels, key="hidden_view_users")
        edit_labels = st.multiselect("Hidden test cases edit cheyagalige accounts", list(user_labels.keys()), default=current_edit_labels, key="hidden_edit_users")
        if st.button("Save Hidden Test Access", type="primary", use_container_width=True):
            supabase.table("users").update({
                "hidden_test_view_user_ids": [user_labels[label] for label in view_labels],
                "hidden_test_edit_user_ids": [user_labels[label] for label in edit_labels],
            }).eq("id", st.session_state.user_id).execute()
            st.success("Access saved.")
            st.rerun()

        st.markdown("#### Hidden Test Case View Limits")
        st.caption("Oka hidden test case view chesthe 1 limit consume avuthundi. Already viewed case malli chuste limit taggadu.")
        for u in student_users:
            used = get_hidden_case_views_used(u.get("id"))
            limit_val = int(u.get("hidden_test_view_limit") or 0)
            c_user, c_limit, c_save = st.columns([4, 2, 1])
            with c_user:
                st.markdown(f"**{u.get('name','')}** ({u.get('email','')})")
                st.caption(f"Used: {used} | Remaining: {max(0, limit_val - used)}")
            with c_limit:
                new_limit = st.number_input("Limit", min_value=0, max_value=1000, value=limit_val, step=1, key=f"hidden_limit_{u['id']}", label_visibility="collapsed")
            with c_save:
                if st.button("Save", key=f"save_hidden_limit_{u['id']}", use_container_width=True):
                    supabase.table("users").update({"hidden_test_view_limit": int(new_limit)}).eq("id", u["id"]).execute()
                    st.success("Limit saved.")
                    st.rerun()
        return

    if menu == "Manage Course Content":
        tab1, tab2, tab3 = st.tabs(["Modules Setup", "Submodules Setup", "Live/Recorded Classes"])

        with tab1:
            st.subheader("Manage Core Modules")
            with st.form("add_module_form", clear_on_submit=True):
                module_name = st.text_input("New Module Title")
                if st.form_submit_button(" Save Module"):
                    if module_name.strip():
                        supabase.table("modules").insert({"title": module_name}).execute()
                        st.success("Module Added!")
                        st.rerun()
            st.divider()
            modules = supabase.table("modules").select("*").execute().data
            for m in modules:
                col1, col2, col3 = st.columns([4, 1, 1])
                with col1:
                    new_m_title = st.text_input("Module Name", value=m["title"], key=f"mod_t_{m['id']}")
                with col2:
                    if st.button("Update", key=f"mod_u_{m['id']}", use_container_width=True):
                        supabase.table("modules").update({"title": new_m_title}).eq("id", m["id"]).execute()
                        st.success("Updated!"); st.rerun()
                with col3:
                    if st.button("Delete", key=f"mod_d_{m['id']}", type="secondary", use_container_width=True):
                        supabase.table("modules").delete().eq("id", m["id"]).execute()
                        st.warning("Deleted!"); st.rerun()

        with tab2:
            st.subheader("Manage Submodules")
            modules_list = supabase.table("modules").select("*").execute().data
            mod_options = {m["title"]: m["id"] for m in modules_list} if modules_list else {}
            with st.form("add_sub_form", clear_on_submit=True):
                sel_mod = st.selectbox("Select Parent Module", list(mod_options.keys()) or ["No modules yet"])
                sub_name = st.text_input("Submodule Title")
                if st.form_submit_button(" Save Submodule"):
                    if sub_name.strip() and sel_mod in mod_options:
                        supabase.table("submodules").insert({"module_id": mod_options[sel_mod], "title": sub_name}).execute()
                        st.success("Linked!"); st.rerun()
            st.divider()
            submodules = supabase.table("submodules").select("*").execute().data
            for s in submodules:
                p_module = supabase.table("modules").select("title").eq("id", s["module_id"]).execute().data
                p_title = p_module[0]["title"] if p_module else "Unknown"
                col1, col2, col3, col4 = st.columns([2, 3, 1, 1])
                with col1: st.caption(f"Parent: {p_title}")
                with col2:
                    new_s_title = st.text_input("Edit Title", value=s["title"], key=f"sub_t_{s['id']}", label_visibility="collapsed")
                with col3:
                    if st.button("Update", key=f"sub_u_{s['id']}", use_container_width=True):
                        supabase.table("submodules").update({"title": new_s_title}).eq("id", s["id"]).execute()
                        st.success("Updated!"); st.rerun()
                with col4:
                    if st.button("Delete", key=f"sub_d_{s['id']}", type="secondary", use_container_width=True):
                        supabase.table("submodules").delete().eq("id", s["id"]).execute()
                        st.warning("Deleted!"); st.rerun()

        with tab3:
            st.subheader("Manage Stream/Video Classes")
            sub_list = supabase.table("submodules").select("*").execute().data
            sub_options = {s["title"]: s["id"] for s in sub_list} if sub_list else {}
            with st.expander(" Add New Class Room"):
                with st.form("add_class_form", clear_on_submit=True):
                    sel_sub = st.selectbox("Link to Submodule", list(sub_options.keys()) or ["No submodules yet"])
                    c_title = st.text_input("Class Title")
                    c_link = st.text_input("Live Stream Link")
                    v_link = st.text_input("Recorded Video URL")
                    p_link = st.text_input("Notes PDF URL")
                    if st.form_submit_button(" Deploy Class"):
                        if sel_sub in sub_options:
                            supabase.table("classes").insert({
                                "submodule_id": sub_options[sel_sub], "title": c_title,
                                "class_link": c_link, "recorded_video": v_link, "notes_pdf": p_link
                            }).execute()
                            st.success("Class Broadcasted!"); st.rerun()
            st.divider()
            all_users = supabase.table("users").select("id, role").execute().data
            total_users = len([u for u in all_users if u.get("role") == "user"])
            classes = supabase.table("classes").select("*").execute().data
            for cls in classes:
                comp_data = supabase.table("class_completions").select("user_id").eq("class_id", cls["id"]).execute().data
                comp_count = len(comp_data)
                pct = int((comp_count / total_users * 100)) if total_users > 0 else 0
                with st.expander(f" {cls['title']}    {comp_count}/{total_users} students  ({pct}%)"):
                    col_prog, col_num = st.columns([5, 1])
                    with col_prog: st.progress(pct / 100)
                    with col_num: st.markdown(f"**{pct}%**")
                    if comp_count > 0:
                        with st.expander(f" {comp_count}  complete "):
                            for row in comp_data:
                                u = supabase.table("users").select("name, email").eq("id", row["user_id"]).execute().data
                                if u: st.caption(f" {u[0]['name']} ({u[0]['email']})")
                    st.divider()
                    ec_title = st.text_input("Title", value=cls["title"], key=f"ct_{cls['id']}")
                    ec_link = st.text_input("Live Link", value=cls.get("class_link",""), key=f"cl_{cls['id']}")
                    ev_link = st.text_input("Video Link", value=cls.get("recorded_video",""), key=f"cv_{cls['id']}")
                    ep_link = st.text_input("PDF Link", value=cls.get("notes_pdf",""), key=f"cp_{cls['id']}")
                    b1, b2 = st.columns(2)
                    with b1:
                        if st.button("Save Changes", key=f"cu_{cls['id']}", type="primary", use_container_width=True):
                            supabase.table("classes").update({"title": ec_title, "class_link": ec_link, "recorded_video": ev_link, "notes_pdf": ep_link}).eq("id", cls["id"]).execute()
                            st.success("Saved!"); st.rerun()
                    with b2:
                        if st.button("Remove Class", key=f"cd_{cls['id']}", use_container_width=True):
                            supabase.table("classes").delete().eq("id", cls["id"]).execute()
                            st.warning("Deleted!"); st.rerun()

    elif menu == "Manage Exams & Questions":
        ex_tab1, ex_tab2, ex_tab3, ex_tab4, ex_tab5, ex_tab6, ex_tab7 = st.tabs([
            " Exams Setup", " Add Questions", " Review Papers", " Bulk Upload (CSV)", " AI Gen", " Exam Folders", " Multi Round Exams"
        ])

        with ex_tab1:
            st.subheader("Setup Dynamic Exams")
            classes_list = supabase.table("classes").select("*").execute().data
            cls_options = {c["title"]: c["id"] for c in classes_list} if classes_list else {}
            exam_folder_options = build_exam_folder_options(fetch_exam_folders())
            with st.form("create_exam_form", clear_on_submit=True):
                sel_cls = st.selectbox("Link with Lesson Class", list(cls_options.keys()) or ["No classes yet"])
                sel_exam_folder = st.selectbox("Show under Exam Folder", list(exam_folder_options.keys()), key="create_exam_folder")
                e_title = st.text_input("Exam Sheet Name")
                e_duration = st.number_input("Exam Duration (Minutes)", min_value=1, max_value=900, value=30, help="15 hours varaku set cheyyachu (900 minutes).")
                e_pwd = st.text_input("Exam Password (Optional)", type="password")
                c_en = st.checkbox("Turn On Exam", value=True)
                c_ans = st.checkbox("Enable Answers Visibility")
                c_proctor = st.checkbox("Enable Camera Proctoring")
                section_json = st.text_area(
                    "Sectional Timer Config (Optional JSON)",
                    value='[{"title":"Aptitude","duration_mins":10,"question_count":10},{"title":"Programming","duration_mins":45,"question_count":2},{"title":"Reasoning","duration_mins":20,"question_count":15}]',
                    help="Blank ga vadilesthe normal single timer exam. question_count order lo questions sections ki assign avuthayi.",
                    height=110,
                )
                if st.form_submit_button(" Generate Exam Layout"):
                    if sel_cls in cls_options:
                        parsed_sections = parse_exam_section_config(section_json) if section_json.strip() else []
                        create_exam_record({
                            "class_id": cls_options[sel_cls], "title": e_title,
                            "duration_mins": int(e_duration),
                            "password": e_pwd.strip() if e_pwd.strip() else None,
                            "enabled": c_en, "show_answers": c_ans,
                            "proctoring_enabled": c_proctor,
                            "section_config": parsed_sections,
                            "folder_id": exam_folder_options.get(sel_exam_folder)
                        })
                        st.success("Exam Created!"); st.rerun()
            with st.expander("Programming Exam Builder - existing questions nundi create cheyyandi"):
                st.caption("Already add chesina programming questions select chesi, new exam create cheyyachu.")
                prog_questions = supabase.table("questions").select("*").eq("type", "programming").execute().data or []
                if not prog_questions:
                    st.info("Existing programming questions levu. First Add Questions tab lo programming question add cheyyandi.")
                else:
                    builder_cls = st.selectbox("Class select cheyyandi", list(cls_options.keys()) or ["No classes yet"], key="prog_builder_class")
                    builder_folder = st.selectbox("Exam Folder select cheyyandi", list(exam_folder_options.keys()), key="prog_builder_folder")
                    builder_title = st.text_input("Programming Exam Name", key="prog_builder_title")
                    builder_duration = st.number_input("Duration (Minutes)", min_value=1, max_value=900, value=60, key="prog_builder_duration", help="15 hours varaku set cheyyachu (900 minutes).")
                    builder_pwd = st.text_input("Password (Optional)", type="password", key="prog_builder_pwd")
                    builder_enabled = st.checkbox("Turn On Exam", value=True, key="prog_builder_enabled")
                    builder_show_answers = st.checkbox("Enable Answers Visibility", value=True, key="prog_builder_show_answers")
                    builder_proctor = st.checkbox("Enable Camera Proctoring", value=True, key="prog_builder_proctor")
                    builder_section_json = st.text_area(
                        "Sectional Timer Config (Optional JSON)",
                        value="",
                        key="prog_builder_sections",
                        help="Programming-only exam ki blank ok. Mixed exam ki sections JSON use cheyyandi.",
                        height=90,
                    )

                    q_options = {}
                    for q in prog_questions:
                        marks = get_question_max_marks(q)
                        label = f"{q.get('question', 'Untitled')}  |  {marks} marks  |  QID: {q.get('id')}"
                        q_options[label] = q
                    selected_labels = st.multiselect("Programming Questions select cheyyandi", list(q_options.keys()), key="prog_builder_questions")

                    if st.button("Selected Questions tho Exam Create", type="primary", use_container_width=True, key="prog_builder_create"):
                        if builder_cls not in cls_options:
                            st.error("Class select cheyyandi.")
                        elif not builder_title.strip():
                            st.error("Exam name enter cheyyandi.")
                        elif not selected_labels:
                            st.error("At least one programming question select cheyyandi.")
                        else:
                            created = create_exam_record({
                                "class_id": cls_options[builder_cls],
                                "title": builder_title.strip(),
                                "duration_mins": int(builder_duration),
                                "password": builder_pwd.strip() if builder_pwd.strip() else None,
                                "enabled": builder_enabled,
                                "show_answers": builder_show_answers,
                                "proctoring_enabled": builder_proctor,
                                "section_config": parse_exam_section_config(builder_section_json) if builder_section_json.strip() else [],
                                "folder_id": exam_folder_options.get(builder_folder),
                            })
                            new_exam = created[0] if created else None
                            if not new_exam:
                                matches = supabase.table("exams").select("*").eq("title", builder_title.strip()).eq("class_id", cls_options[builder_cls]).execute().data or []
                                new_exam = matches[-1] if matches else None
                            if not new_exam:
                                st.error("Exam create ayindi kani id fetch avvaledu. Page refresh chesi check cheyyandi.")
                            else:
                                for label in selected_labels:
                                    src = q_options[label]
                                    supabase.table("questions").insert({
                                        "exam_id": new_exam["id"],
                                        "question": src.get("question", ""),
                                        "type": "programming",
                                        "option_a": src.get("option_a", ""),
                                        "option_b": src.get("option_b", ""),
                                        "option_c": src.get("option_c", ""),
                                        "option_d": src.get("option_d", ""),
                                        "correct_answer": src.get("correct_answer", "AUTO"),
                                        "hint": src.get("hint", ""),
                                        "image_url": src.get("image_url"),
                                        "explanation": src.get("explanation"),
                                    }).execute()
                                st.success("Programming exam create ayyindi!")
                                st.rerun()
            st.divider()
            st.write("###  Live Exam Controls")
            exams_all = supabase.table("exams").select("*").execute().data
            for ex in exams_all:
                with st.container(border=True):
                    st.markdown(f"####  **{ex['title']}**")
                    col_e1, col_e2, col_e3 = st.columns([2, 2, 2])
                    with col_e1:
                        updated_dur = st.number_input("Duration (Mins)", min_value=1, max_value=900, value=min(int(ex.get("duration_mins",30) or 30), 900), key=f"dur_{ex['id']}", help="15 hours varaku set cheyyachu (900 minutes).")
                    with col_e2:
                        updated_pwd = st.text_input("Password", value=str(ex.get("password","") or ""), key=f"pwd_ed_{ex['id']}")
                    with col_e3:
                        t_active = st.toggle("Active", value=ex["enabled"], key=f"tog_en_{ex['id']}")
                        t_ans = st.toggle("Show Answers", value=ex["show_answers"], key=f"tog_ans_{ex['id']}")
                        t_proctor = st.toggle("Camera Proctoring", value=bool(ex.get("proctoring_enabled")), key=f"tog_proc_{ex['id']}")
                    current_sections = parse_exam_section_config(ex.get("section_config"))
                    sections_text = json.dumps(current_sections, ensure_ascii=False, indent=2) if current_sections else ""
                    updated_sections_text = st.text_area(
                        "Sectional Timer Config JSON",
                        value=sections_text,
                        key=f"sections_{ex['id']}",
                        height=110,
                        placeholder='[{"title":"Aptitude","duration_mins":10,"question_count":10},{"title":"Programming","duration_mins":45,"question_count":2}]',
                    )
                    col_btn1, col_btn2 = st.columns(2)
                    with col_btn1:
                        if st.button("Save", key=f"up_ex_{ex['id']}", type="primary", use_container_width=True):
                            old_pwd = str(ex.get("password") or "")
                            new_pwd = updated_pwd.strip()
                            parsed_sections = parse_exam_section_config(updated_sections_text) if updated_sections_text.strip() else []
                            update_exam_record(ex["id"], {
                                "duration_mins": int(updated_dur),
                                "password": new_pwd if new_pwd else None,
                                "enabled": t_active, "show_answers": t_ans,
                                "proctoring_enabled": t_proctor,
                                "section_config": parsed_sections,
                            })
                            if new_pwd and new_pwd != old_pwd:
                                send_notification(f" '{ex['title']}' exam  password set : {new_pwd}")
                            st.success("Updated!"); st.rerun()
                    with col_btn2:
                        if st.button("Delete Exam", key=f"del_ex_{ex['id']}", type="secondary", use_container_width=True):
                            supabase.table("exams").delete().eq("id", ex["id"]).execute()
                            st.warning("Deleted!"); st.rerun()

        with ex_tab2:
            st.subheader("Add Questions")
            exams_q = supabase.table("exams").select("*").execute().data
            ex_options = {e["title"]: e["id"] for e in exams_q} if exams_q else {}
            st.markdown("####  Add Question")
            sel_ex = st.selectbox("Select Exam", list(ex_options.keys()) or ["No exams yet"], key="add_q_exam")

            #  Image upload / URL 
            st.caption("Upload image to automatically extract question and options.")
            img_col1, img_col2 = st.columns(2)
            with img_col1:
                img_url_input = st.text_input("Image URL", key="add_img_url", placeholder="https://...")
            with img_col2:
                img_file = st.file_uploader("Image upload", type=["jpg","jpeg","png","gif","webp"], key="add_img_file")

            # Image preview
            if img_file:
                st.image(img_file, width=380)
            elif img_url_input.strip():
                try: st.image(img_url_input.strip(), width=380)
                except Exception: pass

            extract_source = img_file if img_file else (img_url_input.strip() or None)
            if extract_source:
                if st.button("Extract Question from Image", use_container_width=True, key="extract_ocr_btn"):
                    with st.spinner("OCR processing..."):
                        extracted = extract_question_from_image(extract_source)
                    if extracted:
                        # Set widget keys DIRECTLY before they render  widgets not yet on screen
                        st.session_state["aq_q_text"] = extracted.get("question", "")
                        st.session_state["aq_opt_A"]  = extracted.get("option_a", "")
                        st.session_state["aq_opt_B"]  = extracted.get("option_b", "")
                        st.session_state["aq_opt_C"]  = extracted.get("option_c", "")
                        st.session_state["aq_opt_D"]  = extracted.get("option_d", "")
                        st.session_state["aq_hint"]   = extracted.get("hint", "")
                        ans = extracted.get("correct_answer","").strip().upper()
                        if ans in ["A","B","C","D"]:
                            st.session_state["aq_correct_lbl"] = ans
                        q_type_ocr = extracted.get("type","mcq")
                        st.session_state["aq_q_type_idx"] = ["mcq","blank","programming"].index(q_type_ocr) if q_type_ocr in ["mcq","blank","programming"] else 0
                        st.rerun()  # ONE rerun  widgets will now render with pre-filled values

            st.divider()

            #  Question Type 
            type_idx_default = st.session_state.get("aq_q_type_idx", 0)
            q_type = st.selectbox("Question Type", ["mcq","blank","programming"],
                                   index=type_idx_default, key="aq_q_type")

            #  Question Text 
            q_text = st.text_area("Question / Title", key="aq_q_text")

            opt_vals = {"A": "", "B": "", "C": "", "D": ""}
            correct_lbl = st.session_state.get("aq_correct_lbl", "")
            h_text = st.text_input(" Hint", key="aq_hint")
            exp_text = ""
            prog_description = ""
            prog_test_cases = []

            if q_type == "programming":
                prog_language_label = st.selectbox("Programming Language", list(PROGRAMMING_LANGUAGE_LABELS.keys()), key="aq_prog_language")
                prog_language = PROGRAMMING_LANGUAGE_LABELS[prog_language_label]
                prog_description = st.text_area("Programming Description", key="aq_prog_desc",
                    placeholder="Problem statement clear ga rayandi:\n- Task / objective\n- Input format\n- Output format\n- Constraints\n- Sample explanation optional")
                st.markdown("**Test Cases & Marks**")
                tc_count = st.number_input("Number of test cases", min_value=1, max_value=10, value=3, step=1, key="aq_tc_count")
                for idx in range(int(tc_count)):
                    with st.container(border=True):
                        st.markdown(f"##### Test Case {idx + 1}")
                        c_in, c_out, c_marks, c_hidden = st.columns([3, 3, 1, 1])
                        with c_in:
                            tc_input = st.text_area("Input", key=f"aq_tc_input_{idx}", height=90)
                        with c_out:
                            tc_output = st.text_area("Expected Output", key=f"aq_tc_output_{idx}", height=90)
                        with c_marks:
                            tc_marks = st.number_input("Marks", min_value=1, max_value=100, value=1, step=1, key=f"aq_tc_marks_{idx}")
                        with c_hidden:
                            tc_hidden = st.checkbox("Hidden", value=(idx > 0), key=f"aq_tc_hidden_{idx}")
                        prog_test_cases.append({
                            "input": tc_input,
                            "expected_output": tc_output,
                            "marks": int(tc_marks),
                            "hidden": bool(tc_hidden),
                        })
                correct_lbl = "AUTO"
            else:
                #  4 Options +  Set Correct button 
                if q_type == "mcq":
                    st.markdown("**Options**   option  ** Set Correct** ")
                    for lbl in ["A","B","C","D"]:
                        c1, c2, c3 = st.columns([1, 6, 2])
                        with c1:
                            st.markdown(f"<div style='padding-top:8px;font-weight:700;font-size:1rem;'>{lbl}.</div>",
                                        unsafe_allow_html=True)
                        with c2:
                            v = st.text_input(f"Option {lbl}", key=f"aq_opt_{lbl}", label_visibility="collapsed")
                            opt_vals[lbl] = v
                        with c3:
                            cur_correct = st.session_state.get("aq_correct_lbl", "")
                            is_correct = cur_correct == lbl
                            if st.button(
                                " Correct" if is_correct else " Set Correct",
                                key=f"aq_setcor_{lbl}",
                                use_container_width=True,
                                type="primary" if is_correct else "secondary"
                            ):
                                st.session_state["aq_correct_lbl"] = lbl
                                st.rerun()

                    correct_lbl = st.session_state.get("aq_correct_lbl", "")
                    if correct_lbl:
                        opt_text = opt_vals.get(correct_lbl, "")
                        st.info(f" Correct Answer: **{correct_lbl}. {opt_text}**")
                else:
                    correct_lbl = st.text_input("Correct Answer", key="aq_blank_correct")

                exp_text = st.text_area(" Answer Explanation (optional)", key="aq_explanation",
                                         placeholder="   correct  ...")

            if st.button("Add Question", type="primary", key="add_q_btn", use_container_width=True):
                if sel_ex in ex_options and q_text.strip():
                    final_img_url = None
                    if img_file:
                        final_img_url = upload_image_to_imgbb(img_file)
                    elif img_url_input.strip():
                        final_img_url = img_url_input.strip()
                    explanation_value = make_programming_meta(prog_description, prog_test_cases, prog_language) if q_type == "programming" else (exp_text.strip() if exp_text.strip() else None)
                    supabase.table("questions").insert({
                        "exam_id": ex_options[sel_ex],
                        "question": q_text,
                        "type": q_type,
                        "option_a": opt_vals.get("A",""),
                        "option_b": opt_vals.get("B",""),
                        "option_c": opt_vals.get("C",""),
                        "option_d": opt_vals.get("D",""),
                        "correct_answer": correct_lbl,
                        "hint": h_text,
                        "image_url": final_img_url,
                        "explanation": explanation_value
                    }).execute()
                    # Clear all aq_ keys
                    for k in list(st.session_state.keys()):
                        if str(k).startswith("aq_tc_"):
                            st.session_state.pop(k, None)
                    for k in ["aq_q_text","aq_opt_A","aq_opt_B","aq_opt_C","aq_opt_D",
                              "aq_hint","aq_explanation","aq_prog_desc","aq_blank_correct","aq_correct_lbl","aq_q_type_idx"]:
                        st.session_state.pop(k, None)
                    st.success("Question added.")
                    st.rerun()
                else:
                    st.error("Exam select   question text enter .")

        with ex_tab4:
            st.subheader("Bulk Upload Questions (CSV)")
            exams = supabase.table("exams").select("id, title").execute()
            exam_options = {ex["title"]: ex["id"] for ex in exams.data} if exams.data else {}
            if exam_options:
                selected_exam = st.selectbox("Select Exam:", list(exam_options.keys()))
                exam_id_bulk = exam_options[selected_exam]
                st.caption("CSV supports mcq, blank, and programming rows. Programming rows can use test_input_1/test_output_1/test_marks_1/test_hidden_1 columns up to 10, or a test_cases_json list.")
                with st.expander("Programming CSV format - quick guide", expanded=True):
                    st.markdown(
                        """
                        **Required columns:** `question`, `type`

                        For programming questions:
                        `type` should be `programming`, `language` should be `java`, `c`, or `python`, and `description` should include the full problem statement, input format, output format, and constraints.

                        Add test cases using repeated columns:
                        `test_input_1`, `test_output_1`, `test_marks_1`, `test_hidden_1`.
                        Continue as `test_input_2`, `test_output_2` ... up to 10 cases.

                        Alternative advanced format:
                        Put a JSON list in `test_cases_json`, for example `[{"input":"2\\n4 5","expected_output":"9","marks":2,"hidden":false}]`.
                        """
                    )
                    st.dataframe(pd.DataFrame([
                        {"Column": "question", "Programming example": "Sum of two numbers"},
                        {"Column": "type", "Programming example": "programming"},
                        {"Column": "language", "Programming example": "java"},
                        {"Column": "description", "Programming example": "Read two integers and print their sum. Input: a b. Output: sum."},
                        {"Column": "test_input_1", "Programming example": "4 5"},
                        {"Column": "test_output_1", "Programming example": "9"},
                        {"Column": "test_marks_1", "Programming example": "2"},
                        {"Column": "test_hidden_1", "Programming example": "false"},
                    ]), hide_index=True, use_container_width=True)
                template_rows = pd.DataFrame([
                    {
                        "question": "Two Sum",
                        "type": "programming",
                        "language": "java",
                        "description": "Read n and n integers, print the sum of the two numbers.",
                        "test_input_1": "2\n4 5",
                        "test_output_1": "9",
                        "test_marks_1": 2,
                        "test_hidden_1": "false",
                        "test_input_2": "2\n10 15",
                        "test_output_2": "25",
                        "test_marks_2": 3,
                        "test_hidden_2": "true",
                        "correct_answer": "AUTO",
                    },
                    {
                        "question": "Capital of India?",
                        "type": "mcq",
                        "option_a": "Delhi",
                        "option_b": "Mumbai",
                        "option_c": "Chennai",
                        "option_d": "Kolkata",
                        "correct_answer": "A",
                        "hint": "",
                        "explanation": "",
                    },
                ])
                st.download_button(
                    "Download CSV Template",
                    template_rows.to_csv(index=False).encode("utf-8"),
                    file_name="question_upload_template.csv",
                    mime="text/csv",
                    use_container_width=True,
                )
                uploaded_file = st.file_uploader("Upload CSV", type=["csv"])
                if uploaded_file is not None:
                    import io as _io
                    try:
                        raw = uploaded_file.read()
                        try: df = pd.read_csv(_io.StringIO(raw.decode("utf-8")))
                        except UnicodeDecodeError: df = pd.read_csv(_io.StringIO(raw.decode("latin1")))
                        required = ["question", "type"]
                        missing = [col for col in required if col not in df.columns]
                        if missing:
                            st.error(f"CSV   columns : {missing}")
                            st.caption("Expected minimum: question, type. MCQ/blank use correct_answer. Programming use language, description, and test cases columns.")
                        else:
                            df = df.fillna("")
                            st.success(f" {len(df)} rows loaded!")
                            st.write("Preview:", df.head())
                            errors = []
                            for row_no, row in df.iterrows():
                                q_type = str(row.get("type", "")).strip().lower()
                                if not str(row.get("question", "")).strip():
                                    errors.append(f"Row {row_no + 2}: question empty")
                                if q_type == "programming" and not parse_programming_csv_cases(row):
                                    errors.append(f"Row {row_no + 2}: programming row needs at least one test case")
                                if q_type in {"mcq", "blank"} and not str(row.get("correct_answer", "")).strip():
                                    errors.append(f"Row {row_no + 2}: correct_answer required")
                            if errors:
                                st.error("CSV validation failed:")
                                for err in errors[:12]:
                                    st.caption(err)
                                if len(errors) > 12:
                                    st.caption(f"...and {len(errors) - 12} more")
                                st.stop()
                            if st.button("Upload to DB"):
                                try:
                                    for _, row in df.iterrows():
                                        supabase.table("questions").insert(
                                            build_question_payload_from_csv_row(row, exam_id_bulk)
                                        ).execute()
                                    st.success(f" {len(df)} questions uploaded!")
                                except Exception as e:
                                    st.error(f"Upload Error: {e}")
                    except Exception as e:
                        st.error(f"CSV : {e}")
            else:
                st.warning("Create an exam first.")

        with ex_tab5:
            st.subheader("AI Question Generator (Gemini)")
            exams_ai = supabase.table("exams").select("*").execute().data
            ai_ex_options = {e["title"]: e["id"] for e in exams_ai} if exams_ai else {}
            sel_ai_ex = st.selectbox("Save to Exam", list(ai_ex_options.keys()) or ["No exams yet"], key="ai_gen_exam")
            lesson_text = st.text_area("Paste Lesson Content here:")
            if st.button("Generate Questions"):
                if not lesson_text.strip():
                    st.warning("Paste lesson text first.")
                else:
                    try:
                        prompt = (
                            "Convert this text into 5 MCQ questions in JSON format. "
                            "Return ONLY a JSON array, no markdown fences. Each item must have: "
                            "question, option_a, option_b, option_c, option_d, correct_answer, explanation. "
                            "correct_answer must exactly match the text of the correct option. "
                            "explanation should be 1-2 sentences explaining why that answer is correct. "
                            f"Text: {lesson_text}"
                        )
                        model = genai.GenerativeModel(model_name="gemini-2.0-flash-lite")
                        response = model.generate_content(prompt)
                        raw_text = response.text.strip().strip("`")
                        if raw_text.lower().startswith("json"):
                            raw_text = raw_text[4:].strip()
                        parsed_qs = json.loads(raw_text)
                        st.session_state.ai_generated_qs = parsed_qs
                        st.success(f" {len(parsed_qs)} questions generated!")
                    except Exception as e:
                        st.error(f"AI Error: {e}")

            if st.session_state.get("ai_generated_qs"):
                st.subheader("Generated Questions Preview")
                for gi, gq in enumerate(st.session_state.ai_generated_qs):
                    with st.container(border=True):
                        st.markdown(f"**Q{gi+1}. {gq.get('question','')}**")
                        st.caption(f"A: {gq.get('option_a','')} | B: {gq.get('option_b','')} | C: {gq.get('option_c','')} | D: {gq.get('option_d','')}")
                        st.caption(f" Correct: {gq.get('correct_answer','')}")
                        if gq.get("explanation"):
                            st.caption(f" {gq.get('explanation','')}")
                if sel_ai_ex in ai_ex_options:
                    if st.button("Save to DB", type="primary", use_container_width=True):
                        try:
                            for gq in st.session_state.ai_generated_qs:
                                supabase.table("questions").insert({
                                    "exam_id": ai_ex_options[sel_ai_ex],
                                    "question": gq.get("question",""), "type": "mcq",
                                    "option_a": gq.get("option_a",""), "option_b": gq.get("option_b",""),
                                    "option_c": gq.get("option_c",""), "option_d": gq.get("option_d",""),
                                    "correct_answer": gq.get("correct_answer",""), "hint": "",
                                    "explanation": gq.get("explanation","") or None
                                }).execute()
                            st.success(f" {len(st.session_state.ai_generated_qs)} questions saved!")
                            st.session_state.ai_generated_qs = None
                            st.rerun()
                        except Exception as e:
                            st.error(f"Save Error: {e}")

        with ex_tab3:
            st.subheader("Review Existing Exam Papers")
            exams_all_edit = supabase.table("exams").select("*").execute().data
            if not exams_all_edit:
                st.info("No exams created yet.")
            else:
                exam_edit_options = {ex["title"]: ex["id"] for ex in exams_all_edit}
                selected_exam_title = st.selectbox("Choose Exam to Review", list(exam_edit_options.keys()))
                selected_exam_id = exam_edit_options[selected_exam_title]
                current_questions = supabase.table("questions").select("*").eq("exam_id", selected_exam_id).execute().data

                exp_reqs = supabase.table("explain_requests").select("*").eq("exam_id", selected_exam_id).execute().data
                q_requesters = {}
                for req in exp_reqs:
                    qids = json.loads(req.get("question_ids") or "[]")
                    uinfo = supabase.table("users").select("name").eq("id", req["user_id"]).execute().data
                    uname = uinfo[0]["name"] if uinfo else "Unknown"
                    for qid in qids:
                        q_requesters.setdefault(qid, [])
                        if uname not in q_requesters[qid]:
                            q_requesters[qid].append(uname)

                marked_q_ids = set(q_requesters.keys())
                marked_questions = [q for q in current_questions if q["id"] in marked_q_ids]

                st.write(f"### Questions in **{selected_exam_title}** ({len(current_questions)} total)")
                col1, col2, col3 = st.columns(3)
                col1.metric("Total Questions", len(current_questions))
                col2.metric("Explain Requested", len(marked_questions))
                col3.metric("Students Requested", len(set(req["user_id"] for req in exp_reqs)))

                dl_col1, dl_col2 = st.columns(2)
                with dl_col1:
                    if current_questions and st.button("All Questions PPT", use_container_width=True):
                        with st.spinner("PPT generate ..."):
                            ppt_bytes = generate_exam_ppt(current_questions, selected_exam_title, q_requesters=q_requesters)
                            if ppt_bytes:
                                st.download_button(" Download", data=ppt_bytes,
                                    file_name=f"{selected_exam_title[:25]}_all.pptx",
                                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                                    key="ppt_all_btn")
                with dl_col2:
                    if marked_questions and st.button(f" Marked Questions PPT ({len(marked_questions)})", use_container_width=True, type="primary"):
                        with st.spinner("PPT generate ..."):
                            ppt_bytes = generate_exam_ppt(marked_questions, f"{selected_exam_title} - Explain", q_requesters=q_requesters)
                            if ppt_bytes:
                                st.download_button(" Download", data=ppt_bytes,
                                    file_name=f"{selected_exam_title[:25]}_marked.pptx",
                                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                                    key="ppt_marked_btn")

                st.divider()
                for idx, q in enumerate(current_questions):
                    with st.container(border=True):
                        col_q1, col_q2, col_q3 = st.columns([5, 1, 1])
                        with col_q1:
                            st.markdown(f"**Q{idx+1}. {q['question']}** `({str(q['type']).upper()})`")
                            if q.get("image_url"): st.image(q["image_url"], width=200)
                            if q["type"] == "mcq":
                                st.caption(f"A: {q['option_a']} | B: {q['option_b']} | C: {q['option_c']} | D: {q['option_d']}")
                            st.caption(f" Answer: {q['correct_answer']} |  Hint: {q.get('hint','')}")
                            if q.get("explanation"):
                                st.caption(f" Explanation: {q['explanation']}")
                        with col_q2:
                            if st.button("Edit", key=f"edit_btn_{q['id']}", use_container_width=True):
                                st.session_state[f"editing_{q['id']}"] = True; st.rerun()
                        with col_q3:
                            if st.button("Delete", key=f"del_q_{q['id']}", type="secondary", use_container_width=True):
                                supabase.table("questions").delete().eq("id", q["id"]).execute()
                                st.success(f"Q{idx+1} Deleted!"); st.rerun()

                        if st.session_state.get(f"editing_{q['id']}", False):
                            with st.form(key=f"edit_form_{q['id']}"):
                                st.markdown("#####  Question Edit")
                                eq_type = st.selectbox("Type", ["mcq","blank","programming"],
                                    index=["mcq","blank","programming"].index(q["type"]) if q["type"] in ["mcq","blank","programming"] else 0,
                                    key=f"eq_type_{q['id']}")
                                eq_text = st.text_area("Question / Title", value=q["question"], key=f"eq_text_{q['id']}", height=90)
                                eq_hint = st.text_input("Hint", value=q.get("hint","") or "", key=f"eq_hint_{q['id']}")
                                eq_img_url = st.text_input("Image URL", value=q.get("image_url","") or "", key=f"eq_img_{q['id']}")
                                eq_a = eq_b = eq_c = eq_d = ""
                                eq_ans = q.get("correct_answer", "") or ""
                                eq_explanation = q.get("explanation", "") or ""
                                if eq_type == "programming":
                                    prog_meta = get_programming_meta(q)
                                    current_prog_label = get_programming_language_meta(prog_meta.get("language", "java"))["label"]
                                    prog_labels = list(PROGRAMMING_LANGUAGE_LABELS.keys())
                                    eq_prog_label = st.selectbox(
                                        "Programming Language",
                                        prog_labels,
                                        index=prog_labels.index(current_prog_label) if current_prog_label in prog_labels else 0,
                                        key=f"eq_prog_lang_{q['id']}",
                                    )
                                    eq_prog_language = PROGRAMMING_LANGUAGE_LABELS[eq_prog_label]
                                    eq_prog_description = st.text_area(
                                        "Problem Description",
                                        value=prog_meta.get("description", ""),
                                        key=f"eq_prog_desc_{q['id']}",
                                        height=180,
                                        placeholder="Problem statement, input format, output format, constraints separate ga rayandi.",
                                    )
                                    old_cases = prog_meta.get("test_cases", []) or []
                                    eq_prog_cases = []
                                    eq_case_count = st.number_input("Test cases", min_value=1, max_value=10, value=max(1, min(len(old_cases) or 3, 10)), step=1, key=f"eq_tc_count_{q['id']}")
                                    for tc_idx in range(int(eq_case_count)):
                                        old_tc = old_cases[tc_idx] if tc_idx < len(old_cases) else {}
                                        with st.container(border=True):
                                            st.markdown(f"##### Test Case {tc_idx + 1}")
                                            tc_in_col, tc_out_col, tc_marks_col, tc_hidden_col = st.columns([3, 3, 1, 1])
                                            with tc_in_col:
                                                tc_input = st.text_area("Input", value=old_tc.get("input", ""), key=f"eq_tc_input_{q['id']}_{tc_idx}", height=80)
                                            with tc_out_col:
                                                tc_output = st.text_area("Expected Output", value=old_tc.get("expected_output", ""), key=f"eq_tc_output_{q['id']}_{tc_idx}", height=80)
                                            with tc_marks_col:
                                                tc_marks = st.number_input("Marks", min_value=1, max_value=100, value=int(old_tc.get("marks", 1) or 1), step=1, key=f"eq_tc_marks_{q['id']}_{tc_idx}")
                                            with tc_hidden_col:
                                                tc_hidden = st.checkbox("Hidden", value=bool(old_tc.get("hidden", tc_idx > 0)), key=f"eq_tc_hidden_{q['id']}_{tc_idx}")
                                            eq_prog_cases.append({"input": tc_input, "expected_output": tc_output, "marks": int(tc_marks), "hidden": bool(tc_hidden)})
                                    eq_ans = "AUTO"
                                    eq_explanation = make_programming_meta(eq_prog_description, eq_prog_cases, eq_prog_language)
                                else:
                                    col1, col2 = st.columns(2)
                                    with col1:
                                        eq_a = st.text_input("Option A", value=q.get("option_a","") or "", key=f"eq_a_{q['id']}")
                                        eq_b = st.text_input("Option B", value=q.get("option_b","") or "", key=f"eq_b_{q['id']}")
                                    with col2:
                                        eq_c = st.text_input("Option C", value=q.get("option_c","") or "", key=f"eq_c_{q['id']}")
                                        eq_d = st.text_input("Option D", value=q.get("option_d","") or "", key=f"eq_d_{q['id']}")
                                    eq_ans = st.text_input("Correct Answer", value=q.get("correct_answer","") or "", key=f"eq_ans_{q['id']}")
                                    eq_explanation = st.text_area("Answer Explanation", value=q.get("explanation","") or "", key=f"eq_exp_{q['id']}")
                                if q.get("image_url"): st.image(q["image_url"], width=150)
                                save_col, cancel_col = st.columns(2)
                                with save_col:
                                    saved = st.form_submit_button(" Save", use_container_width=True, type="primary")
                                with cancel_col:
                                    cancelled = st.form_submit_button(" Cancel", use_container_width=True)
                                if saved:
                                    supabase.table("questions").update({
                                        "type": eq_type, "question": eq_text,
                                        "option_a": eq_a, "option_b": eq_b, "option_c": eq_c, "option_d": eq_d,
                                        "correct_answer": eq_ans, "hint": eq_hint,
                                        "explanation": eq_explanation.strip() if eq_explanation.strip() else None,
                                        "image_url": eq_img_url.strip() if eq_img_url.strip() else None
                                    }).eq("id", q["id"]).execute()
                                    st.session_state[f"editing_{q['id']}"] = False
                                    st.success("Updated!"); st.rerun()
                                if cancelled:
                                    st.session_state[f"editing_{q['id']}"] = False; st.rerun()

                st.divider()
                st.markdown("####  Quick Add Question")
                with st.form("quick_add_question_form", clear_on_submit=True):
                    q_type_new = st.selectbox("Type", ["mcq","blank","programming"], key="new_q_type")
                    q_text_new = st.text_area("Question Text", key="new_q_text")
                    col_opts1, col_opts2 = st.columns(2)
                    with col_opts1:
                        a_new = st.text_input("Option A", key="new_a"); b_new = st.text_input("Option B", key="new_b")
                    with col_opts2:
                        c_new = st.text_input("Option C", key="new_c"); d_new = st.text_input("Option D", key="new_d")
                    h_text_new = st.text_input("Hint", key="new_hint")
                    c_ans_new = st.text_input("Correct Answer", key="new_ans")
                    exp_new = st.text_area(" Explanation (optional)", key="new_explanation")
                    if st.form_submit_button(" Add Question"):
                        if q_text_new.strip():
                            supabase.table("questions").insert({
                                "exam_id": selected_exam_id, "question": q_text_new, "type": q_type_new,
                                "option_a": a_new, "option_b": b_new, "option_c": c_new, "option_d": d_new,
                                "correct_answer": c_ans_new if q_type_new != "programming" else "Manual Review Required",
                                "hint": h_text_new,
                                "explanation": exp_new.strip() if exp_new.strip() else None
                            }).execute()
                            st.success("Added!"); st.rerun()
                        else:
                            st.error("Question text empty!")

        with ex_tab6:
            show_admin_exam_folders_tab()

        with ex_tab7:
            show_admin_exam_series_tab()

    elif menu == "Student Results & Ranks":
        r_tab1, r_tab2, r_tab3, r_tab4, r_tab5, r_tab6, r_tab7, r_tab8, r_tab9 = st.tabs([
            " Leaderboards", " Manual Evaluation", " Score Summary",
            " Re-Exam Requests", " Explain Requests", " Hint Requests", " Test Case Reports", " Attendance", " Live Programming Exams"
        ])

        with r_tab1:
            st.title("Leaderboard")
            exams = supabase.table("exams").select("*").execute().data
            if exams:
                sel_ex_lb = st.selectbox("Select Exam", [e["title"] for e in exams])
                target_ex = next((e for e in exams if e["title"] == sel_ex_lb), None)
                if target_ex:
                    board = get_exam_leaderboard(target_ex["id"])
                    if board:
                        for rank, st_row in enumerate(board):
                            medal = "" if rank==0 else "" if rank==1 else "" if rank==2 else f"{rank+1}."
                            st.write(f"{medal} **{st_row['Name']}** ({st_row['Email']})  Score: **{st_row['Score']}**")
                    else:
                        st.info("No attempts yet.")

        with r_tab2:
            st.title("Manual Evaluator")
            attempts_to_eval = supabase.table("exam_attempts").select("*").execute().data
            if not attempts_to_eval:
                st.info("No submissions available.")
            else:
                for att in attempts_to_eval:
                    u_data = supabase.table("users").select("*").eq("id", att["user_id"]).execute().data
                    e_data = supabase.table("exams").select("*").eq("id", att["exam_id"]).execute().data
                    if u_data and e_data:
                        with st.container(border=True):
                            col_s1, col_s2 = st.columns([3, 1])
                            with col_s1:
                                st.markdown(f"#####  **{u_data[0]['name']}** |  **{e_data[0]['title']}**")
                                st.code(att.get("submitted_answers","# No code submitted."), language="python")
                            with col_s2:
                                current_score = int(att.get("score") or 0)
                                try:
                                    exam_questions = supabase.table("questions").select("*").eq("exam_id", att["exam_id"]).execute().data or []
                                    exam_max_score = get_exam_max_marks(exam_questions)
                                except Exception:
                                    exam_max_score = 100
                                score_max = max(100, int(exam_max_score or 0), current_score)
                                new_score = st.number_input("Score", min_value=0, max_value=score_max, value=current_score, step=1, key=f"score_in_{att['id']}")
                                if st.button("Save", key=f"btn_score_{att['id']}", type="primary", use_container_width=True):
                                    supabase.table("exam_attempts").update({"score": new_score}).eq("id", att["id"]).execute()
                                    st.success("Score Saved!"); st.rerun()

        with r_tab3:
            st.title("Score Logger")
            attempts = supabase.table("exam_attempts").select("*").execute().data
            if not attempts:
                st.warning("No submissions yet.")
            else:
                for att in attempts:
                    u_prof = supabase.table("users").select("*").eq("id", att["user_id"]).execute().data
                    e_prof = supabase.table("exams").select("*").eq("id", att["exam_id"]).execute().data
                    if u_prof and e_prof:
                        st.markdown(f" **{u_prof[0]['name']}**  **{e_prof[0]['title']}**  Score: **{att['score']}**")
                        st.divider()

        with r_tab4:
            st.title("Re-Exam Requests")
            requests_data = supabase.table("exam_retake_requests").select("*").eq("status","pending").order("requested_at",desc=True).execute().data
            if not requests_data:
                st.info("Pending requests .")
            else:
                for req in requests_data:
                    u_info = supabase.table("users").select("name, email").eq("id", req["user_id"]).execute().data
                    e_info = supabase.table("exams").select("title").eq("id", req["exam_id"]).execute().data
                    if u_info and e_info:
                        with st.container(border=True):
                            col1, col2, col3 = st.columns([4, 1, 1])
                            with col1:
                                st.markdown(f"** {u_info[0]['name']}** ({u_info[0]['email']})")
                                st.caption(f" {e_info[0]['title']} | {req['requested_at'][:10]}")
                            with col2:
                                if st.button("Approve", key=f"apr_{req['id']}", type="primary", use_container_width=True):
                                    supabase.table("exam_retake_requests").update({"status":"approved","reviewed_at":"now()"}).eq("id",req["id"]).execute()
                                    st.success("Approved!"); st.rerun()
                            with col3:
                                if st.button("Reject", key=f"rej_{req['id']}", use_container_width=True):
                                    supabase.table("exam_retake_requests").update({"status":"rejected","reviewed_at":"now()"}).eq("id",req["id"]).execute()
                                    st.warning("Rejected!"); st.rerun()

        with r_tab5:
            st.title("Explain Requests")
            exp_requests = supabase.table("explain_requests").select("*").order("created_at",desc=True).execute().data
            if not exp_requests:
                st.info("Explain requests .")
            else:
                for req in exp_requests:
                    u_info = supabase.table("users").select("name, email").eq("id", req["user_id"]).execute().data
                    e_info = supabase.table("exams").select("title").eq("id", req["exam_id"]).execute().data
                    uname = u_info[0]["name"] if u_info else "Unknown"
                    ename = e_info[0]["title"] if e_info else "Unknown Exam"
                    qids = json.loads(req["question_ids"]) if req.get("question_ids") else []
                    status = req.get("status","pending")
                    status_color = {"pending":"","done":"","rejected":""}.get(status,"")
                    with st.container(border=True):
                        col1, col2, col3 = st.columns([4, 1, 1])
                        with col1:
                            st.markdown(f"** {uname}** ({u_info[0]['email'] if u_info else ''})")
                            st.caption(f" {ename} | {len(qids)} questions | {status_color} {status} | {str(req.get('created_at',''))[:10]}")
                        with col2:
                            if qids:
                                marked_qs = supabase.table("questions").select("*").in_("id", qids).execute().data
                                if marked_qs:
                                    all_reqs_for_exam = supabase.table("explain_requests").select("*").eq("exam_id",req["exam_id"]).execute().data
                                    qr = {}
                                    for r2 in all_reqs_for_exam:
                                        qids2 = json.loads(r2.get("question_ids") or "[]")
                                        ui2 = supabase.table("users").select("name").eq("id",r2["user_id"]).execute().data
                                        un2 = ui2[0]["name"] if ui2 else "Unknown"
                                        for qid2 in qids2:
                                            qr.setdefault(qid2, [])
                                            if un2 not in qr[qid2]: qr[qid2].append(un2)
                                    ppt_bytes = generate_exam_ppt(marked_qs, f"{ename} - Explain", q_requesters=qr)
                                    if ppt_bytes:
                                        st.download_button(" PPT", data=ppt_bytes,
                                            file_name=f"{ename[:20]}_explain.pptx",
                                            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                                            key=f"exp_ppt_{req['id']}")
                        with col3:
                            if status == "pending":
                                if st.button("Done", key=f"exp_done_{req['id']}", use_container_width=True, type="primary"):
                                    supabase.table("explain_requests").update({"status":"done"}).eq("id",req["id"]).execute()
                                    st.rerun()

        with r_tab6:
            st.title("Hint Requests")
            try:
                hint_requests = supabase.table("hint_requests").select("*").order("created_at", desc=True).execute().data or []
            except Exception as e:
                hint_requests = []
                st.warning(f"Hint requests table ready ledu: {e}")
            if not hint_requests:
                st.info("Hint requests levu.")
            for req in hint_requests:
                u_info = supabase.table("users").select("name, email").eq("id", req["user_id"]).execute().data or [{}]
                q_info = supabase.table("questions").select("*").eq("id", req["question_id"]).execute().data or [{}]
                e_info = supabase.table("exams").select("title").eq("id", req.get("exam_id")).execute().data or [{}]
                with st.container(border=True):
                    st.markdown(f"**{u_info[0].get('name','Unknown')}** ({u_info[0].get('email','')})")
                    st.caption(f"{e_info[0].get('title','Unknown Exam')} | Status: {req.get('status','pending')}")
                    st.markdown(f"Question: **{q_info[0].get('question','')}**")
                    admin_hint = st.text_area("Admin custom hint", value=req.get("admin_hint") or q_info[0].get("hint", "") or "", key=f"admin_hint_{req['id']}", height=90)
                    approve_col, reject_col = st.columns(2)
                    with approve_col:
                        if st.button("Approve Hint", key=f"approve_hint_{req['id']}", type="primary", use_container_width=True):
                            supabase.table("hint_requests").update({
                                "status": "approved",
                                "admin_hint": admin_hint.strip(),
                                "reviewed_at": "now()",
                            }).eq("id", req["id"]).execute()
                            st.success("Hint approved.")
                            st.rerun()
                    with reject_col:
                        if st.button("Reject", key=f"reject_hint_{req['id']}", use_container_width=True):
                            supabase.table("hint_requests").update({"status": "rejected", "reviewed_at": "now()"}).eq("id", req["id"]).execute()
                            st.warning("Hint rejected.")
                            st.rerun()

        with r_tab7:
            st.title("Test Case Reports")
            try:
                reports = supabase.table("test_case_reports").select("*").order("created_at", desc=True).execute().data or []
            except Exception as e:
                reports = []
                st.warning(f"Test case reports table ready ledu: {e}")
            if not reports:
                st.info("Reports levu.")
            for rep in reports:
                u_info = supabase.table("users").select("name, email").eq("id", rep["user_id"]).execute().data or [{}]
                q_rows = supabase.table("questions").select("*").eq("id", rep["question_id"]).execute().data or []
                if not q_rows:
                    continue
                q = q_rows[0]
                meta = get_programming_meta(q)
                case_no = int(rep.get("case_no") or 0)
                tc = (meta.get("test_cases", []) or [])[case_no - 1] if 0 < case_no <= len(meta.get("test_cases", []) or []) else {}
                with st.container(border=True):
                    st.markdown(f"**{u_info[0].get('name','Unknown')}** ({u_info[0].get('email','')})")
                    st.caption(f"Q: {q.get('question','')} | Test Case {case_no} | Status: {rep.get('status','pending')}")
                    if rep.get("report_text"):
                        st.info(rep.get("report_text"))
                    payload = rep.get("result_payload") or {}
                    if payload:
                        c_exp, c_act = st.columns(2)
                        with c_exp:
                            st.caption("Expected")
                            st.code(payload.get("expected_output", ""), language="text")
                        with c_act:
                            st.caption("Actual")
                            st.code(payload.get("actual_output", ""), language="text")
                    if st.button("Edit Question/Test Case", key=f"open_report_edit_{rep['id']}", use_container_width=True):
                        st.session_state[f"report_edit_{rep['id']}"] = not st.session_state.get(f"report_edit_{rep['id']}", False)
                        st.rerun()
                    if st.session_state.get(f"report_edit_{rep['id']}", False):
                        with st.form(f"report_edit_form_{rep['id']}"):
                            new_q_title = st.text_area("Question / Title", value=q.get("question", ""), height=80)
                            new_desc = st.text_area("Description", value=meta.get("description", ""), height=140)
                            new_input = st.text_area("Test Input", value=tc.get("input", ""), height=80)
                            new_output = st.text_area("Expected Output", value=tc.get("expected_output", ""), height=80)
                            new_marks = st.number_input("Marks", min_value=1, max_value=100, value=int(tc.get("marks", 1) or 1), step=1)
                            new_hidden = st.checkbox("Hidden", value=bool(tc.get("hidden", False)))
                            note = st.text_input("Admin note", value=rep.get("admin_note") or "")
                            if st.form_submit_button("Save Fix & Close Report", type="primary", use_container_width=True):
                                cases = meta.get("test_cases", []) or []
                                if 0 < case_no <= len(cases):
                                    cases[case_no - 1] = {
                                        "input": new_input,
                                        "expected_output": new_output,
                                        "marks": int(new_marks),
                                        "hidden": bool(new_hidden),
                                    }
                                new_explanation = make_programming_meta(new_desc, cases, meta.get("language", "java"))
                                supabase.table("questions").update({"question": new_q_title, "explanation": new_explanation}).eq("id", q["id"]).execute()
                                supabase.table("test_case_reports").update({"status": "fixed", "admin_note": note, "reviewed_at": "now()"}).eq("id", rep["id"]).execute()
                                st.success("Question/test case updated.")
                                st.rerun()

        with r_tab8:
            st.title("Student Attendance")
            all_students = supabase.table("users").select("id, name, email").eq("role","user").execute().data
            if not all_students:
                st.info("Students .")
            else:
                sel_student = st.selectbox("Student select ", [f"{s['name']} ({s['email']})" for s in all_students])
                sel_idx = [f"{s['name']} ({s['email']})" for s in all_students].index(sel_student)
                sel_uid = all_students[sel_idx]["id"]
                show_attendance_tab(sel_uid)

        with r_tab9:
            st.title("Live Programming Exams")
            with st.expander("Programming session database SQL setup"):
                st.code(get_programming_session_sql(), language="sql")
            try:
                sessions = supabase.table("programming_exam_sessions").select("*").eq("status", "active").order("updated_at", desc=True).execute().data or []
            except Exception as e:
                sessions = []
                st.warning(f"Session table ready ledu: {e}")
            if not sessions:
                st.info("Active programming exams levu.")
            for sess in sessions:
                u_data = supabase.table("users").select("name, email").eq("id", sess["user_id"]).execute().data or [{}]
                e_data = supabase.table("exams").select("title").eq("id", sess["exam_id"]).execute().data or [{}]
                uname = u_data[0].get("name") or u_data[0].get("email") or "Student"
                ename = e_data[0].get("title") or "Programming Exam"
                with st.container(border=True):
                    c1, c2, c3 = st.columns([4, 1.5, 1.5])
                    with c1:
                        st.markdown(f"**{uname}** - {ename}")
                        st.caption(f"Question: {int(sess.get('question_index') or 0) + 1} | Malpractice: {sess.get('malpractice_count') or 0} | Last: {sess.get('last_malpractice_reason') or 'None'}")
                    with c2:
                        if sess.get("force_submit"):
                            st.warning("Force submit pending")
                        else:
                            if st.button("Submit Exam", key=f"force_submit_{sess['user_id']}_{sess['exam_id']}", type="primary", use_container_width=True):
                                supabase.table("programming_exam_sessions").update({"force_submit": True, "updated_at": "now()"}).eq("user_id", sess["user_id"]).eq("exam_id", sess["exam_id"]).execute()
                                send_notification(f"Admin requested submit for {ename}. Mee exam automatic ga submit avuthundi.", sess["user_id"])
                                st.success("Student exam auto submit ki mark ayyindi.")
                                st.rerun()
                    with c3:
                        if st.button("Refresh", key=f"refresh_sess_{sess['user_id']}_{sess['exam_id']}", use_container_width=True):
                            st.rerun()


    elif menu == "Group Chat":
        with st.expander(" Broadcast Notification "):
            notif_msg = st.text_input("Message", key="broadcast_msg")
            if st.button("Send to All Users", type="primary"):
                if notif_msg.strip():
                    send_notification(notif_msg.strip())
                    st.success(" !"); st.rerun()
                else:
                    st.warning("Message enter .")
            st.divider()
            group_chat()

# =========================
# CREDIT CARD SHARING PORTAL
# =========================
CARD_CONFIG = {
    "card_1": {"label": "Card 1", "bill_day": 13},
    "card_2": {"label": "Card 2", "bill_day": 22},
}


def get_card_label(card_code):
    return CARD_CONFIG.get(card_code, {}).get("label", card_code)


def get_card_bill_day(card_code):
    return int(CARD_CONFIG.get(card_code, {}).get("bill_day", 1))


def add_months(source_date, months):
    month = source_date.month - 1 + months
    year = source_date.year + month // 12
    month = month % 12 + 1
    max_days = [31, 29 if year % 4 == 0 and (year % 100 != 0 or year % 400 == 0) else 28, 31, 30, 31, 30, 31, 31, 30, 31, 30, 31]
    return source_date.replace(year=year, month=month, day=min(source_date.day, max_days[month - 1]))


def get_statement_period(card_code, anchor=None):
    anchor = anchor or date.today()
    bill_day = get_card_bill_day(card_code)
    this_month_bill = date(anchor.year, anchor.month, min(bill_day, 28 if anchor.month == 2 else 30 if anchor.month in [4, 6, 9, 11] else 31))
    if anchor >= this_month_bill:
        end_date = this_month_bill
    else:
        end_date = add_months(this_month_bill, -1)
    start_date = add_months(end_date, -1) + timedelta(days=1)
    return start_date, end_date


def money_value(value):
    try:
        return float(value or 0)
    except Exception:
        return 0.0


def show_credit_card_sql_help():
    with st.expander("Credit Cards database SQL setup"):
        st.code("""
create table if not exists card_users (
  id uuid primary key default gen_random_uuid(),
  name text not null,
  mobile text,
  app_pin text not null unique,
  active boolean not null default true,
  created_at timestamptz default now()
);

create table if not exists card_user_cards (
  id uuid primary key default gen_random_uuid(),
  user_id uuid not null references card_users(id) on delete cascade,
  card_code text not null check (card_code in ('card_1','card_2')),
  created_at timestamptz default now(),
  unique(user_id, card_code)
);

create table if not exists card_transactions (
  id uuid primary key default gen_random_uuid(),
  user_id uuid not null references card_users(id) on delete cascade,
  card_code text not null check (card_code in ('card_1','card_2')),
  purpose text not null,
  amount numeric not null default 0,
  transaction_date date not null,
  proof_url text,
  source_key text unique,
  status text not null default 'pending' check (status in ('pending','approved','rejected')),
  admin_note text,
  approved_by uuid references users(id),
  approved_at timestamptz,
  created_at timestamptz default now()
);

create table if not exists card_payments (
  id uuid primary key default gen_random_uuid(),
  user_id uuid not null references card_users(id) on delete cascade,
  card_code text not null check (card_code in ('card_1','card_2')),
  billing_month text not null,
  bill_start date not null,
  bill_end date not null,
  amount numeric not null default 0,
  proof_url text not null,
  status text not null default 'pending' check (status in ('pending','paid','rejected')),
  admin_note text,
  approved_by uuid references users(id),
  approved_at timestamptz,
  created_at timestamptz default now(),
  unique(user_id, card_code, billing_month)
);

create table if not exists card_recurring_transactions (
  id uuid primary key default gen_random_uuid(),
  user_id uuid not null references card_users(id) on delete cascade,
  card_code text not null check (card_code in ('card_1','card_2')),
  purpose text not null,
  amount numeric not null default 0,
  start_date date not null default current_date,
  end_date date,
  active boolean not null default true,
  created_at timestamptz default now()
);

alter table card_transactions add column if not exists source_key text;
create unique index if not exists card_transactions_source_key_idx
  on card_transactions(source_key)
  where source_key is not null;
""", language="sql")


def get_card_users():
    try:
        return supabase.table("card_users").select("id, name, mobile, app_pin, active").order("name").execute().data or []
    except Exception as e:
        st.error(f"Card users load failed: {e}")
        return []


def get_assigned_cards(user_id):
    try:
        rows = supabase.table("card_user_cards").select("card_code").eq("user_id", user_id).execute().data or []
        return [r.get("card_code") for r in rows]
    except Exception:
        return []


def upload_optional_image(uploaded_file):
    if not uploaded_file:
        return ""
    return upload_image_to_imgbb(uploaded_file) or ""


def generate_recurring_transactions(user_id=None):
    try:
        query = supabase.table("card_recurring_transactions").select("*").eq("active", True)
        if user_id:
            query = query.eq("user_id", user_id)
        recurring_rows = query.execute().data or []
    except Exception as e:
        st.error(f"Monthly auto payments load failed: {e}")
        return 0

    created_count = 0
    today = date.today()
    for item in recurring_rows:
        card_code = item.get("card_code")
        if not card_code:
            continue
        start_date, end_date = get_statement_period(card_code, today)
        item_start = date.fromisoformat(str(item.get("start_date") or today))
        item_end = date.fromisoformat(str(item["end_date"])) if item.get("end_date") else None
        if item_start > end_date:
            continue
        if item_end and item_end < start_date:
            continue
        billing_month = end_date.strftime("%Y-%m")
        source_key = f"recurring:{item['id']}:{billing_month}"
        txn_date = max(item_start, start_date)
        if item_end:
            txn_date = min(txn_date, item_end)
        try:
            existing = supabase.table("card_transactions").select("id").eq("source_key", source_key).execute().data
            if existing:
                continue
            supabase.table("card_transactions").insert({
                "user_id": item["user_id"],
                "card_code": card_code,
                "purpose": item.get("purpose") or "Monthly payment",
                "amount": money_value(item.get("amount")),
                "transaction_date": str(txn_date),
                "proof_url": "",
                "source_key": source_key,
                "status": "approved",
                "admin_note": "Auto monthly payment",
                "approved_by": st.session_state.user_id if st.session_state.get("role") == "admin" else None,
                "approved_at": "now()",
            }).execute()
            created_count += 1
        except Exception as e:
            st.error(f"Monthly auto payment create failed: {e}")
    return created_count


def admin_credit_cards_dashboard():
    st.title("Credit Cards Admin")
    show_credit_card_sql_help()
    generate_recurring_transactions()
    tab_users, tab_txns, tab_payments, tab_recurring, tab_manual = st.tabs(["Users & Cards", "Approve Transactions", "Payment Screenshots", "Monthly Auto Payments", "Add/Edit Transactions"])

    with tab_users:
        st.subheader("Add card user PIN login")
        with st.form("add_card_user_form", clear_on_submit=True):
            name = st.text_input("User name")
            mobile = st.text_input("Mobile / note")
            pin = st.text_input("Login PIN", type="password", max_chars=6)
            assigned = st.multiselect("Cards allowed", ["card_1", "card_2"], format_func=get_card_label, default=["card_1"])
            submitted = st.form_submit_button("Create PIN Login", type="primary")
        created_user = None
        if submitted:
            if not name.strip() or not pin.strip() or not assigned:
                st.error("Name, PIN, cards required.")
            elif not pin.isdigit() or len(pin) < 4:
                st.error("PIN 4 to 6 digits undali.")
            else:
                try:
                    existing_lms = supabase.table("users").select("id").eq("app_pin", pin).execute().data
                    existing_card = supabase.table("card_users").select("id").eq("app_pin", pin).execute().data
                    if existing_lms or existing_card:
                        st.error("Ee PIN already another user ki undi. Vere PIN try cheyyandi.")
                    else:
                        created = supabase.table("card_users").insert({
                            "name": name.strip(),
                            "mobile": mobile.strip(),
                            "app_pin": pin,
                        }).execute()
                        created_data = created.data or []
                        if not created_data:
                            created_data = supabase.table("card_users").select("*").eq("app_pin", pin).execute().data or []
                        new_user_id = created_data[0]["id"] if created_data else None
                        if not new_user_id:
                            raise Exception("Card user create failed.")
                        for card_code in assigned:
                            supabase.table("card_user_cards").insert({"user_id": new_user_id, "card_code": card_code}).execute()
                        created_user = created_data[0] if created_data else {"id": new_user_id, "name": name.strip(), "mobile": mobile.strip(), "app_pin": pin, "active": True}
                        st.success(f"Card user created: {name.strip()} | PIN: {pin}")
                except Exception as e:
                    st.error(f"Create user failed: {e}")

        st.divider()
        users = get_card_users()
        if created_user and not any(u.get("id") == created_user.get("id") for u in users):
            users = [created_user] + users
        if not users:
            st.info("No card users yet.")
        for user in users:
            with st.expander(f"{user.get('name','User')} - {user.get('mobile','')}"):
                current_cards = get_assigned_cards(user["id"])
                selected_cards = st.multiselect("Assigned cards", ["card_1", "card_2"], default=current_cards, format_func=get_card_label, key=f"cards_{user['id']}")
                if st.button("Save cards", key=f"save_cards_{user['id']}"):
                    try:
                        supabase.table("card_user_cards").delete().eq("user_id", user["id"]).execute()
                        for card_code in selected_cards:
                            supabase.table("card_user_cards").insert({"user_id": user["id"], "card_code": card_code}).execute()
                        st.success("Cards updated.")
                        st.rerun()
                    except Exception as e:
                        st.error(f"Cards update failed: {e}")

    with tab_txns:
        st.subheader("Pending user transactions")
        try:
            rows = supabase.table("card_transactions").select("*").eq("status", "pending").order("transaction_date", desc=True).execute().data or []
        except Exception as e:
            st.error(f"Transactions load failed: {e}")
            rows = []
        if not rows:
            st.info("Pending transactions levu.")
        for row in rows:
            user = supabase.table("card_users").select("name, mobile").eq("id", row["user_id"]).execute().data or [{}]
            with st.expander(f"{get_card_label(row['card_code'])} | {user[0].get('name','User')} | Rs.{money_value(row.get('amount')):.2f} | {row.get('transaction_date')}"):
                st.write(row.get("purpose", ""))
                if row.get("proof_url"):
                    st.link_button("Open screenshot", row["proof_url"])
                    st.image(row["proof_url"], width=300)
                note = st.text_input("Admin note", value=row.get("admin_note") or "", key=f"txn_note_{row['id']}")
                col1, col2 = st.columns(2)
                with col1:
                    if st.button("Approve", type="primary", key=f"txn_ok_{row['id']}"):
                        supabase.table("card_transactions").update({"status": "approved", "admin_note": note, "approved_by": st.session_state.user_id, "approved_at": "now()"}).eq("id", row["id"]).execute()
                        st.success("Approved.")
                        st.rerun()
                with col2:
                    if st.button("Reject", key=f"txn_no_{row['id']}"):
                        supabase.table("card_transactions").update({"status": "rejected", "admin_note": note}).eq("id", row["id"]).execute()
                        st.warning("Rejected.")
                        st.rerun()

    with tab_payments:
        st.subheader("User bill payment screenshots")
        st.markdown("### All users pending bills")
        try:
            all_card_users = get_card_users()
            pending_bill_rows = []
            for card_user in all_card_users:
                user_cards = get_assigned_cards(card_user["id"])
                for card_code in user_cards:
                    start_date, end_date = get_statement_period(card_code)
                    billing_month = end_date.strftime("%Y-%m")
                    approved = supabase.table("card_transactions").select("*").eq("user_id", card_user["id"]).eq("card_code", card_code).eq("status", "approved").gte("transaction_date", str(start_date)).lte("transaction_date", str(end_date)).execute().data or []
                    total = sum(money_value(r.get("amount")) for r in approved)
                    if total <= 0:
                        continue
                    payments_for_bill = supabase.table("card_payments").select("*").eq("user_id", card_user["id"]).eq("card_code", card_code).eq("billing_month", billing_month).order("created_at", desc=True).execute().data or []
                    paid = any(p.get("status") == "paid" for p in payments_for_bill)
                    if paid:
                        continue
                    pending_pay = any(p.get("status") == "pending" for p in payments_for_bill)
                    pending_bill_rows.append({
                        "user": card_user.get("name", "User"),
                        "mobile": card_user.get("mobile", ""),
                        "card": get_card_label(card_code),
                        "month": billing_month,
                        "period": f"{start_date} to {end_date}",
                        "amount": round(total, 2),
                        "status": "Payment uploaded - approval pending" if pending_pay else "Not paid",
                    })
            if pending_bill_rows:
                c1, c2 = st.columns(2)
                c1.metric("Pending bills", len(pending_bill_rows))
                c2.metric("Total receivable", f"Rs.{sum(r['amount'] for r in pending_bill_rows):.2f}")
                st.dataframe(pending_bill_rows, use_container_width=True, hide_index=True)
            else:
                st.info("All users pending bills levu.")
        except Exception as e:
            st.error(f"Pending bills load failed: {e}")
        st.divider()
        try:
            payments = supabase.table("card_payments").select("*").order("created_at", desc=True).execute().data or []
        except Exception as e:
            st.error(f"Payments load failed: {e}")
            payments = []
        for pay in payments:
            user = supabase.table("card_users").select("name, mobile").eq("id", pay["user_id"]).execute().data or [{}]
            with st.expander(f"{pay.get('status','pending').upper()} | {get_card_label(pay['card_code'])} | {user[0].get('name','User')} | {pay.get('billing_month')} | Rs.{money_value(pay.get('amount')):.2f}"):
                st.write(f"Bill period: {pay.get('bill_start')} to {pay.get('bill_end')}")
                st.link_button("Open payment screenshot", pay["proof_url"])
                st.image(pay["proof_url"], width=320)
                note = st.text_input("Payment admin note", value=pay.get("admin_note") or "", key=f"pay_note_{pay['id']}")
                col1, col2 = st.columns(2)
                with col1:
                    if st.button("Mark Paid", type="primary", key=f"pay_ok_{pay['id']}"):
                        supabase.table("card_payments").update({"status": "paid", "admin_note": note, "approved_by": st.session_state.user_id, "approved_at": "now()"}).eq("id", pay["id"]).execute()
                        st.success("Payment marked paid.")
                        st.rerun()
                with col2:
                    if st.button("Reject Payment", key=f"pay_no_{pay['id']}"):
                        supabase.table("card_payments").update({"status": "rejected", "admin_note": note}).eq("id", pay["id"]).execute()
                        st.warning("Payment rejected.")
                        st.rerun()

    with tab_recurring:
        st.subheader("One-time setup for monthly payments")
        users = get_card_users()
        if not users:
            st.info("Create card users first.")
        else:
            labels = {f"{u.get('name','User')} ({u.get('mobile','')})": u for u in users}
            selected_label = st.selectbox("User", list(labels.keys()), key="rec_user")
            selected_user = labels[selected_label]
            user_cards = get_assigned_cards(selected_user["id"]) or ["card_1", "card_2"]
            with st.form("recurring_payment_form", clear_on_submit=True):
                rec_card = st.selectbox("Card", user_cards, format_func=get_card_label, key="rec_card")
                rec_purpose = st.text_input("Purpose", placeholder="Example: EMI, Netflix, Rent")
                rec_amount = st.number_input("Monthly amount", min_value=0.0, step=1.0, key="rec_amount")
                rec_start = st.date_input("Start from", value=date.today(), key="rec_start")
                has_end = st.checkbox("End date unda?")
                rec_end = st.date_input("End date", value=add_months(date.today(), 12), key="rec_end") if has_end else None
                rec_submit = st.form_submit_button("Add Monthly Auto Payment", type="primary")
            if rec_submit:
                if not rec_purpose.strip() or rec_amount <= 0:
                    st.error("Purpose and amount required.")
                elif rec_end and rec_end < rec_start:
                    st.error("End date start date kante mundu undakudadhu.")
                else:
                    try:
                        supabase.table("card_recurring_transactions").insert({
                            "user_id": selected_user["id"],
                            "card_code": rec_card,
                            "purpose": rec_purpose.strip(),
                            "amount": rec_amount,
                            "start_date": str(rec_start),
                            "end_date": str(rec_end) if rec_end else None,
                            "active": True,
                        }).execute()
                        made = generate_recurring_transactions(selected_user["id"])
                        st.success(f"Monthly auto payment added. Current month entries created: {made}")
                        st.rerun()
                    except Exception as e:
                        st.error(f"Monthly auto payment add failed: {e}")

        st.divider()
        st.subheader("Manage monthly auto payments")
        try:
            recurring_rows = supabase.table("card_recurring_transactions").select("*").order("created_at", desc=True).execute().data or []
        except Exception as e:
            st.error(f"Monthly auto payments load failed: {e}")
            recurring_rows = []
        if not recurring_rows:
            st.info("Monthly auto payments levu.")
        for rec in recurring_rows:
            user = supabase.table("card_users").select("name, mobile").eq("id", rec["user_id"]).execute().data or [{}]
            status = "Active" if rec.get("active") else "Paused"
            with st.expander(f"{status} | {user[0].get('name','User')} | {get_card_label(rec.get('card_code'))} | {rec.get('purpose')} | Rs.{money_value(rec.get('amount')):.2f}"):
                c1, c2, c3 = st.columns([2, 1, 1])
                new_purpose = c1.text_input("Purpose", value=rec.get("purpose") or "", key=f"rec_purpose_{rec['id']}")
                new_amount = c2.number_input("Amount", min_value=0.0, value=money_value(rec.get("amount")), step=1.0, key=f"rec_amount_{rec['id']}")
                new_active = c3.checkbox("Active", value=bool(rec.get("active")), key=f"rec_active_{rec['id']}")
                c4, c5 = st.columns(2)
                new_start = c4.date_input("Start date", value=date.fromisoformat(str(rec.get("start_date"))), key=f"rec_start_{rec['id']}")
                default_end = date.fromisoformat(str(rec["end_date"])) if rec.get("end_date") else add_months(date.today(), 12)
                use_end = c5.checkbox("Use end date", value=bool(rec.get("end_date")), key=f"rec_use_end_{rec['id']}")
                new_end = c5.date_input("End date", value=default_end, key=f"rec_end_{rec['id']}") if use_end else None
                col_save, col_delete = st.columns(2)
                with col_save:
                    if st.button("Save auto payment", key=f"save_rec_{rec['id']}", type="primary"):
                        supabase.table("card_recurring_transactions").update({
                            "purpose": new_purpose.strip(),
                            "amount": new_amount,
                            "start_date": str(new_start),
                            "end_date": str(new_end) if new_end else None,
                            "active": new_active,
                        }).eq("id", rec["id"]).execute()
                        st.success("Monthly auto payment updated.")
                        st.rerun()
                with col_delete:
                    if st.button("Delete auto payment", key=f"del_rec_{rec['id']}"):
                        supabase.table("card_recurring_transactions").delete().eq("id", rec["id"]).execute()
                        st.warning("Monthly auto payment deleted.")
                        st.rerun()

    with tab_manual:
        st.subheader("Add transaction for a user")
        users = get_card_users()
        if users:
            labels = {f"{u.get('name','User')} ({u.get('mobile','')})": u for u in users}
            selected_label = st.selectbox("User", list(labels.keys()), key="manual_card_user")
            selected_user = labels[selected_label]
            user_cards = get_assigned_cards(selected_user["id"]) or ["card_1", "card_2"]
            with st.form("manual_card_txn_form", clear_on_submit=True):
                card_code = st.selectbox("Card", user_cards, format_func=get_card_label)
                purpose = st.text_input("Purpose")
                amount = st.number_input("Amount", min_value=0.0, step=1.0)
                txn_date = st.date_input("Date", value=date.today())
                proof_file = st.file_uploader("Screenshot optional", type=["png", "jpg", "jpeg", "webp", "pdf"])
                submit_txn = st.form_submit_button("Add Approved Transaction", type="primary")
            if submit_txn:
                if not purpose.strip() or amount <= 0:
                    st.error("Purpose and amount required.")
                else:
                    proof_url = upload_optional_image(proof_file)
                    supabase.table("card_transactions").insert({
                        "user_id": selected_user["id"],
                        "card_code": card_code,
                        "purpose": purpose.strip(),
                        "amount": amount,
                        "transaction_date": str(txn_date),
                        "proof_url": proof_url,
                        "status": "approved",
                        "approved_by": st.session_state.user_id,
                        "approved_at": "now()",
                    }).execute()
                    st.success("Transaction added.")
                    st.rerun()

        st.divider()
        st.subheader("Edit existing transactions")
        try:
            txns = supabase.table("card_transactions").select("*").order("transaction_date", desc=True).limit(100).execute().data or []
        except Exception as e:
            st.error(f"Transactions load failed: {e}")
            txns = []
        for row in txns:
            user = supabase.table("card_users").select("name, mobile").eq("id", row["user_id"]).execute().data or [{}]
            with st.expander(f"{get_card_label(row['card_code'])} | {user[0].get('name','User')} | {row.get('purpose','')}"):
                col1, col2, col3 = st.columns([2, 1, 1])
                new_purpose = col1.text_input("Purpose", value=row.get("purpose") or "", key=f"edit_purpose_{row['id']}")
                new_amount = col2.number_input("Amount", min_value=0.0, value=money_value(row.get("amount")), step=1.0, key=f"edit_amount_{row['id']}")
                new_status = col3.selectbox("Status", ["pending", "approved", "rejected"], index=["pending", "approved", "rejected"].index(row.get("status", "pending")), key=f"edit_status_{row['id']}")
                new_date = st.date_input("Date", value=date.fromisoformat(str(row.get("transaction_date"))), key=f"edit_date_{row['id']}")
                if row.get("proof_url"):
                    st.link_button("Open screenshot", row["proof_url"], key=f"edit_link_{row['id']}")
                if st.button("Save transaction", key=f"save_txn_{row['id']}"):
                    supabase.table("card_transactions").update({
                        "purpose": new_purpose.strip(),
                        "amount": new_amount,
                        "transaction_date": str(new_date),
                        "status": new_status,
                    }).eq("id", row["id"]).execute()
                    st.success("Transaction updated.")
                    st.rerun()


def card_user_dashboard():
    st.sidebar.title("Credit Card Portal")
    persist_browser_login()
    if st.sidebar.button("Logout", use_container_width=True):
        for key in defaults:
            st.session_state[key] = defaults[key]
        show_logout_redirect()
    st.sidebar.divider()
    pages = ["Monthly Bill", "Pending Bills", "Paid", "Add Transaction", "History"]
    for pg in pages:
        if st.sidebar.button(pg, use_container_width=True, type="primary" if st.session_state.card_user_page == pg else "secondary", key=f"card_nav_{pg}"):
            st.session_state.card_user_page = pg
            st.rerun()

    user_id = st.session_state.user_id
    assigned_cards = get_assigned_cards(user_id)
    if not assigned_cards:
        st.warning("No cards assigned yet. Admin ni contact cheyyandi.")
        return
    generate_recurring_transactions(user_id)

    st.title("Credit Card Portal")
    if st.session_state.card_user_page == "Add Transaction":
        with st.form("user_add_card_txn", clear_on_submit=True):
            card_code = st.selectbox("Card", assigned_cards, format_func=get_card_label)
            purpose = st.text_input("Purpose")
            amount = st.number_input("Amount", min_value=0.0, step=1.0)
            txn_date = st.date_input("Transaction date", value=date.today())
            proof_file = st.file_uploader("Payment/transaction screenshot", type=["png", "jpg", "jpeg", "webp"])
            submitted = st.form_submit_button("Submit for admin approval", type="primary")
        if submitted:
            if not purpose.strip() or amount <= 0:
                st.error("Purpose and amount required.")
            else:
                proof_url = upload_optional_image(proof_file)
                supabase.table("card_transactions").insert({
                    "user_id": user_id,
                    "card_code": card_code,
                    "purpose": purpose.strip(),
                    "amount": amount,
                    "transaction_date": str(txn_date),
                    "proof_url": proof_url,
                    "status": "pending",
                }).execute()
                st.success("Transaction sent to admin for approval.")
                st.rerun()
        return

    if st.session_state.card_user_page == "Pending Bills":
        st.subheader("Pending Bills")
        found_pending_bill = False
        for card_code in assigned_cards:
            start_date, end_date = get_statement_period(card_code)
            billing_month = end_date.strftime("%Y-%m")
            approved = supabase.table("card_transactions").select("*").eq("user_id", user_id).eq("card_code", card_code).eq("status", "approved").gte("transaction_date", str(start_date)).lte("transaction_date", str(end_date)).order("transaction_date", desc=False).execute().data or []
            total = sum(money_value(r.get("amount")) for r in approved)
            payments = supabase.table("card_payments").select("*").eq("user_id", user_id).eq("card_code", card_code).eq("billing_month", billing_month).order("created_at", desc=True).execute().data or []
            paid = any(p.get("status") == "paid" for p in payments)
            pending_pay = any(p.get("status") == "pending" for p in payments)
            if total <= 0 or paid:
                continue
            found_pending_bill = True
            status_text = "Payment Waiting Admin Approval" if pending_pay else "PAY"
            with st.container(border=True):
                st.subheader(f"{get_card_label(card_code)} - {billing_month}")
                c1, c2, c3 = st.columns(3)
                c1.metric("Bill period", f"{start_date} to {end_date}")
                c2.metric("Bill amount", f"Rs.{total:.2f}")
                c3.metric("Status", status_text)
                st.dataframe([{
                    "date": r.get("transaction_date"),
                    "purpose": r.get("purpose"),
                    "amount": money_value(r.get("amount")),
                } for r in approved], use_container_width=True, hide_index=True)
                if not pending_pay:
                    with st.form(f"pending_pay_form_{card_code}_{billing_month}"):
                        pay_file = st.file_uploader("Upload bill paid screenshot", type=["png", "jpg", "jpeg", "webp"], key=f"pending_pay_file_{card_code}_{billing_month}")
                        submit_pay = st.form_submit_button("Pay / Submit Screenshot", type="primary")
                    if submit_pay:
                        proof_url = upload_optional_image(pay_file)
                        if not proof_url:
                            st.error("Screenshot upload required.")
                        else:
                            supabase.table("card_payments").upsert({
                                "user_id": user_id,
                                "card_code": card_code,
                                "billing_month": billing_month,
                                "bill_start": str(start_date),
                                "bill_end": str(end_date),
                                "amount": total,
                                "proof_url": proof_url,
                                "status": "pending",
                            }, on_conflict="user_id,card_code,billing_month").execute()
                            st.success("Payment screenshot admin approval ki sent.")
                            st.rerun()
        if not found_pending_bill:
            st.info("Pending bills levu.")
        return

    if st.session_state.card_user_page == "Paid":
        st.subheader("Paid Bills")
        paid_rows = supabase.table("card_payments").select("*").eq("user_id", user_id).eq("status", "paid").order("bill_end", desc=True).execute().data or []
        if not paid_rows:
            st.info("Paid bills levu.")
        else:
            st.dataframe([{
                "card": get_card_label(r.get("card_code")),
                "month": r.get("billing_month"),
                "period": f"{r.get('bill_start')} to {r.get('bill_end')}",
                "amount": money_value(r.get("amount")),
                "paid_at": str(r.get("approved_at") or "")[:19],
            } for r in paid_rows], use_container_width=True, hide_index=True)
            for row in paid_rows:
                if row.get("proof_url"):
                    st.link_button(f"Open screenshot - {get_card_label(row.get('card_code'))} {row.get('billing_month')}", row["proof_url"])
        return

    if st.session_state.card_user_page == "History":
        rows = supabase.table("card_transactions").select("*").eq("user_id", user_id).order("transaction_date", desc=True).execute().data or []
        if not rows:
            st.info("Transactions levu.")
        else:
            st.dataframe([{
                "date": r.get("transaction_date"),
                "card": get_card_label(r.get("card_code")),
                "purpose": r.get("purpose"),
                "amount": money_value(r.get("amount")),
                "status": r.get("status"),
                "admin_note": r.get("admin_note") or "",
            } for r in rows], use_container_width=True, hide_index=True)
        return

    for card_code in assigned_cards:
        start_date, end_date = get_statement_period(card_code)
        billing_month = end_date.strftime("%Y-%m")
        approved = supabase.table("card_transactions").select("*").eq("user_id", user_id).eq("card_code", card_code).eq("status", "approved").gte("transaction_date", str(start_date)).lte("transaction_date", str(end_date)).order("transaction_date", desc=False).execute().data or []
        total = sum(money_value(r.get("amount")) for r in approved)
        payments = supabase.table("card_payments").select("*").eq("user_id", user_id).eq("card_code", card_code).eq("billing_month", billing_month).order("created_at", desc=True).execute().data or []
        paid = any(p.get("status") == "paid" for p in payments)
        pending_pay = any(p.get("status") == "pending" for p in payments)
        status_text = "PAID" if paid else "Payment Waiting Admin Approval" if pending_pay else "PAY"
        with st.container(border=True):
            st.subheader(f"{get_card_label(card_code)} - Bill date {get_card_bill_day(card_code)}")
            c1, c2, c3 = st.columns(3)
            c1.metric("Bill period", f"{start_date} to {end_date}")
            c2.metric("Approved amount", f"Rs.{total:.2f}")
            c3.metric("Status", status_text)
            if approved:
                st.dataframe([{
                    "date": r.get("transaction_date"),
                    "purpose": r.get("purpose"),
                    "amount": money_value(r.get("amount")),
                } for r in approved], use_container_width=True, hide_index=True)
            else:
                st.info("Ee bill period lo approved transactions levu.")
            if not paid and total > 0:
                with st.form(f"pay_form_{card_code}_{billing_month}"):
                    pay_file = st.file_uploader("Upload bill paid screenshot", type=["png", "jpg", "jpeg", "webp"], key=f"pay_file_{card_code}_{billing_month}")
                    submit_pay = st.form_submit_button("Pay / Submit Screenshot", type="primary")
                if submit_pay:
                    proof_url = upload_optional_image(pay_file)
                    if not proof_url:
                        st.error("Screenshot upload required.")
                    else:
                        supabase.table("card_payments").upsert({
                            "user_id": user_id,
                            "card_code": card_code,
                            "billing_month": billing_month,
                            "bill_start": str(start_date),
                            "bill_end": str(end_date),
                            "amount": total,
                            "proof_url": proof_url,
                            "status": "pending",
                        }, on_conflict="user_id,card_code,billing_month").execute()
                        st.success("Payment screenshot admin approval ki sent.")
                        st.rerun()

# =========================
# USER DASHBOARD
# =========================
def user_dashboard(preview_mode=False):
    if not preview_mode:
        st.sidebar.title("User Workspace")
        persist_browser_login()
        if st.sidebar.button("Logout", use_container_width=True):
            for key in defaults:
                st.session_state[key] = defaults[key]
            show_logout_redirect()
        st.sidebar.divider()
        pages = ["My Classes", "Exams", "Progress", "Code Practice", "Group Chat", "Attendance"]
        if user_has_suprabhatam_access(st.session_state.user_id):
            pages.append("Suprabhatam")
        for pg in pages:
            if pg == "Group Chat":
                unread = get_unread_count(st.session_state.user_id)
                label = f"Group Chat ({unread})" if unread > 0 else "Group Chat"
            else:
                label = pg
            if st.sidebar.button(label, use_container_width=True,
                    type="primary" if st.session_state.user_page == pg else "secondary",
                    key=f"nav_{pg}"):
                st.session_state.user_page = pg
                st.rerun()
        user_page = st.session_state.user_page
        # Mark today's attendance on every login
        mark_today_attendance(st.session_state.user_id)
    else:
        st.info("Student Preview Mode")
        user_page = "My Classes"

    if not preview_mode:
        show_notification_banner(st.session_state.user_id)

    if user_page == "Group Chat":
        group_chat(); return
    if user_page == "Progress":
        show_student_progress_tab(st.session_state.user_id); return
    if user_page in ["Java Practice", "Code Practice"]:
        show_code_practice_tab(); return
    if user_page == "Programming":
        st.session_state.user_page = "Exams"
        user_page = "Exams"
    if user_page == "Exams":
        show_student_exams_tab(st.session_state.user_id); return
    if user_page == "Attendance":
        show_attendance_tab(st.session_state.user_id); return
    if user_page == "Suprabhatam":
        if user_has_suprabhatam_access(st.session_state.user_id):
            render_suprabhatam_reader(); return
        st.error("Suprabhatam access ledu. Admin ni contact cheyyandi."); return

    # My Classes
    modules = supabase.table("modules").select("*").execute().data
    modules = modules or []
    if st.session_state.completed_ids is None:
        all_completions = supabase.table("class_completions").select("class_id").eq("user_id", st.session_state.user_id).execute().data
        st.session_state.completed_ids = {str(c["class_id"]) for c in all_completions}
    completed_ids = st.session_state.completed_ids
    focus_class_id = str(st.session_state.get("focus_class_id", ""))
    focus_exam_id = str(st.session_state.get("focus_exam_id", ""))
    inject_student_home_styles()
    try:
        all_classes_count = len(supabase.table("classes").select("id").execute().data or [])
    except Exception:
        all_classes_count = 0
    done_count = len(completed_ids)
    pending_count = max(all_classes_count - done_count, 0)
    st.markdown(
        f"""
        <div class="student-home-hero">
            <h2>Learning Workspace</h2>
            <p>Modules, classes, videos, and notes anni oka place lo clean ga track cheyyandi.</p>
        </div>
        <div class="home-stat-row">
            <div class="home-stat"><small>Modules</small><strong>{len(modules)}</strong></div>
            <div class="home-stat"><small>Classes Done</small><strong>{done_count}/{all_classes_count}</strong></div>
            <div class="home-stat"><small>Pending Classes</small><strong>{pending_count}</strong></div>
        </div>
        """,
        unsafe_allow_html=True,
    )

    for module in modules:
        module_submodules = supabase.table("submodules").select("id").eq("module_id", module["id"]).execute().data
        sub_ids = [s["id"] for s in module_submodules]
        module_has_focus = False
        module_total = 0; module_done = 0
        for sid in sub_ids:
            cls_list = supabase.table("classes").select("id").eq("submodule_id", sid).execute().data
            module_total += len(cls_list)
            module_done += sum(1 for c in cls_list if str(c["id"]) in completed_ids)
            if focus_class_id and any(str(c["id"]) == focus_class_id for c in cls_list):
                module_has_focus = True
        pct = int((module_done / module_total * 100)) if module_total > 0 else 0

        with st.expander(f"{module['title']}    {module_done}/{module_total} classes  ({pct}%)", expanded=module_has_focus):
            st.markdown(
                f"""
                <div class="module-strip">
                    <h3>{html.escape(str(module.get('title', 'Module')))}</h3>
                    <span>{module_done}/{module_total} classes complete</span>
                </div>
                """,
                unsafe_allow_html=True,
            )
            if module_total > 0:
                st.progress(pct / 100, text=f"Module Progress: {pct}%")
            submodules = supabase.table("submodules").select("*").eq("module_id", module["id"]).execute().data
            for sub in submodules:
                sub_classes = supabase.table("classes").select("id").eq("submodule_id", sub["id"]).execute().data
                sub_total = len(sub_classes)
                sub_done = sum(1 for c in sub_classes if str(c["id"]) in completed_ids)
                sub_pct = int((sub_done / sub_total * 100)) if sub_total > 0 else 0
                st.subheader(f"{sub['title']}   {sub_done}/{sub_total}")
                if sub_total > 0: st.progress(sub_pct / 100)
                classes = supabase.table("classes").select("*").eq("submodule_id", sub["id"]).execute().data
                for cls in classes:
                    is_done = str(cls.get("id")) in completed_ids
                    is_focused_class = focus_class_id and str(cls.get("id")) == focus_class_id
                    st.markdown(
                        f"""
                        <div class="class-card">
                            <div class="class-title">{html.escape(str(cls.get('title', 'Class')))}</div>
                            <span class="exam-chip">{'Completed' if is_done else 'In progress'}</span>
                        </div>
                        """,
                        unsafe_allow_html=True,
                    )
                    if is_focused_class:
                        st.info("Progress tab nundi open chesina class idi.")
                    col_link1, col_link2, col_link3 = st.columns(3)
                    with col_link1:
                        if cls.get("class_link"): st.link_button("Join Class", cls["class_link"], use_container_width=True)
                    with col_link2:
                        if cls.get("recorded_video"): st.link_button("Watch Video", cls["recorded_video"], use_container_width=True)
                    with col_link3:
                        if cls.get("notes_pdf"): st.link_button("Notes PDF", cls["notes_pdf"], use_container_width=True)

                    class_id = cls.get("id")
                    if class_id:
                        try: cid = int(class_id)
                        except (ValueError, TypeError): cid = str(class_id)
                        if str(class_id) in completed_ids:
                            st.success("     !")
                        else:
                            if st.button("Mark as Completed", key=f"btn_done_{cls['id']}"):
                                try:
                                    supabase.table("class_completions").insert({"user_id": str(st.session_state.user_id), "class_id": cid}).execute()
                                    st.session_state.completed_ids.add(str(class_id))
                                    st.success("   !")
                                except Exception as e:
                                    st.error(f"Insert Error: {e}")

# =========================
# EXAM WORKSPACE
# =========================
def exam_workspace_view():
    questions = st.session_state.current_questions
    total_questions = len(questions)
    is_prog_exam = is_programming_exam(st.session_state.exam_id) if st.session_state.get("exam_id") else False
    is_proctored = bool(st.session_state.get("exam_proctoring_enabled"))
    inject_programming_exam_shell(is_prog_exam and not st.session_state.exam_submitted)
    if is_prog_exam and not st.session_state.exam_submitted:
        enable_fullscreen_exam_lock()
    if (is_prog_exam or is_proctored) and not st.session_state.exam_submitted:
        qp = st.query_params
        if qp.get("malpractice") == "1":
            reason = qp.get("mal_reason", "Tab switched or exam window lost focus")
            report_programming_malpractice(reason)
            try:
                del st.query_params["malpractice"]
                del st.query_params["mal_reason"]
            except Exception:
                pass
    if is_prog_exam and not st.session_state.exam_submitted:
        session = get_active_programming_session(st.session_state.user_id, st.session_state.exam_id)
        if session and session.get("force_submit"):
            st.error("Admin force submit requested. Exam automatic ga submit avuthundi...")
            try:
                submit_exam_attempt(questions, include_time=True, require_programming_submitted=False)
                st.session_state.exam_submitted = True
                st.rerun()
            except Exception as e:
                st.error(f"Auto submit failed: {e}")
                return
        save_programming_exam_session(status="active")
    if total_questions == 0:
        st.warning("No questions in this exam.")
        if st.button("Go Back"): st.session_state.start_exam = False; st.rerun()
        return

    remaining_time = 0
    if not st.session_state.exam_submitted:
        sections = st.session_state.get("exam_sections") or []
        current_section = get_current_section()
        section_start = int(current_section.get("start", 0) or 0)
        section_end = int(current_section.get("end", total_questions) or total_questions)
        if st.session_state.question_index < section_start or st.session_state.question_index >= section_end:
            st.session_state.question_index = section_start
        overall_remaining = int(st.session_state.exam_end_time - time.time())
        section_remaining = int((st.session_state.get("exam_section_end_time") or st.session_state.exam_end_time) - time.time())
        remaining_time = min(overall_remaining, section_remaining) if len(sections) > 1 else overall_remaining
        if remaining_time <= 0:
            next_section_index = int(st.session_state.get("exam_section_index", 0) or 0) + 1
            if len(sections) > 1 and move_to_exam_section(next_section_index):
                st.warning("Section time over. Next section start ayyindi.")
                time.sleep(1)
                st.rerun()
            else:
                st.error("Time out. Submitting exam...")
                time.sleep(1)
                try:
                    submit_exam_attempt(questions, include_time=False)
                    st.session_state.exam_submitted = True; st.rerun()
                except Exception as e:
                    st.error(f"Submit failed: {e}")
                    return

    if st.session_state.exam_submitted:
        exam_data = supabase.table("exams").select("*").eq("id", st.session_state.exam_id).execute().data or [{}]
        current_exam = exam_data[0]
        current_exam.setdefault("id", st.session_state.exam_id)
        current_exam.setdefault("title", st.session_state.exam_title)
        db_attempt = supabase.table("exam_attempts").select("*").eq("user_id", st.session_state.user_id).eq("exam_id", st.session_state.exam_id).execute().data
        if db_attempt and st.session_state.get("last_attempt_id"):
            last_id = st.session_state.last_attempt_id
            db_attempt = sorted(db_attempt, key=lambda att: 0 if att.get("id") == last_id else 1)
        if db_attempt:
            render_exam_result_summary(current_exam, questions, db_attempt[0], show_return=True)
        else:
            st.warning("Attempt details dorakaledu.")
            if st.button("Return to Exams", type="primary"):
                st.session_state.start_exam = False
                st.session_state.exam_submitted = False
                st.rerun()
    else:
        current = st.session_state.question_index
        question = questions[current]
        sections = st.session_state.get("exam_sections") or []
        current_section = get_current_section()
        section_start = int(current_section.get("start", 0) or 0)
        section_end = int(current_section.get("end", total_questions) or total_questions)
        section_label = str(current_section.get("title") or st.session_state.exam_title)
        mins, secs = divmod(remaining_time, 60)
        if is_prog_exam:
            st.markdown(
                f"""
                <div class="exam-topbar">
                    <div class="exam-title">{clean_ui_text(st.session_state.exam_title)}</div>
                    <div class="exam-status">
                        <span class="lock-pill">Fullscreen locked</span>
                        <span>{clean_ui_text(section_label)}</span>
                        <span>Question {current + 1}/{total_questions}</span>
                        <span>{mins:02d}:{secs:02d}</span>
                    </div>
                </div>
                """,
                unsafe_allow_html=True,
            )
        else:
            st.title(st.session_state.exam_title)
        left, right = st.columns([4, 1])
        with right:
            if is_proctored:
                render_live_proctoring_panel(enabled=True)
                st.divider()
            st.components.v1.html(f"""
                <div id="timer" style="font-size:2rem;font-weight:600;text-align:center;padding:12px;border-radius:8px;
                    background:{'#fff3cd' if remaining_time<300 else '#e8f4fd'};
                    color:{'#856404' if remaining_time<300 else '#0c63e4'};
                    border:1px solid {'#ffc107' if remaining_time<300 else '#b6d4fe'};">
                     <span id="countdown">{mins:02d}:{secs:02d}</span></div>
                <script>
                    var total={remaining_time};
                    function tick(){{if(total<=0){{document.getElementById('countdown').innerText="00:00";return;}}
                    total--;var m=Math.floor(total/60).toString().padStart(2,'0');var s=(total%60).toString().padStart(2,'0');
                    document.getElementById('countdown').innerText=m+':'+s;
                    if(total<300){{var el=document.getElementById('timer');el.style.background='#fff3cd';el.style.color='#856404';}}}}
                    setInterval(tick,1000);
                    function flagMalpractice(reason){{
                        try {{
                            var url = new URL(window.parent.location.href);
                            url.searchParams.set('malpractice', '1');
                            url.searchParams.set('mal_reason', reason);
                            window.parent.location.href = url.toString();
                        }} catch(e) {{}}
                    }}
                    document.addEventListener('visibilitychange', function(){{
                        if (document.hidden) flagMalpractice('Tab switched or minimized');
                    }});
                </script>""", height=80)
            st.divider()
            if len(sections) > 1:
                st.subheader(section_label)
                st.caption(f"Section {int(st.session_state.get('exam_section_index', 0) or 0) + 1}/{len(sections)}")
            st.subheader("Questions")
            cols = st.columns(3)
            for i in range(section_start, section_end):
                with cols[i % 3]:
                    q_id = questions[i]["id"]
                    label = f" {i+1}" if i==current else (f" {i+1}" if q_id in st.session_state.answers and st.session_state.answers[q_id] else f" {i+1}")
                    if st.button(label, key=f"qnav_{i}", use_container_width=True):
                        st.session_state.question_index = i; st.rerun()
            if is_prog_exam:
                st.divider()
                if st.button("Submit Final Exam", type="primary", key="submit_exam_top_right", use_container_width=True):
                    qid_cur = question["id"]
                    if qid_cur in st.session_state.question_start_time:
                        elapsed = int(time.time() - st.session_state.question_start_time[qid_cur])
                        prev = st.session_state.question_time_log.get(qid_cur, 0)
                        st.session_state.question_time_log[qid_cur] = prev + elapsed
                        del st.session_state.question_start_time[qid_cur]
                    try:
                        save_programming_exam_session(status="active")
                        submit_exam_attempt(questions, include_time=True, require_programming_submitted=False)
                        st.session_state.question_time_log = {}
                        st.session_state.question_start_time = {}
                        st.session_state.program_run_results = {}
                        st.session_state.program_submissions = {}
                        st.session_state.exam_submitted = True; st.rerun()
                    except Exception as e:
                        st.error(f"Submit failed: {e}")
            if len(sections) > 1 and int(st.session_state.get("exam_section_index", 0) or 0) < len(sections) - 1:
                st.divider()
                if st.button("Next Section", key="next_section_manual", use_container_width=True):
                    move_to_exam_section(int(st.session_state.get("exam_section_index", 0) or 0) + 1)
                    save_programming_exam_session(status="active")
                    st.rerun()

        with left:
            hcol1, hcol2 = st.columns([4, 1])
            with hcol1: st.subheader(f"Question {current+1}/{total_questions}")
            with hcol2:
                qid = question["id"]
                already_spent = st.session_state.question_time_log.get(qid, 0)
                st.components.v1.html(f"""
                    <div style="background:#f0f4ff;border:1px solid #b6d4fe;border-radius:8px;padding:6px 10px;
                        text-align:center;font-family:monospace;font-size:1.1rem;font-weight:600;color:#0c63e4;margin-top:8px;">
                         <span id="qtimer">00:00</span></div>
                    <script>var elapsed={already_spent};var qtimer=document.getElementById('qtimer');
                    function qtick(){{elapsed++;var m=Math.floor(elapsed/60).toString().padStart(2,'0');var s=(elapsed%60).toString().padStart(2,'0');qtimer.innerText=m+':'+s;}}
                    setInterval(qtick,1000);</script>""", height=55)

            if question["type"] != "programming":
                st.write(question["question"])
            if question.get("image_url"): st.image(question["image_url"], width=350)
            if qid not in st.session_state.question_start_time:
                st.session_state.question_start_time[qid] = time.time()
            stored_ans = st.session_state.answers.get(question["id"],"")

            if question["type"] == "mcq":
                opts = [
                    ("A", question.get("option_a","")),
                    ("B", question.get("option_b","")),
                    ("C", question.get("option_c","")),
                    ("D", question.get("option_d","")),
                ]
                # Custom CSS for option buttons
                st.markdown("""
                <style>
                div[data-testid="stButton"] > button[kind="primary"] {
                    background-color: #1a73e8 !important;
                    color: white !important;
                    border: 2px solid #1a73e8 !important;
                    border-radius: 10px !important;
                    padding: 12px 18px !important;
                    font-size: 1rem !important;
                    text-align: left !important;
                    white-space: normal !important;
                    height: auto !important;
                }
                div[data-testid="stButton"] > button[kind="secondary"] {
                    background-color: #ffffff !important;
                    color: #2c3e50 !important;
                    border: 2px solid #d0d8e8 !important;
                    border-radius: 10px !important;
                    padding: 12px 18px !important;
                    font-size: 1rem !important;
                    text-align: left !important;
                    white-space: normal !important;
                    height: auto !important;
                }
                </style>""", unsafe_allow_html=True)
                st.markdown("")
                for lbl, otxt in opts:
                    if not otxt:
                        continue
                    is_selected = (stored_ans == lbl or stored_ans == otxt)
                    btn_label = f"{' ' if is_selected else ''}{lbl}. {otxt}"
                    if st.button(btn_label, key=f"opt_{question['id']}_{lbl}",
                                 use_container_width=True,
                                 type="primary" if is_selected else "secondary"):
                        st.session_state.answers[question["id"]] = lbl
                        save_programming_exam_session(status="active")
                        st.rerun()

            elif question["type"] == "blank":
                answer = st.text_input("Your Answer", value=stored_ans, key=f"text_{question['id']}")
                if answer != stored_ans:
                    st.session_state.answers[question["id"]] = answer
                    save_programming_exam_session(status="active")
            else:
                meta = get_programming_meta(question)
                max_marks = get_question_max_marks(question)
                stored_code, stored_language = parse_program_answer(stored_ans, meta.get("language", "java"))
                language_options = list(PROGRAMMING_LANGUAGE_LABELS.keys())
                current_language_label = get_programming_language_meta(stored_language)["label"]
                p_left, p_right = st.columns([1, 1])
                with p_left:
                    with st.container(height=660):
                        st.markdown("<div class='problem-pane'>", unsafe_allow_html=True)
                        st.markdown(f"### {question['question']}")
                        if meta.get("description"):
                            st.markdown(meta["description"])
                        st.caption(f"Marks: {max_marks}")
                        hint_req = get_user_question_hint(question["id"]) if st.session_state.get("role") != "admin" else None
                        if hint_req and hint_req.get("status") == "approved" and hint_req.get("admin_hint"):
                            st.info(f"Hint: {hint_req.get('admin_hint')}")
                        elif hint_req and hint_req.get("status") == "pending":
                            st.caption("Hint request pending with admin.")
                        elif st.session_state.get("role") != "admin":
                            if st.button("Request Hint", key=f"hint_request_{question['id']}", use_container_width=True):
                                ok, msg = submit_hint_request(st.session_state.exam_id, question["id"])
                                if ok:
                                    st.success(msg)
                                    st.rerun()
                                else:
                                    st.warning(msg)
                        visible_cases = [tc for tc in meta.get("test_cases", []) if not tc.get("hidden", False)]
                        if visible_cases:
                            st.markdown("#### Samples")
                            for idx, tc in enumerate(visible_cases, start=1):
                                st.markdown(
                                    f"""
                                    <div class="sample-case">
                                        <strong>Sample {idx}</strong>
                                    </div>
                                    """,
                                    unsafe_allow_html=True,
                                )
                                st.caption("Input")
                                st.code(tc.get("input", ""), language="text")
                                st.caption("Expected Output")
                                st.code(tc.get("expected_output", ""), language="text")
                        hidden_cases = [tc for tc in meta.get("test_cases", []) if tc.get("hidden", False)]
                        if hidden_cases and current_user_can_request_hidden_access("view"):
                            qid_key = str(question["id"])
                            if not st.session_state.hidden_test_access.get(qid_key):
                                view_pin = st.text_input(
                                    "Hidden test cases see PIN",
                                    type="password",
                                    key=f"hidden_view_pin_{question['id']}",
                                    max_chars=6,
                                )
                                if st.button("Unlock Hidden Test Cases", key=f"unlock_hidden_view_{question['id']}", use_container_width=True):
                                    ok, msg = verify_hidden_test_pin(view_pin, "view")
                                    if ok:
                                        st.session_state.hidden_test_access[qid_key] = True
                                        st.success("Hidden test cases unlocked.")
                                        st.rerun()
                                    else:
                                        st.error(msg)
                            if st.session_state.hidden_test_access.get(qid_key):
                                st.markdown("#### Hidden Test Cases")
                                if st.session_state.get("role") != "admin":
                                    st.caption(f"Remaining views: {get_hidden_case_views_remaining()}")
                                hidden_seen = 0
                                for case_no, tc in enumerate(meta.get("test_cases", []), start=1):
                                    if not tc.get("hidden", False):
                                        continue
                                    hidden_seen += 1
                                    viewed = has_hidden_case_viewed(question["id"], case_no)
                                    with st.container(border=True):
                                        st.caption(f"Hidden Case {hidden_seen} (Case {case_no})")
                                        if not viewed:
                                            if st.button("View This Case", key=f"view_hidden_case_{question['id']}_{case_no}", use_container_width=True):
                                                ok, msg = consume_hidden_case_view(question["id"], case_no)
                                                if ok:
                                                    st.success("Case unlocked.")
                                                    st.rerun()
                                                else:
                                                    st.error(msg)
                                        if viewed:
                                            st.caption("Input")
                                            st.code(tc.get("input", ""), language="text")
                                            st.caption("Expected Output")
                                            st.code(tc.get("expected_output", ""), language="text")

                        if hidden_cases and current_user_can_request_hidden_access("edit"):
                            qid_key = str(question["id"])
                            if not st.session_state.hidden_test_edit_access.get(qid_key):
                                edit_pin = st.text_input(
                                    "Hidden test cases edit PIN",
                                    type="password",
                                    key=f"hidden_edit_pin_{question['id']}",
                                    max_chars=6,
                                )
                                if st.button("Unlock Hidden Test Editing", key=f"unlock_hidden_edit_{question['id']}", use_container_width=True):
                                    ok, msg = verify_hidden_test_pin(edit_pin, "edit")
                                    if ok:
                                        st.session_state.hidden_test_edit_access[qid_key] = True
                                        st.success("Hidden test editing unlocked.")
                                        st.rerun()
                                    else:
                                        st.error(msg)
                            if st.session_state.hidden_test_edit_access.get(qid_key):
                                with st.form(f"hidden_case_edit_form_{question['id']}"):
                                    st.markdown("#### Edit Hidden Test Cases")
                                    rebuilt_cases = []
                                    hidden_seen = 0
                                    for tc_idx, tc in enumerate(meta.get("test_cases", [])):
                                        if not tc.get("hidden", False):
                                            rebuilt_cases.append(tc)
                                            continue
                                        hidden_seen += 1
                                        with st.container(border=True):
                                            st.caption(f"Hidden Case {hidden_seen}")
                                            edit_in = st.text_area("Input", value=tc.get("input", ""), key=f"hidden_edit_input_{question['id']}_{tc_idx}", height=75)
                                            edit_out = st.text_area("Expected Output", value=tc.get("expected_output", ""), key=f"hidden_edit_output_{question['id']}_{tc_idx}", height=75)
                                            edit_marks = st.number_input("Marks", min_value=1, max_value=100, value=int(tc.get("marks", 1) or 1), key=f"hidden_edit_marks_{question['id']}_{tc_idx}")
                                        rebuilt_cases.append({
                                            "input": edit_in,
                                            "expected_output": edit_out,
                                            "marks": int(edit_marks),
                                            "hidden": True,
                                        })
                                    if st.form_submit_button("Save Hidden Test Cases", type="primary", use_container_width=True):
                                        new_explanation = make_programming_meta(meta.get("description", ""), rebuilt_cases, selected_language if 'selected_language' in locals() else meta.get("language", "java"))
                                        supabase.table("questions").update({"explanation": new_explanation}).eq("id", question["id"]).execute()
                                        st.success("Hidden test cases updated.")
                                        st.rerun()
                        st.markdown("</div>", unsafe_allow_html=True)
                with p_right:
                    selected_language_label = st.selectbox(
                        "Language",
                        language_options,
                        index=language_options.index(current_language_label) if current_language_label in language_options else 0,
                        key=f"prog_language_{question['id']}"
                    )
                    selected_language = PROGRAMMING_LANGUAGE_LABELS[selected_language_label]
                    lang_meta = get_programming_language_meta(selected_language)
                    enable_textarea_tab_support()
                    enable_ide_textarea_behaviour(question["id"], selected_language)
                    editor_value = stored_code if stored_code else lang_meta.get("default_code", "")
                    st.markdown(
                        f"""
                        <div class="vscode-shell">
                            <div class="vscode-titlebar">
                                <div class="vscode-dots"><span></span><span></span><span></span></div>
                                <div class="vscode-file">{lang_meta.get('file_name', 'solution.txt')}</div>
                                <div class="vscode-hints">Tab: indent/snippet | brackets auto-close</div>
                            </div>
                        </div>
                        """,
                        unsafe_allow_html=True,
                    )
                    answer = st.text_area(f"{lang_meta['label']} Program", value=editor_value, key=f"code_{question['id']}", height=395)
                    st.session_state.answers[question["id"]] = {"code": answer, "language": selected_language}
                    save_programming_exam_session(status="active")
                    st.markdown("<div class='console-card'>", unsafe_allow_html=True)
                    custom_input = st.text_area(
                        "Custom Parameters / stdin",
                        key=f"custom_input_{question['id']}",
                        height=82,
                        placeholder="Input values ikkada enter cheyyandi..."
                    )
                    run_col, submit_prog_col, custom_col = st.columns(3)
                    with run_col:
                        run_tests_clicked = st.button("Compile", key=f"run_prog_{question['id']}", use_container_width=True)
                    with submit_prog_col:
                        submit_program_clicked = st.button("Execute & Submit", key=f"submit_prog_{question['id']}", type="primary", use_container_width=True)
                    with custom_col:
                        run_custom_clicked = st.button("Run Custom", key=f"run_custom_{question['id']}", use_container_width=True)

                    if run_tests_clicked:
                        with st.spinner("Test suite running..."):
                            st.session_state.program_run_results[str(question["id"])] = run_programming_test_cases(question, answer, selected_language)

                    if submit_program_clicked:
                        with st.spinner("Score saving..."):
                            score_data = run_programming_test_cases(question, answer, selected_language)
                            st.session_state.program_run_results[str(question["id"])] = score_data
                            st.session_state.program_submissions[str(question["id"])] = {"code": answer, "language": selected_language, "score_data": score_data}
                            save_programming_exam_session(status="active")
                            st.success(f"Saved: {score_data['earned']}/{score_data['total']} marks")

                    if run_custom_clicked:
                        with st.spinner("Custom run executing..."):
                            custom_result = run_programming_code(answer, custom_input, selected_language)
                            custom_result["language"] = selected_language
                            st.session_state.program_custom_results[str(question["id"])] = custom_result

                    saved_prog = st.session_state.program_submissions.get(str(question["id"]), {})
                    if saved_prog.get("code") == answer and normalize_programming_language(saved_prog.get("language", selected_language)) == selected_language and saved_prog.get("score_data"):
                        saved_score = saved_prog["score_data"]
                        st.success(f"Final submit ready: {saved_score['earned']}/{saved_score['total']} marks")
                    elif saved_prog:
                        st.warning("Code changed after score save. Execute & Submit again to update this program marks.")
                    else:
                        st.info("Compile code check chestundi. Execute & Submit marks update chestundi. Final submit lo unsaved programs 0 score ga count avuthayi.")

                    custom_data = st.session_state.program_custom_results.get(str(question["id"]))
                    if custom_data and normalize_programming_language(custom_data.get("language", selected_language)) != selected_language:
                        custom_data = None
                    if custom_data:
                        st.caption(f"Custom run: {custom_data.get('status','')}")
                        if custom_data.get("stdout"):
                            st.code(custom_data.get("stdout", ""), language="text")
                        if custom_data.get("stderr"):
                            st.code(custom_data.get("stderr", ""), language="text")

                    run_data = st.session_state.program_run_results.get(str(question["id"]))
                    if run_data and normalize_programming_language(run_data.get("language", selected_language)) != selected_language:
                        run_data = None
                    if run_data:
                        st.markdown(f"**Suite Result:** {run_data['earned']}/{run_data['total']} marks ({run_data['percentage']}%)")
                        for res in run_data["results"]:
                            is_hidden = bool(res.get("hidden", False))
                            status_class = "console-log-pass" if res["passed"] else "console-log-fail"
                            badge = "PASS" if res["passed"] else "FAIL"
                            title = f"Hidden Case {res['case']}" if is_hidden else f"Case {res['case']}"
                            st.markdown(
                                f"<div class='{status_class}'><strong>{badge}</strong> - {title} - {res['marks']} marks</div>",
                                unsafe_allow_html=True,
                            )
                            if is_hidden:
                                can_show_hidden_result = (
                                    st.session_state.hidden_test_access.get(str(question["id"]))
                                    and has_hidden_case_viewed(question["id"], int(res.get("case") or 0))
                                )
                                if can_show_hidden_result:
                                    st.caption(f"Status: {res.get('status','')}")
                                    c_in, c_exp, c_act = st.columns(3)
                                    with c_in:
                                        st.caption("Input")
                                        st.code(res.get("input", ""), language="text")
                                    with c_exp:
                                        st.caption("Expected")
                                        st.code(res.get("expected_output", ""), language="text")
                                    with c_act:
                                        st.caption("Actual")
                                        st.code(res.get("actual_output", ""), language="text")
                                    if res.get("error"):
                                        st.caption("Runtime / compiler log")
                                        st.code(res["error"], language="text")
                                else:
                                    st.caption("Hidden input/output counted in scoring.")
                            elif not res["passed"]:
                                st.caption(f"Status: {res.get('status','')}")
                                c_exp, c_act = st.columns(2)
                                with c_exp:
                                    st.caption("Expected")
                                    st.code(res.get("expected_output", ""), language="text")
                                with c_act:
                                    st.caption("Actual")
                                    st.code(res.get("actual_output", ""), language="text")
                                if res.get("error"):
                                    st.caption("Runtime / compiler log")
                                    st.code(res["error"], language="text")
                            if st.session_state.get("role") != "admin":
                                with st.expander(f"Report Case {res.get('case')}"):
                                    report_text = st.text_area("Issue / doubt", key=f"report_text_{question['id']}_{res.get('case')}", height=70)
                                    if st.button("Send Report", key=f"send_report_{question['id']}_{res.get('case')}", use_container_width=True):
                                        ok, msg = submit_test_case_report(
                                            st.session_state.exam_id,
                                            question["id"],
                                            int(res.get("case") or 0),
                                            report_text,
                                            res,
                                        )
                                        if ok:
                                            st.success(msg)
                                        else:
                                            st.error(msg)
                    st.markdown("</div>", unsafe_allow_html=True)

            def save_current_q_time():
                qid_cur = question["id"]
                if qid_cur in st.session_state.question_start_time:
                    elapsed = int(time.time() - st.session_state.question_start_time[qid_cur])
                    prev = st.session_state.question_time_log.get(qid_cur, 0)
                    st.session_state.question_time_log[qid_cur] = prev + elapsed
                    del st.session_state.question_start_time[qid_cur]

            nav_col1, nav_col2, pause_col, submit_col = st.columns([1, 1, 1.2, 2])
            with nav_col1:
                if st.button("Previous", disabled=(current <= section_start), use_container_width=True):
                    save_current_q_time(); st.session_state.question_index -= 1; save_programming_exam_session(status="active"); st.rerun()
            with nav_col2:
                if st.button("Next ", disabled=(current >= section_end - 1), use_container_width=True):
                    save_current_q_time(); st.session_state.question_index += 1; save_programming_exam_session(status="active"); st.rerun()
            with pause_col:
                if st.button("Pause & Back", use_container_width=True):
                    save_current_q_time()
                    if is_prog_exam:
                        st.session_state.exam_end_time = -max(1, int(st.session_state.exam_end_time - time.time()))
                    save_programming_exam_session(status="active")
                    st.session_state.start_exam = False
                    st.session_state.current_questions = []
                    st.rerun()
            with submit_col:
                if st.button("Submit Exam", type="primary", use_container_width=True):
                    save_current_q_time()
                    try:
                        save_programming_exam_session(status="active")
                        submit_exam_attempt(questions, include_time=True, require_programming_submitted=False)
                        st.session_state.question_time_log = {}
                        st.session_state.question_start_time = {}
                        st.session_state.program_run_results = {}
                        st.session_state.program_submissions = {}
                        st.session_state.exam_submitted = True; st.rerun()
                    except Exception as e:
                        st.error(f"Submit failed: {e}")

# =========================
# MAIN ROUTING
# =========================
if not st.session_state.logged_in:
    if st.session_state.user_id_temp:
        pin_screen()
    else:
        login()
else:
    if st.session_state.get("exam_lock_message"):
        st.error(st.session_state.exam_lock_message)
        st.session_state.exam_lock_message = ""
    if st.session_state.role == "admin":
        admin_dashboard()
    elif st.session_state.role == "card_user":
        card_user_dashboard()
    elif st.session_state.start_exam:
        exam_workspace_view()
    else:
        user_dashboard(preview_mode=False)










